from pathlib import Path
import shutil

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_parts2_3 import (
    title, footer, rect, text, formula, arrow, image,
    GREEN, GREEN_DARK, GREEN_LIGHT, BLUE, BLUE_LIGHT, ORANGE,
    ORANGE_LIGHT, CHARCOAL, MID, LIGHT, WHITE,
)


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_working_08.pptx"
OUTPUT = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_working_10.pptx"
ASSETS = ROOT / "Presentation" / "vector_field_foundations"
ASSETS.mkdir(parents=True, exist_ok=True)


def style_axis(ax):
    ax.set_aspect("equal")
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_xlabel("latent dimension 1")
    ax.set_ylabel("latent dimension 2")
    ax.grid(alpha=.12)


def vector_plot(path, mode="local"):
    grid = np.linspace(-2.5, 2.5, 17)
    xx, yy = np.meshgrid(grid, grid)
    A = np.array([[-.30, -1.0], [.90, -.30]])
    u = A[0, 0] * xx + A[0, 1] * yy
    v = A[1, 0] * xx + A[1, 1] * yy
    fig, ax = plt.subplots(figsize=(6.2, 5.0))
    ax.quiver(xx, yy, u, v, color="#B8DCCB", alpha=.92,
              angles="xy", scale_units="xy", scale=5.5, width=.006)
    if mode == "local":
        p = np.array([1.35, .65]); vel = A @ p
        ax.scatter(*p, s=120, color="#EE772F", zorder=5)
        ax.arrow(p[0], p[1], vel[0] * .55, vel[1] * .55,
                 width=.035, head_width=.18, color="#008A45", zorder=6)
        ax.annotate("current state  z", p + [.12, -.34], fontsize=12, weight="bold")
        ax.annotate("local change  f(z)", p + vel * .55 + [.08, .03],
                    fontsize=12, color="#006633", weight="bold")
    else:
        dt = .055; z = np.array([2.15, .25]); traj = [z.copy()]
        for _ in range(125):
            z = z + dt * (A @ z)
            traj.append(z.copy())
        traj = np.asarray(traj)
        ax.plot(traj[:, 0], traj[:, 1], color="#008A45", lw=3.0, zorder=4)
        ax.scatter(*traj[0], s=105, color="#EE772F", zorder=5)
        ax.scatter(0, 0, s=110, marker="X", color="#333537", zorder=5)
        ax.annotate("initial state", traj[0] + [.08, .15], fontsize=11, weight="bold")
        ax.annotate("fixed point", [.12, -.32], fontsize=11, weight="bold")
    style_axis(ax); ax.set_xlim(-2.7, 2.7); ax.set_ylim(-2.7, 2.7)
    fig.tight_layout(); fig.savefig(path, dpi=220, facecolor="white"); plt.close(fig)


def motif_plot(path):
    mats = [
        (np.array([[-.8, 0], [0, -.55]]), "point attractor", "memory / return"),
        (np.array([[.7, 0], [0, -.7]]), "saddle point", "decision boundary"),
        (np.array([[-.10, -1], [1, -.10]]), "rotation", "motor dynamics"),
        (np.array([[0, 0], [0, -.75]]), "line attractor", "integration"),
    ]
    grid = np.linspace(-2, 2, 13); xx, yy = np.meshgrid(grid, grid)
    fig, axes = plt.subplots(1, 4, figsize=(12.0, 3.1))
    for ax, (A, name, example) in zip(axes, mats):
        u = A[0, 0]*xx + A[0, 1]*yy; v = A[1, 0]*xx + A[1, 1]*yy
        speed = np.sqrt(u*u + v*v); speed[speed == 0] = 1
        ax.quiver(xx, yy, u/speed, v/speed, color="#74B99A", alpha=.9,
                  angles="xy", scale_units="xy", scale=5.4, width=.009)
        ax.scatter(0, 0, marker="X", s=50, color="#333537")
        if name == "line attractor": ax.axhline(0, color="#EE772F", lw=2.5)
        ax.set_title(name, fontsize=12, weight="bold")
        ax.text(.5, -.18, example, transform=ax.transAxes, ha="center", fontsize=10, color="#666")
        ax.set_xlim(-2.2, 2.2); ax.set_ylim(-2.2, 2.2); ax.set_aspect("equal")
        ax.set_xticks([]); ax.set_yticks([])
        for sp in ax.spines.values(): sp.set_color("#D8DDDA")
    fig.tight_layout(w_pad=1.2); fig.savefig(path, dpi=220, facecolor="white", bbox_inches="tight"); plt.close(fig)


def piecewise_plot(path):
    grid = np.linspace(-2.5, 2.5, 15); xx, yy = np.meshgrid(grid, grid)
    left = xx < 0
    u = np.where(left, -.65*(xx+1.15), -.65*(xx-1.15))
    v = -.70*yy
    fig, ax = plt.subplots(figsize=(7.2, 4.5))
    ax.axvspan(-2.7, 0, color="#EAF2FB"); ax.axvspan(0, 2.7, color="#FFF0E5")
    colors = np.where(left.ravel(), "#3F78B5", "#EE772F")
    ax.quiver(xx, yy, u, v, color=colors, alpha=.92, angles="xy", scale_units="xy", scale=5.0)
    ax.axvline(0, color="#777", lw=1.5, ls="--")
    ax.scatter([-1.15, 1.15], [0, 0], marker="X", s=90, color=["#3F78B5", "#EE772F"])
    ax.text(-2.35, 2.15, "regime 1:  A₁, b₁", color="#3F78B5", weight="bold", fontsize=13)
    ax.text(.35, 2.15, "regime 2:  A₂, b₂", color="#C45C1B", weight="bold", fontsize=13)
    style_axis(ax); ax.set_xlim(-2.6, 2.6); ax.set_ylim(-2.6, 2.6)
    fig.tight_layout(); fig.savefig(path, dpi=220, facecolor="white"); plt.close(fig)


def build():
    vector_plot(ASSETS / "local_arrow.png", "local")
    vector_plot(ASSETS / "trajectory.png", "trajectory")
    motif_plot(ASSETS / "motifs.png")
    piecewise_plot(ASSETS / "piecewise.png")

    shutil.copy2(SOURCE, OUTPUT)
    prs = Presentation(OUTPUT)
    original_plds_id = prs.slides._sldIdLst[16]
    blank_layout = prs.slide_layouts[0]
    new_ids = []

    # 17 — define a point before defining dynamics.
    s = prs.slides.add_slide(blank_layout); new_ids.append(prs.slides._sldIdLst[-1]); rect(s, 0, 0, 13.333, 7.5, WHITE, None, False)
    title(s, "A point in latent space is a neural population state", 17, "VECTOR FIELD FOUNDATIONS")
    text(s, "observed population activity", .82, 1.52, 3.05, .34, 16, CHARCOAL, True, PP_ALIGN.CENTER)
    vals = [.82, .28, .68, .43, .19]
    for i, val in enumerate(vals):
        rect(s, .88, 1.98+i*.58, 2.85*val, .30, BLUE, BLUE, False)
        text(s, f"neuron {i+1}", 3.38, 1.94+i*.58, .72, .30, 11, MID, False)
    formula(s, "yₜ ∈ ℝᴺ", 1.20, 5.18, 2.45, .65, 23, BLUE_LIGHT, BLUE)
    arrow(s, 4.20, 3.25, 5.45, 3.25, GREEN, 2.3)
    text(s, "dimensionality\nreduction / inference", 4.20, 2.47, 1.25, .62, 13, GREEN_DARK, True, PP_ALIGN.CENTER)
    rect(s, 5.82, 1.45, 6.32, 4.55, WHITE, GREEN, True, 1.1)
    arrow(s, 8.82, 5.25, 8.82, 2.05, MID, 1.5)
    arrow(s, 6.55, 4.92, 11.32, 4.92, MID, 1.5)
    text(s, "latent 2", 6.05, 2.22, .85, .30, 12, MID, False)
    text(s, "latent 1", 10.73, 5.10, .85, .30, 12, MID, False)
    rect(s, 9.45, 3.25, .26, .26, ORANGE, ORANGE, True)
    text(s, "zₜ", 9.78, 3.17, .60, .34, 20, ORANGE, True)
    text(s, "one time point → one location", 7.25, 5.52, 3.80, .34, 17, GREEN_DARK, True, PP_ALIGN.CENTER)
    rect(s, 1.02, 6.18, 11.28, .52, GREEN_LIGHT, GREEN, True, .8)
    text(s, "A trajectory is simply the sequence  z₁, z₂, …, zₜ  traced through this space.", 1.28, 6.27, 10.76, .32, 17, GREEN_DARK, True, PP_ALIGN.CENTER)
    footer(s, 17)

    # 18 — vector field as a local rule.
    s = prs.slides.add_slide(blank_layout); new_ids.append(prs.slides._sldIdLst[-1]); rect(s, 0, 0, 13.333, 7.5, WHITE, None, False)
    title(s, "A vector field answers one local question", 18, "VECTOR FIELD FOUNDATIONS")
    formula(s, "żₜ = f(zₜ)", .72, 1.35, 3.75, .82, 28, GREEN_LIGHT, GREEN)
    text(s, "Given the current state zₜ,\nwhich direction—and how fast—does it change?", .75, 2.55, 3.75, 1.15, 20, CHARCOAL, True, PP_ALIGN.CENTER)
    rect(s, .83, 4.15, 3.55, 1.20, ORANGE_LIGHT, ORANGE, True, 1.0)
    text(s, "point = where the system is\narrow = its instantaneous velocity", 1.05, 4.42, 3.10, .72, 17, ORANGE, True, PP_ALIGN.CENTER)
    image(s, ASSETS / "local_arrow.png", 4.80, 1.05, 7.45, 5.25)
    rect(s, 1.10, 6.18, 11.10, .52, GREEN_LIGHT, GREEN, True, .8)
    text(s, "The field is not the trajectory: it is the rule defined at every possible state.", 1.35, 6.27, 10.60, .32, 17, GREEN_DARK, True, PP_ALIGN.CENTER)
    footer(s, 18)

    # 19 — integrate local changes to obtain a trajectory.
    s = prs.slides.add_slide(blank_layout); new_ids.append(prs.slides._sldIdLst[-1]); rect(s, 0, 0, 13.333, 7.5, WHITE, None, False)
    title(s, "A trajectory is generated by repeatedly following the field", 19, "VECTOR FIELD FOUNDATIONS")
    formula(s, "zₜ₊₁ ≈ zₜ + Δt · f(zₜ)", .62, 1.28, 4.28, .78, 25, GREEN_LIGHT, GREEN)
    text(s, "Start from an initial state", .82, 2.38, 3.75, .34, 18, CHARCOAL, True, PP_ALIGN.CENTER)
    formula(s, "z₀", 1.65, 2.90, 2.10, .62, 23, ORANGE_LIGHT, ORANGE)
    text(s, "follow the local arrow", .82, 3.82, 3.75, .34, 18, CHARCOAL, True, PP_ALIGN.CENTER)
    formula(s, "z₀ → z₁ → z₂ → ···", .82, 4.35, 3.75, .72, 22, LIGHT, None)
    text(s, "The same field can generate many trajectories\nfrom different initial states.", .80, 5.38, 3.80, .70, 16, MID, False, PP_ALIGN.CENTER)
    image(s, ASSETS / "trajectory.png", 4.88, 1.04, 7.30, 5.35)
    rect(s, 1.08, 6.20, 11.12, .50, BLUE_LIGHT, BLUE, True, .8)
    text(s, "Data give us trajectories; dynamical modeling tries to infer the field that generated them.", 1.30, 6.28, 10.67, .30, 17, BLUE, True, PP_ALIGN.CENTER)
    footer(s, 19)

    # 20 — interpret fields through motifs and fixed points.
    s = prs.slides.add_slide(blank_layout); new_ids.append(prs.slides._sldIdLst[-1]); rect(s, 0, 0, 13.333, 7.5, WHITE, None, False)
    title(s, "Field geometry reveals candidate neural computations", 20, "DYNAMICAL MOTIFS")
    formula(s, "fixed point  z* :  f(z*) = 0", .73, 1.48, 4.30, .68, 23, ORANGE_LIGHT, ORANGE)
    text(s, "At a fixed point the expected velocity is zero.\nNearby arrows tell us whether it is stable or unstable.", 5.40, 1.66, 6.58, .58, 17, CHARCOAL, True)
    image(s, ASSETS / "motifs.png", .68, 2.28, 11.95, 3.48)
    rect(s, 1.02, 6.10, 11.30, .62, GREEN_LIGHT, GREEN, True, .8)
    text(s, "These are hypotheses about computation—not merely shapes in a dimensionality-reduction plot.", 1.27, 6.23, 10.80, .32, 17, GREEN_DARK, True, PP_ALIGN.CENTER)
    footer(s, 20, "Motifs adapted conceptually from the SLDS lecture; diagrams redrawn for this presentation")

    # 21 — one matrix generates a complete linear field.
    s = prs.slides.add_slide(blank_layout); new_ids.append(prs.slides._sldIdLst[-1]); rect(s, 0, 0, 13.333, 7.5, WHITE, None, False)
    title(s, "An LDS assumes one linear rule across the whole space", 21, "LINEAR DYNAMICAL SYSTEM")
    formula(s, "żₜ = A zₜ + b", .70, 1.30, 4.32, .78, 28, GREEN_LIGHT, GREEN)
    text(s, "A shapes the flow", .75, 2.48, 4.18, .34, 18, CHARCOAL, True, PP_ALIGN.CENTER)
    items = [("negative real part", "contraction / stability"), ("positive real part", "expansion / instability"), ("imaginary part", "rotation / oscillation")]
    for i, (a, b) in enumerate(items):
        y = 3.02 + i*.83
        rect(s, .78, y, 4.08, .62, WHITE, [BLUE, ORANGE, GREEN][i], True, .8)
        text(s, a, .96, y+.10, 1.70, .30, 14, [BLUE, ORANGE, GREEN][i], True)
        text(s, b, 2.62, y+.10, 2.05, .30, 14, CHARCOAL, False, PP_ALIGN.RIGHT)
    image(s, ASSETS / "trajectory.png", 5.30, 1.13, 6.75, 4.85)
    formula(s, "z* = −A⁻¹b   (when A is invertible)", 5.62, 5.35, 6.18, .62, 20, ORANGE_LIGHT, ORANGE)
    rect(s, 1.05, 6.20, 11.22, .50, LIGHT, MID, True, .7)
    text(s, "Strength: interpretable and data-efficient.  Limitation: one global linear field cannot express every nonlinear computation.", 1.27, 6.28, 10.78, .30, 16, CHARCOAL, True, PP_ALIGN.CENTER)
    footer(s, 21)

    # 22 — SLDS combines discrete regimes with continuous dynamics.
    s = prs.slides.add_slide(blank_layout); new_ids.append(prs.slides._sldIdLst[-1]); rect(s, 0, 0, 13.333, 7.5, WHITE, None, False)
    title(s, "SLDS: discrete regimes, continuous trajectories", 22, "SWITCHING LDS")
    formula(s, "sₜ ∈ {1,…,K}", .55, 1.48, 3.45, .65, 22, ORANGE_LIGHT, ORANGE)
    formula(s, "zₜ₊₁ = Aₛₜ zₜ + bₛₜ + εₜ", .55, 2.27, 3.45, .80, 21, GREEN_LIGHT, GREEN)
    formula(s, "yₜ ∼ p(yₜ | zₜ)", .55, 3.20, 3.45, .68, 21, BLUE_LIGHT, BLUE)
    text(s, "sₜ selects the local dynamical rule\nzₜ remains a continuous neural state", .68, 4.18, 3.20, .82, 17, CHARCOAL, True, PP_ALIGN.CENTER)
    rect(s, .72, 5.32, 3.12, .62, LIGHT, MID, True, .8)
    text(s, "HMM state + LDS trajectory", .88, 5.47, 2.80, .30, 16, CHARCOAL, True, PP_ALIGN.CENTER)
    image(s, ASSETS / "piecewise.png", 4.25, 1.05, 8.10, 5.18)
    rect(s, 1.00, 6.15, 11.35, .55, GREEN_LIGHT, GREEN, True, .8)
    text(s, "Piecewise-linear fields approximate nonlinear dynamics while preserving interpretable local regimes.", 1.25, 6.25, 10.85, .32, 17, GREEN_DARK, True, PP_ALIGN.CENTER)
    footer(s, 22)

    # Move the six new slides immediately before the original PLDS slide.
    slide_ids = prs.slides._sldIdLst
    for sid in new_ids:
        slide_ids.remove(sid)
    insertion = list(slide_ids).index(original_plds_id)
    for offset, sid in enumerate(new_ids):
        slide_ids.insert(insertion + offset, sid)

    # Renumber visible slide-number labels after reordering.
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            value = shape.text.strip()
            if value.isdigit() and 1 <= int(value) <= 80:
                shape.text_frame.paragraphs[0].runs[0].text = str(i)
                if shape.width < Inches(.42):
                    shape.width = Inches(.42)
    # Clarify the landing slide after the foundations sequence.
    plds = prs.slides[22]
    for shape in plds.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.startswith("PLDS:"):
            shape.text_frame.paragraphs[0].runs[0].text = "PLDS: observe continuous neural dynamics through spikes"
            break

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
