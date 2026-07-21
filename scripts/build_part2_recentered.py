from pathlib import Path
import shutil

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

from build_parts2_3 import (
    clear, title, footer, rect, text, formula, arrow, image, card,
    GREEN, GREEN_DARK, GREEN_LIGHT, BLUE, BLUE_LIGHT, ORANGE,
    ORANGE_LIGHT, CHARCOAL, MID, LIGHT, WHITE, MATH,
)


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_working_04.pptx"
OUTPUT = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_working_06.pptx"
TUTORIAL = ROOT / "Presentation" / "tutorial_assets_part2"
ASSETS = ROOT / "Presentation" / "assets_v2"


def node(slide, label, x, y, color, w=.66, h=.52):
    rect(slide, x, y, w, h, WHITE, color, True, 1.2)
    text(slide, label, x+.03, y+.05, w-.06, h-.10, 17, color, True, font=MATH)


def build():
    shutil.copy2(SOURCE, OUTPUT)
    prs = Presentation(OUTPUT)

    # 11 — orient the audience before introducing individual models.
    s=prs.slides[10]; clear(s)
    title(s,"One inverse problem, three ways to structure the latent state",11,"MODEL LADDER")
    text(s,"observations",.62,1.10,2.0,.34,14,MID,True)
    formula(s,"y",.65,1.56,1.20,.66,25,BLUE_LIGHT,BLUE)
    arrow(s,2.05,1.89,2.92,1.89,GREEN,2.0)
    text(s,"infer",2.10,1.42,.78,.30,12,GREEN_DARK,True,align=1)
    formula(s,"z",3.05,1.56,1.20,.66,25,GREEN_LIGHT,GREEN)
    text(s,"What kind of object should the hidden state z be?",4.68,1.52,7.35,.72,25,CHARCOAL,True)

    cols=[
        ("FINITE MIXTURE","independent discrete assignment","zₙ ∈ {1,…,K}",BLUE),
        ("FINITE-STATE HMM","temporally dependent discrete state","zₜ₋₁ → zₜ",ORANGE),
        ("LDS","temporally dependent continuous state","zₜ = A zₜ₋₁ + εₜ",GREEN),
    ]
    for i,(name,desc,eq,col) in enumerate(cols):
        x=.55+i*4.22
        rect(s,x,2.75,3.78,2.75,WHITE,col,True,1.1)
        text(s,name,x+.22,3.02,3.34,.36,18,col,True,align=1)
        text(s,desc,x+.25,3.58,3.28,.38,16,CHARCOAL,True,align=1)
        formula(s,eq,x+.42,4.26,2.94,.66,19,LIGHT,None)
    arrow(s,4.36,4.13,4.66,4.13,GREEN,1.6)
    arrow(s,8.58,4.13,8.88,4.13,GREEN,1.6)
    rect(s,1.25,5.90,10.85,.63,GREEN_LIGHT,GREEN,True,.8)
    text(s,"independent assignment  →  temporal state  →  continuous dynamics",1.48,6.03,10.40,.35,19,GREEN_DARK,True,align=1)
    footer(s,11)

    # 12 — mixture: Bayes responsibilities plus spatial geometry.
    s=prs.slides[11]; clear(s)
    title(s,"Finite mixtures infer uncertain component identity",12,"FINITE MIXTURE")
    text(s,"Generative story",.62,1.10,2.45,.34,16,CHARCOAL,True)
    formula(s,"zₙ ∼ Categorical(π)",.62,1.55,4.05,.70,22,BLUE_LIGHT,BLUE)
    formula(s,"xₙ | zₙ=k ∼ p(x | θₖ)",.62,2.48,4.05,.70,21,GREEN_LIGHT,GREEN)
    text(s,"Bayes rule gives a soft assignment",.62,3.54,4.08,.42,17,CHARCOAL,True)
    formula(s,"rₙₖ = p(zₙ=k | xₙ)\n= πₖ p(xₙ|θₖ) / ∑ⱼ πⱼ p(xₙ|θⱼ)",.62,4.05,4.10,1.22,19,ORANGE_LIGHT,ORANGE)
    image(s,TUTORIAL/"mix_cell9_0.png",5.12,1.22,7.55,4.62)
    rect(s,.92,5.92,11.55,.62,GREEN_LIGHT,GREEN,True,.8)
    text(s,"Geometry: posterior confidence falls near overlapping component boundaries.",1.18,6.05,11.02,.34,18,GREEN_DARK,True,align=1)
    footer(s,12,"Tutorial figure: Part1_Basics/02b_mixtures_em.ipynb")

    # 13 — HMM: temporal graph plus posterior heat map.
    s=prs.slides[12]; clear(s)
    title(s,"A finite-state HMM connects discrete assignments through time",13,"HMM")
    text(s,"Latent states",.56,1.24,1.10,.30,12,GREEN_DARK,True)
    xs=[1.52,2.76,4.00,5.24]
    for i,x in enumerate(xs):
        node(s,f"z{i+1}",x,1.12,GREEN)
        node(s,f"y{i+1}",x,2.32,BLUE)
        arrow(s,x+.33,1.67,x+.33,2.25,BLUE,1.5)
        if i<3: arrow(s,x+.70,1.38,xs[i+1]-.06,1.38,GREEN,1.7)
    text(s,"observations",.56,2.43,1.10,.30,12,BLUE,True)
    formula(s,"p(z₁:T,y₁:T) = p(z₁) ∏ₜ p(zₜ|zₜ₋₁) ∏ₜ p(yₜ|zₜ)",.58,3.14,5.63,.76,19,LIGHT,None)
    text(s,"The posterior is no longer one label per point:",.72,4.34,5.20,.36,16,CHARCOAL,True,align=1)
    formula(s,"p(zₜ=k | y₁:T)",1.42,4.88,3.70,.68,23,GREEN_LIGHT,GREEN)
    image(s,TUTORIAL/"hmm_cell8_0.png",6.48,1.10,6.32,4.95)
    rect(s,1.05,6.05,11.20,.55,GREEN_LIGHT,GREEN,True,.8)
    text(s,"Filtering uses the past; smoothing uses the complete sequence to revise state uncertainty.",1.30,6.15,10.70,.32,17,GREEN_DARK,True,align=1)
    footer(s,13,"Tutorial figure: Part2_Dynamics/03a_hmm_foundations.ipynb")

    # 14 — LDS: equations alongside continuous phase-space geometry.
    s=prs.slides[13]; clear(s)
    title(s,"LDS replaces state labels with a continuous latent trajectory",14,"LDS")
    text(s,"Dynamics model",.62,1.15,2.25,.34,16,GREEN_DARK,True)
    formula(s,"zₜ = A zₜ₋₁ + εₜ",.62,1.58,4.22,.75,24,GREEN_LIGHT,GREEN)
    text(s,"Observation model",.62,2.66,2.25,.34,16,BLUE,True)
    formula(s,"yₜ = C zₜ + ηₜ",.62,3.09,4.22,.75,24,BLUE_LIGHT,BLUE)
    card(s,"A  shapes the flow","rotation · contraction · expansion · timescale",.62,4.22,4.22,1.25,GREEN,15)
    image(s,ASSETS/"lds_flow.png",5.20,1.15,7.43,4.95)
    formula(s,"posterior = prediction + Kalman gain × prediction error",1.30,6.05,10.72,.58,19,ORANGE_LIGHT,ORANGE,ORANGE)
    footer(s,14,"Phase-space illustration: continuous state, vector field, and latent trajectory")

    # 15 — transition: keep the scientific question central.
    s=prs.slides[14]; clear(s)
    title(s,"Classical models trade flexibility for interpretable structure",15,"MODEL TRANSITION")
    rows=[
        ("FINITE MIXTURE","Who generated this observation?","independent discrete assignments",BLUE),
        ("FINITE-STATE HMM","Which discrete regime are we in?","temporally dependent discrete states",ORANGE),
        ("LDS","Where are we in a continuous flow?","linear dynamics + Gaussian noise",GREEN),
    ]
    for i,(name,q,lim,col) in enumerate(rows):
        y=1.22+i*1.37
        rect(s,.62,y,11.98,1.02,WHITE,col,True,.9)
        text(s,name,.92,y+.18,1.42,.36,15,col,True)
        text(s,q,2.52,y+.17,4.30,.38,18,CHARCOAL,True)
        text(s,lim,7.08,y+.17,4.95,.38,16,MID,False,align=1)
    text(s,"Neural dynamics may be continuous, nonlinear, stochastic, and observed through spikes.",.82,5.52,11.70,.48,21,CHARCOAL,True,align=1)
    rect(s,1.42,6.05,10.48,.72,GREEN_LIGHT,GREEN,True,.9)
    text(s,"Next question: how can inference remain tractable when the model becomes flexible?",1.70,6.15,9.92,.52,18,GREEN_DARK,True,align=1)
    footer(s,15)

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
