"""Build the 60-minute English group-talk deck from executed tutorial outputs."""

from __future__ import annotations

import base64
from pathlib import Path

import matplotlib.pyplot as plt
import nbformat
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentation"
ASSETS = OUT / "assets"
OUT.mkdir(exist_ok=True)
ASSETS.mkdir(exist_ok=True)

W, H = Inches(13.333), Inches(7.5)
BG = RGBColor(248, 249, 247)
INK = RGBColor(31, 42, 46)
MUTED = RGBColor(95, 109, 113)
TEAL = RGBColor(0, 128, 122)
BLUE = RGBColor(50, 101, 156)
ORANGE = RGBColor(224, 122, 59)
YELLOW = RGBColor(235, 187, 68)
PALE_TEAL = RGBColor(220, 239, 236)
PALE_BLUE = RGBColor(224, 234, 244)
PALE_ORANGE = RGBColor(249, 231, 216)
WHITE = RGBColor(255, 255, 255)
FONT = "Aptos"


def extract_notebook_image(notebook: str, cell_index: int, filename: str) -> Path:
    path = ROOT / notebook
    nb = nbformat.read(path, as_version=4)
    images = [
        output.data["image/png"]
        for output in nb.cells[cell_index].get("outputs", [])
        if "image/png" in output.get("data", {})
    ]
    if not images:
        raise ValueError(f"No PNG output in {path}, cell {cell_index}")
    target = ASSETS / filename
    target.write_bytes(base64.b64decode(images[0]))
    return target


FIGURES = {
    "mixture": extract_notebook_image(
        "VAE-Tutorial/Part1_Basics/02b_mixtures_em.ipynb", 8, "mixture_em.png"
    ),
    "responsibilities": extract_notebook_image(
        "VAE-Tutorial/Part1_Basics/02b_mixtures_em.ipynb", 9, "responsibilities.png"
    ),
    "hmm_smoothing": extract_notebook_image(
        "VAE-Tutorial/Part2_Dynamics/03a_hmm_foundations.ipynb", 8, "hmm_smoothing.png"
    ),
    "hmm_em": extract_notebook_image(
        "VAE-Tutorial/Part2_Dynamics/03a_hmm_foundations.ipynb", 14, "hmm_em.png"
    ),
    "vae_arch": extract_notebook_image(
        "VAE-Tutorial/Part3_Variational/04_standard_vae.ipynb", 14, "vae_architecture.png"
    ),
    "vae_latent": extract_notebook_image(
        "VAE-Tutorial/Part3_Variational/04_standard_vae.ipynb", 24, "vae_latent.png"
    ),
    "lfads_arch": extract_notebook_image(
        "VAE-Tutorial/Part3_Variational/06_lfads.ipynb", 16, "lfads_architecture.png"
    ),
    "lfads_dynamics": extract_notebook_image(
        "VAE-Tutorial/Part3_Variational/06_lfads.ipynb", 51, "lfads_dynamics.png"
    ),
    "gpslds_arch": extract_notebook_image(
        "VAE-Tutorial/Part4_Advanced/08_gpslds.ipynb", 5, "gpslds_architecture.png"
    ),
    "gpslds_local": extract_notebook_image(
        "VAE-Tutorial/Part4_Advanced/08_gpslds.ipynb", 15, "gpslds_local.png"
    ),
}


def equation_png(tex: str, filename: str, fontsize: int = 24) -> Path:
    target = ASSETS / filename
    fig = plt.figure(figsize=(11, 0.9), facecolor="none")
    fig.text(0.5, 0.5, f"${tex}$", ha="center", va="center", fontsize=fontsize, color="#1f2a2e")
    plt.axis("off")
    fig.savefig(target, dpi=220, transparent=True, bbox_inches="tight", pad_inches=0.08)
    plt.close(fig)
    return target


EQ = {
    "bayes": equation_png(
        r"p(z\mid x,\theta)=\frac{p(x\mid z,\theta)p(z\mid\theta)}{p(x\mid\theta)}",
        "eq_bayes.png",
    ),
    "gradient": equation_png(
        r"\theta_{k+1}=\theta_k-\eta\,\nabla_\theta\mathcal{L}(\theta_k)",
        "eq_gradient.png",
    ),
    "gmm": equation_png(
        r"p(x_n\mid\theta)=\sum_{k=1}^{K}\pi_k\,\mathcal{N}(x_n\mid\mu_k,\Sigma_k)",
        "eq_gmm.png",
    ),
    "elbo": equation_png(
        r"\log p(x)=\mathcal{L}(q,\theta)+\mathrm{KL}\!\left[q(z)\,\|\,p(z\mid x,\theta)\right]",
        "eq_elbo.png",
    ),
    "hmm": equation_png(
        r"p(z_{1:T},x_{1:T})=p(z_1)\prod_{t=2}^{T}p(z_t\mid z_{t-1})\prod_{t=1}^{T}p(x_t\mid z_t)",
        "eq_hmm.png",
    ),
    "vae": equation_png(
        r"\mathcal{L}=\mathbb{E}_{q_\phi(z\mid x)}[\log p_\theta(x\mid z)]-\mathrm{KL}[q_\phi(z\mid x)\|p(z)]",
        "eq_vae.png",
    ),
    "reparam": equation_png(
        r"z=\mu_\phi(x)+\sigma_\phi(x)\odot\epsilon,\qquad \epsilon\sim\mathcal{N}(0,I)",
        "eq_reparam.png",
    ),
}


prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill=WHITE, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def text_box(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=22,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font=FONT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0.03)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def rich_lines(slide, lines, x, y, w, h, size=22, color=INK, gap=9):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.03)
    for i, line in enumerate(lines):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap)
        p.level = 0
    return box


def add_picture_contain(slide, path, x, y, w, h):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    return slide.shapes.add_picture(
        str(path), Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2), Inches(pw), Inches(ph)
    )


def base_slide(title, kicker=None, dark=False):
    slide = prs.slides.add_slide(BLANK)
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = INK if dark else BG
    if kicker:
        text_box(slide, kicker.upper(), 0.65, 0.34, 5.8, 0.3, 10, TEAL if not dark else YELLOW, True)
    text_box(slide, title, 0.65, 0.72 if kicker else 0.47, 12.0, 0.65, 28, WHITE if dark else INK, True)
    return slide


def footer(slide, number, source=None, dark=False):
    color = RGBColor(186, 196, 197) if dark else RGBColor(135, 147, 149)
    if source:
        text_box(slide, source, 0.67, 7.12, 10.9, 0.18, 8, color)
    text_box(slide, f"{number:02d}", 11.95, 7.05, 0.7, 0.25, 9, color, True, PP_ALIGN.RIGHT)


def notes(slide, timing, script):
    slide.notes_slide.notes_text_frame.text = f"Timing: {timing}\n\n{script}"


def title_slide(title, subtitle, script):
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = INK
    rect(slide, 0, 0, 0.17, 7.5, TEAL)
    text_box(slide, "GROUP SEMINAR", 0.85, 0.78, 3.0, 0.3, 11, YELLOW, True)
    text_box(slide, title, 0.85, 1.45, 11.5, 1.65, 38, WHITE, True)
    text_box(slide, subtitle, 0.88, 3.43, 10.8, 0.65, 20, RGBColor(201, 213, 214))
    # A compact visual grammar for the model progression.
    xs = [1.1, 3.25, 5.4, 7.55, 9.7]
    labels = ["latent z", "EM", "HMM", "VAE / LFADS", "gpSLDS"]
    colors = [BLUE, ORANGE, TEAL, YELLOW, RGBColor(188, 88, 95)]
    for i, (x, label, color) in enumerate(zip(xs, labels, colors)):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(5.25), Inches(1.15), Inches(1.15))
        shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.color.rgb = color
        text_box(slide, label, x - 0.14, 6.53, 1.45, 0.25, 10, WHITE, True, PP_ALIGN.CENTER)
        if i < 4:
            text_box(slide, "→", x + 1.35, 5.56, 0.52, 0.35, 21, RGBColor(170, 184, 185), True, PP_ALIGN.CENTER)
    notes(slide, "1:00", script)
    return slide


def section_slide(number, title, subtitle, script):
    slide = base_slide(title, f"Part {number}", dark=True)
    text_box(slide, subtitle, 0.68, 2.0, 10.8, 0.8, 22, RGBColor(203, 215, 216))
    text_box(slide, number, 10.85, 3.55, 1.6, 1.6, 74, TEAL, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    notes(slide, "0:30", script)
    return slide


def image_slide(title, path, caption, script, source, timing="1:30"):
    slide = base_slide(title)
    rect(slide, 0.7, 1.35, 11.95, 5.38, WHITE, RGBColor(226, 230, 228), True)
    add_picture_contain(slide, path, 0.9, 1.52, 11.55, 4.65)
    text_box(slide, caption, 0.96, 6.25, 11.2, 0.28, 11, MUTED, False, PP_ALIGN.CENTER)
    footer(slide, len(prs.slides), source)
    notes(slide, timing, script)
    return slide


def two_column_slide(title, left_title, left_lines, right_title, right_lines, script, timing="1:30"):
    slide = base_slide(title)
    for x, heading, lines, fill in [
        (0.7, left_title, left_lines, PALE_BLUE),
        (6.78, right_title, right_lines, PALE_TEAL),
    ]:
        rect(slide, x, 1.5, 5.82, 4.95, fill, fill, True)
        text_box(slide, heading, x + 0.35, 1.88, 5.1, 0.42, 18, INK, True)
        rich_lines(slide, lines, x + 0.36, 2.58, 5.0, 3.35, 18, INK, 15)
    footer(slide, len(prs.slides))
    notes(slide, timing, script)
    return slide


title_slide(
    "Variational Neural Inference",
    "From probability and gradients to interpretable neural dynamics",
    "Today I will build one continuous story from the basic language of probability to modern latent dynamical models. The goal is not to catalogue algorithms. The goal is to understand which modeling assumption is added at each step, which posterior quantity must be computed, and what scientific interpretation becomes possible.",
)

# 2
slide = base_slide("One question connects the entire talk")
text_box(slide, "How can noisy, high-dimensional observations\nreveal hidden neural computation?", 1.0, 1.52, 11.3, 1.25, 25, INK, True, PP_ALIGN.CENTER)
for x, label, sub, color in [
    (1.15, "observations", "spikes, behavior", BLUE),
    (4.95, "latent structure", "states, factors, dynamics", TEAL),
    (8.75, "interpretation", "mechanism, uncertainty", ORANGE),
]:
    rect(slide, x, 3.5, 3.35, 1.55, WHITE, color, True)
    text_box(slide, label, x + 0.2, 3.83, 2.95, 0.36, 19, color, True, PP_ALIGN.CENTER)
    text_box(slide, sub, x + 0.2, 4.38, 2.95, 0.3, 13, MUTED, False, PP_ALIGN.CENTER)
text_box(slide, "→", 4.53, 4.05, 0.4, 0.4, 25, MUTED, True, PP_ALIGN.CENTER)
text_box(slide, "→", 8.32, 4.05, 0.4, 0.4, 25, MUTED, True, PP_ALIGN.CENTER)
footer(slide, 2)
notes(slide, "1:15", "Neural data are observations, not explanations. Latent-variable models insert an explicit hidden layer between measurement and interpretation. Every model in this talk differs in the structure assigned to that hidden layer and in the method used to infer it.")

# 3
slide = base_slide("Roadmap: each model adds one assumption")
stages = [
    ("Probability", "uncertainty", BLUE),
    ("Mixtures", "discrete latent classes", ORANGE),
    ("HMMs", "temporal dependence", TEAL),
    ("VAEs", "amortized nonlinear inference", YELLOW),
    ("LFADS", "recurrent neural dynamics", RGBColor(160, 91, 140)),
    ("gpSLDS", "interpretable nonlinear flow", RGBColor(188, 88, 95)),
]
for i, (name, sub, color) in enumerate(stages):
    y = 1.45 + i * 0.86
    rect(slide, 1.25, y, 0.18, 0.56, color)
    text_box(slide, name, 1.7, y - 0.02, 2.25, 0.38, 18, INK, True)
    text_box(slide, sub, 4.15, y + 0.02, 6.2, 0.3, 16, MUTED)
footer(slide, 3)
notes(slide, "1:15", "The progression is cumulative. Probability gives us uncertainty. Mixtures introduce hidden assignments. HMMs connect those assignments through time. VAEs replace per-datapoint optimization with a learned inference network. LFADS equips the latent state with recurrent dynamics. gpSLDS then trades some neural-network flexibility for local dynamical interpretation.")

section_slide("1", "Foundations", "Probability, latent variables, gradients, and learning", "I will start from the smallest set of ideas needed for every later model.")

# 5
two_column_slide(
    "A probabilistic model is a data-generating hypothesis",
    "Observed variables",
    ["Measured directly", "Examples: spike counts, images, trajectories", "Written as x"],
    "Latent variables and parameters",
    ["Latent z explains hidden structure", "Parameters θ specify the model", "A joint distribution connects them"],
    "A probabilistic model says how data could have been generated. Observed variables are what the experiment records. Latent variables are explanatory quantities that are not directly measured. Parameters define the strength and shape of the relationships. This distinction prevents us from treating a low-dimensional visualization as if it were automatically a generative explanation.",
    "1:40",
)

# 6
slide = base_slide("Bayes' rule turns a generative model into inference")
add_picture_contain(slide, EQ["bayes"], 1.0, 1.55, 11.3, 1.25)
labels = [
    ("posterior", "what we infer after seeing data", TEAL),
    ("likelihood", "how the latent state explains x", BLUE),
    ("prior", "belief before observing x", ORANGE),
    ("evidence", "normalization; often the hard term", YELLOW),
]
for i, (name, sub, color) in enumerate(labels):
    x = 0.78 + (i % 2) * 6.1; y = 3.3 + (i // 2) * 1.25
    rect(slide, x, y, 5.55, 0.92, WHITE, color, True)
    text_box(slide, name, x + 0.25, y + 0.18, 1.35, 0.3, 16, color, True)
    text_box(slide, sub, x + 1.65, y + 0.18, 3.55, 0.42, 14, MUTED)
footer(slide, 6)
notes(slide, "1:40", "Bayes' rule is the central computational operation in latent-variable modeling. The numerator combines the likelihood and prior. The denominator, the marginal likelihood or evidence, sums or integrates over all latent explanations. In simple models this is analytic. In modern models it is usually the source of intractability.")

# 7
slide = base_slide("Learning and inference are different operations")
for x, title, q, result, color in [
    (0.85, "Inference", "Given θ and x", "estimate z", TEAL),
    (4.8, "Learning", "Given x", "estimate θ", BLUE),
    (8.75, "Prediction", "Given new x", "predict missing or future data", ORANGE),
]:
    rect(slide, x, 1.65, 3.65, 4.35, WHITE, color, True)
    text_box(slide, title, x + 0.25, 2.02, 3.15, 0.4, 20, color, True, PP_ALIGN.CENTER)
    text_box(slide, q, x + 0.35, 3.0, 2.95, 0.38, 15, MUTED, False, PP_ALIGN.CENTER)
    text_box(slide, "↓", x + 1.48, 3.65, 0.65, 0.55, 30, color, True, PP_ALIGN.CENTER)
    text_box(slide, result, x + 0.35, 4.65, 2.95, 0.4, 18, INK, True, PP_ALIGN.CENTER)
footer(slide, 7)
notes(slide, "1:20", "It is useful to separate three operations. Inference estimates hidden states for fixed parameters. Learning estimates parameters from data, usually by repeatedly invoking inference. Prediction applies the fitted model to unseen or future observations. Many apparent disagreements between algorithms disappear once we ask which of these operations they approximate.")

# 8
slide = base_slide("Gradients tell us how to change parameters")
add_picture_contain(slide, EQ["gradient"], 1.0, 1.45, 11.3, 1.05)
rect(slide, 1.0, 3.0, 11.3, 2.5, PALE_BLUE, PALE_BLUE, True)
rich_lines(slide, [
    "∇θ L points toward the steepest local increase in loss",
    "η controls the step size",
    "Backpropagation applies the chain rule through a computation graph",
    "Automatic differentiation computes exact derivatives of the implemented program",
], 1.38, 3.42, 10.5, 1.72, 18, INK, 10)
footer(slide, 8)
notes(slide, "1:50", "A gradient is a local sensitivity. It tells us how much the objective would change under a small parameter change. Gradient descent moves in the negative-gradient direction. Backpropagation is not a separate learning principle; it is an efficient application of the chain rule. Automatic differentiation records the computation and evaluates those derivatives exactly up to numerical precision.")

# 9
slide = base_slide("Optimization does not remove modeling assumptions")
items = [
    ("Objective", "What quantity should improve?", BLUE),
    ("Parameterization", "Which functions can the model express?", TEAL),
    ("Data split", "What evidence tests generalization?", ORANGE),
    ("Identifiability", "Can different parameters explain the same data?", YELLOW),
]
for i, (title, sub, color) in enumerate(items):
    x = 0.85 + (i % 2) * 6.15; y = 1.55 + (i // 2) * 2.38
    rect(slide, x, y, 5.58, 1.75, WHITE, color, True)
    text_box(slide, title, x + 0.3, y + 0.33, 2.1, 0.36, 18, color, True)
    text_box(slide, sub, x + 0.3, y + 0.95, 4.9, 0.42, 15, MUTED)
footer(slide, 9)
notes(slide, "1:25", "Better optimization cannot rescue a poorly specified scientific model. We still need to choose an objective, a parameterization, and a validation strategy. Latent models also face identifiability: labels may permute, latent axes may rotate, and multiple parameter settings may induce the same observations. Interpretation therefore requires constraints and validation beyond training loss.")

section_slide("2", "Mixture models and EM", "The first nontrivial latent-variable learning problem", "Mixture models are the smallest setting where uncertainty over a hidden assignment affects parameter learning.")

# 11
slide = base_slide("A mixture model explains heterogeneity")
add_picture_contain(slide, EQ["gmm"], 0.9, 1.35, 11.5, 1.15)
for x, title, body, color in [
    (0.9, "πk", "component frequency", ORANGE),
    (4.68, "μk", "component center", TEAL),
    (8.46, "Σk", "component shape", BLUE),
]:
    rect(slide, x, 3.0, 3.45, 1.65, WHITE, color, True)
    text_box(slide, title, x + 0.2, 3.35, 3.05, 0.45, 23, color, True, PP_ALIGN.CENTER)
    text_box(slide, body, x + 0.2, 4.05, 3.05, 0.28, 14, MUTED, False, PP_ALIGN.CENTER)
text_box(slide, "Hidden assignment zn selects the component that generated xn.", 1.2, 5.55, 10.9, 0.45, 20, INK, True, PP_ALIGN.CENTER)
footer(slide, 11)
notes(slide, "1:40", "A Gaussian mixture assumes that each observation came from one of K components. The latent variable z is discrete. The mixture weights describe component prevalence; each component has a mean and covariance. The marginal density is a weighted sum because the generating component is unobserved.")

# 12
slide = base_slide("EM alternates posterior inference and parameter learning")
for x, label, title, body, color in [
    (1.0, "E", "Expectation step", "Compute responsibilities\nrnk = p(zn = k | xn, θ)", TEAL),
    (7.1, "M", "Maximization step", "Refit π, μ, Σ using\nposterior-weighted statistics", ORANGE),
]:
    rect(slide, x, 1.75, 5.25, 3.5, WHITE, color, True)
    text_box(slide, label, x + 0.3, 2.08, 0.75, 0.75, 36, color, True, PP_ALIGN.CENTER)
    text_box(slide, title, x + 1.25, 2.2, 3.55, 0.42, 20, INK, True)
    text_box(slide, body, x + 0.55, 3.38, 4.15, 0.95, 17, MUTED, False, PP_ALIGN.CENTER)
text_box(slide, "↔", 6.25, 3.0, 0.8, 0.5, 32, INK, True, PP_ALIGN.CENTER)
text_box(slide, "Each full iteration does not decrease the data log likelihood.", 1.2, 6.0, 10.9, 0.38, 17, INK, True, PP_ALIGN.CENTER)
footer(slide, 12)
notes(slide, "1:55", "If assignments were known, fitting each Gaussian would be easy. If parameters were known, computing assignment probabilities would also be easy. EM exploits this conditional simplicity. The E-step computes soft assignments. The M-step treats those responsibilities as fractional counts. The process monotonically improves or preserves the observed-data likelihood, although it can converge to a local optimum.")

image_slide(
    "Soft assignments preserve uncertainty",
    FIGURES["responsibilities"],
    "Color mixtures represent posterior responsibilities; the right panel shows assignment confidence.",
    "A hard clustering algorithm collapses uncertainty immediately. EM keeps a probability for every component. Points near overlapping regions receive mixed colors and lower confidence. This uncertainty is not cosmetic; it determines how strongly each observation contributes to each component's sufficient statistics.",
    "Source: executed notebook 02b_mixtures_em.ipynb",
    "1:30",
)

# 14
slide = base_slide("EM is coordinate ascent on a lower bound")
add_picture_contain(slide, EQ["elbo"], 0.65, 1.35, 12.0, 1.1)
for x, label, body, color in [
    (0.95, "ELBO", "tractable objective", TEAL),
    (4.95, "+", "", MUTED),
    (6.15, "KL gap", "posterior approximation error", ORANGE),
]:
    if label == "+":
        text_box(slide, label, x, 3.38, 0.5, 0.5, 32, color, True, PP_ALIGN.CENTER)
    else:
        rect(slide, x, 3.0, 4.15, 1.65, WHITE, color, True)
        text_box(slide, label, x + 0.25, 3.33, 3.65, 0.4, 20, color, True, PP_ALIGN.CENTER)
        text_box(slide, body, x + 0.25, 4.02, 3.65, 0.28, 14, MUTED, False, PP_ALIGN.CENTER)
text_box(slide, "Exact E-step: q(z) = p(z | x, θ), so the KL gap becomes zero.", 1.0, 5.65, 11.3, 0.42, 18, INK, True, PP_ALIGN.CENTER)
footer(slide, 14)
notes(slide, "1:55", "The ELBO identity is the bridge from EM to variational inference. The log evidence equals a lower bound plus a nonnegative KL divergence. In ordinary EM, the E-step can set q equal to the exact posterior, making the gap zero. The M-step then optimizes the bound with respect to parameters. Variational inference keeps the same logic when the exact posterior is unavailable.")

image_slide(
    "EM produces a fitted density, not only labels",
    FIGURES["mixture"],
    "Left: recovered Gaussian geometry. Right: monotonic observed-data log likelihood.",
    "The scientific output of a mixture model is more than a colored partition. It is a probability density with component frequencies, centers, and uncertainty geometry. The likelihood trace is a basic implementation check, but held-out density and stability across initializations are stronger validation criteria.",
    "Source: executed notebook 02b_mixtures_em.ipynb",
    "1:25",
)

# 16
slide = base_slide("What a mixture model cannot represent")
text_box(slide, "Exchangeability", 0.9, 1.45, 3.1, 0.5, 23, ORANGE, True)
text_box(slide, "The model gives the same probability after\nany permutation of observations.", 0.92, 2.15, 4.9, 0.85, 16, INK)
for i in range(7):
    x = 6.55 + i * 0.78
    c = [BLUE, BLUE, TEAL, ORANGE, ORANGE, TEAL, BLUE][i]
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.1), Inches(0.48), Inches(0.48))
    shape.fill.solid(); shape.fill.fore_color.rgb = c; shape.line.color.rgb = c
text_box(slide, "same histogram", 8.15, 3.0, 2.7, 0.35, 15, MUTED, False, PP_ALIGN.CENTER)
text_box(slide, "≠", 6.2, 4.15, 0.5, 0.5, 30, INK, True, PP_ALIGN.CENTER)
for i in range(7):
    x = 6.55 + i * 0.78
    c = [BLUE, TEAL, BLUE, ORANGE, TEAL, ORANGE, BLUE][i]
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(4.22), Inches(0.48), Inches(0.48))
    shape.fill.solid(); shape.fill.fore_color.rgb = c; shape.line.color.rgb = c
text_box(slide, "different sequence", 8.05, 5.12, 2.9, 0.35, 15, MUTED, False, PP_ALIGN.CENTER)
text_box(slide, "To model persistence, switching, and prediction, z must depend on its history.", 0.9, 5.75, 11.5, 0.45, 20, TEAL, True, PP_ALIGN.CENTER)
footer(slide, 16)
notes(slide, "1:20", "Mixture observations are exchangeable. The model sees component counts but not order. Two sequences with the same histogram are equivalent. Neural and behavioral data contain persistence, switching, and temporal prediction. The minimal extension is to connect consecutive latent assignments with a Markov transition model.")

section_slide("3", "Hidden Markov models", "Discrete latent states with temporal persistence", "An HMM changes one assumption: the hidden assignment now depends on the previous hidden state.")

# 18
slide = base_slide("An HMM adds a Markov chain to a mixture")
add_picture_contain(slide, EQ["hmm"], 0.55, 1.2, 12.2, 1.1)
for t in range(5):
    x = 1.25 + t * 2.2
    z = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(3.0), Inches(0.74), Inches(0.74))
    z.fill.solid(); z.fill.fore_color.rgb = TEAL; z.line.color.rgb = TEAL
    text_box(slide, f"z{t+1}", x + 0.08, 3.19, 0.58, 0.25, 14, WHITE, True, PP_ALIGN.CENTER)
    obs = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.75), Inches(0.74), Inches(0.74))
    obs.fill.solid(); obs.fill.fore_color.rgb = BLUE; obs.line.color.rgb = BLUE
    text_box(slide, f"x{t+1}", x + 0.08, 4.94, 0.58, 0.25, 14, WHITE, True, PP_ALIGN.CENTER)
    text_box(slide, "↓", x + 0.17, 4.0, 0.4, 0.36, 20, MUTED, True, PP_ALIGN.CENTER)
    if t < 4:
        text_box(slide, "→", x + 1.1, 3.18, 0.55, 0.3, 20, MUTED, True, PP_ALIGN.CENTER)
text_box(slide, "transition model", 4.4, 2.5, 4.0, 0.3, 13, TEAL, True, PP_ALIGN.CENTER)
text_box(slide, "emission model", 4.55, 5.72, 3.7, 0.3, 13, BLUE, True, PP_ALIGN.CENTER)
footer(slide, 18)
notes(slide, "1:45", "The hidden state follows a first-order Markov chain: the next state depends on the current state, not the entire past. Each observation depends only on its current state. These conditional-independence assumptions are strong, but they enable exact dynamic programming rather than summing over K to the T possible state sequences.")

# 19
slide = base_slide("Dynamic programming makes exact sequence inference possible")
for x, title, equation, body, color in [
    (0.78, "Forward message", "αt(k) ∝ p(xt | k) Σj αt−1(j) Ajk", "past evidence", BLUE),
    (6.78, "Backward message", "βt(k) = Σj Akj p(xt+1 | j) βt+1(j)", "future evidence", ORANGE),
]:
    rect(slide, x, 1.55, 5.78, 3.2, WHITE, color, True)
    text_box(slide, title, x + 0.3, 1.95, 5.15, 0.38, 19, color, True, PP_ALIGN.CENTER)
    text_box(slide, equation, x + 0.4, 2.82, 4.95, 0.62, 16, INK, False, PP_ALIGN.CENTER)
    text_box(slide, body, x + 0.35, 3.9, 5.08, 0.35, 14, MUTED, False, PP_ALIGN.CENTER)
text_box(slide, "smoothed posterior ∝ forward × backward", 1.35, 5.55, 10.65, 0.48, 22, TEAL, True, PP_ALIGN.CENTER)
footer(slide, 19)
notes(slide, "2:00", "The forward message summarizes all evidence up to time t. The backward message summarizes evidence after time t. Their product gives the smoothed marginal posterior for each state. This is exact under the HMM assumptions and costs order T K squared, rather than enumerating exponentially many state paths. Scaling or log-space calculations are necessary for numerical stability.")

image_slide(
    "Filtering and smoothing answer different questions",
    FIGURES["hmm_smoothing"],
    "Filtering uses the past and present; smoothing revises state beliefs using the complete sequence.",
    "Filtering is the online posterior and is appropriate for real-time decoding. Smoothing is an offline posterior and can revise ambiguous moments using future observations. The distinction matters experimentally: an offline state estimate should not be reported as if it were available causally at that time.",
    "Source: executed notebook 03a_hmm_foundations.ipynb",
    "1:35",
)

# 21
two_column_slide(
    "Three posterior summaries serve different goals",
    "Marginals",
    ["p(zt | x1:T)", "Uncertainty at each time", "Expected occupancy and transitions"],
    "Paths",
    ["Viterbi: one MAP sequence", "Posterior sampling: distribution over sequences", "A single path hides ambiguity"],
    "HMM inference does not have one universal output. Marginals quantify uncertainty at each time. Viterbi decoding gives the single most probable joint sequence, which is not the same as choosing the most probable state independently at every time. Posterior path samples preserve sequence-level uncertainty and are often more honest for downstream analyses.",
    "1:35",
)

# 22
image_slide(
    "Baum-Welch is EM for an HMM",
    FIGURES["hmm_em"],
    "Forward-backward supplies expected state occupancies and transition counts for the M-step.",
    "Baum-Welch repeats exact HMM inference and parameter updates. The E-step computes smoothed state probabilities and pairwise transition marginals. The M-step updates the initial distribution, transition matrix, and emission parameters. As with mixture EM, labels may permute and local optima remain possible, so recovery is judged after label alignment and across initializations.",
    "Source: executed notebook 03a_hmm_foundations.ipynb",
    "1:45",
)

# 23
slide = base_slide("HMM assumptions determine the scientific interpretation")
rows = [
    ("Gaussian HMM", "state-specific mean/covariance", "recurring static regimes"),
    ("AR-HMM", "state-specific autoregression", "switching local dynamics"),
    ("SLDS", "continuous latent state + switches", "piecewise linear dynamics"),
]
text_box(slide, "Model", 0.85, 1.55, 2.0, 0.3, 13, MUTED, True)
text_box(slide, "Observation / dynamics", 3.25, 1.55, 3.5, 0.3, 13, MUTED, True)
text_box(slide, "Interpretation", 8.05, 1.55, 3.8, 0.3, 13, MUTED, True)
for i, (model, obs, interp) in enumerate(rows):
    y = 2.05 + i * 1.35
    fill = WHITE if i % 2 == 0 else RGBColor(239, 243, 241)
    rect(slide, 0.75, y, 11.85, 1.08, fill, fill)
    text_box(slide, model, 1.0, y + 0.29, 1.9, 0.35, 17, TEAL, True)
    text_box(slide, obs, 3.25, y + 0.29, 4.2, 0.35, 16, INK)
    text_box(slide, interp, 8.05, y + 0.29, 3.9, 0.35, 16, INK)
footer(slide, 23)
notes(slide, "1:30", "A state label means only what the emission and transition assumptions make it mean. In a Gaussian HMM, states represent recurring observation distributions. In an AR-HMM, each state defines local autoregressive dynamics. In an SLDS, a discrete switch selects continuous linear dynamics. This hierarchy is useful, but increasing expressiveness also weakens identifiability and raises computational cost.")

section_slide("4", "Variational inference and VAEs", "Approximate posteriors for nonlinear latent-variable models", "Exact inference breaks when the decoder becomes nonlinear. Variational inference replaces the intractable posterior with an optimized approximation.")

# 25
slide = base_slide("Variational inference turns inference into optimization")
add_picture_contain(slide, EQ["elbo"], 0.65, 1.2, 12.0, 1.05)
for i, (label, body, color) in enumerate([
    ("Choose qφ(z | x)", "a tractable posterior family", BLUE),
    ("Maximize the ELBO", "fit φ and model parameters θ", TEAL),
    ("Check the gap", "approximation quality is not guaranteed", ORANGE),
]):
    x = 0.75 + i * 4.2
    rect(slide, x, 3.0, 3.85, 2.1, WHITE, color, True)
    text_box(slide, label, x + 0.2, 3.38, 3.45, 0.42, 18, color, True, PP_ALIGN.CENTER)
    text_box(slide, body, x + 0.3, 4.25, 3.25, 0.55, 14, MUTED, False, PP_ALIGN.CENTER)
footer(slide, 25)
notes(slide, "1:45", "Variational inference chooses a tractable family q and minimizes its KL divergence to the posterior indirectly by maximizing the ELBO. The approximation can be factorized, Gaussian, sequential, or represented by a neural network. A high ELBO does not prove that q captures every posterior dependency; it is an optimization target and a lower bound, not an automatic certificate of scientific validity.")

# 26
slide = base_slide("Amortization learns a reusable inference rule")
for x, title, lines, color in [
    (0.8, "Classical variational inference", ["new local parameters for every xn", "optimize separately", "accurate but expensive"], ORANGE),
    (6.78, "Amortized inference", ["shared encoder qφ(z | x)", "one forward pass", "fast but may introduce amortization gap"], TEAL),
]:
    rect(slide, x, 1.5, 5.75, 4.85, WHITE, color, True)
    text_box(slide, title, x + 0.35, 1.88, 5.05, 0.45, 19, color, True, PP_ALIGN.CENTER)
    rich_lines(slide, lines, x + 0.55, 2.85, 4.7, 2.55, 17, INK, 18)
footer(slide, 26)
notes(slide, "1:35", "Classical variational methods optimize local variational parameters separately for every observation. A VAE amortizes that work: an encoder predicts the variational parameters from x. This enables fast inference and generalization to new data, but the shared encoder may fail to reach the best variational parameters for every datapoint, creating an amortization gap.")

# 27
slide = base_slide("The VAE objective balances reconstruction and regularization")
add_picture_contain(slide, EQ["vae"], 0.5, 1.25, 12.3, 1.1)
for x, title, body, color in [
    (1.0, "Reconstruction", "Does z explain x through the decoder?", BLUE),
    (7.05, "KL regularization", "Does the inferred code remain close to the prior?", ORANGE),
]:
    rect(slide, x, 3.0, 5.25, 2.1, WHITE, color, True)
    text_box(slide, title, x + 0.3, 3.4, 4.65, 0.4, 19, color, True, PP_ALIGN.CENTER)
    text_box(slide, body, x + 0.45, 4.25, 4.35, 0.55, 15, MUTED, False, PP_ALIGN.CENTER)
text_box(slide, "The balance controls information capacity and can cause posterior collapse.", 1.0, 5.8, 11.3, 0.38, 17, INK, True, PP_ALIGN.CENTER)
footer(slide, 27)
notes(slide, "1:50", "The first ELBO term rewards reconstructions under latent samples from the encoder. The KL term regularizes the approximate posterior toward the prior. Too little regularization gives an unstructured code; too much can make the decoder ignore z, known as posterior collapse. This is especially important with powerful sequential decoders.")

image_slide(
    "A VAE combines a generative model and an inference network",
    FIGURES["vae_arch"],
    "The encoder parameterizes qφ(z | x); the decoder parameterizes pθ(x | z).",
    "The decoder is the generative model: latent z generates observations. The encoder is an auxiliary inference model: it maps each observation to a distribution over z. Keeping these roles separate clarifies why a VAE is not simply an autoencoder with noise. The probabilistic objective and latent prior are essential.",
    "Source: executed notebook 04_standard_vae.ipynb",
    "1:35",
)

# 29
slide = base_slide("The reparameterization trick differentiates through sampling")
add_picture_contain(slide, EQ["reparam"], 0.55, 1.35, 12.2, 1.1)
for x, label, sub, color in [
    (1.1, "encoder", "μ(x), σ(x)", BLUE),
    (5.15, "fixed noise", "ε ~ N(0, I)", ORANGE),
    (9.2, "sample", "differentiable z", TEAL),
]:
    rect(slide, x, 3.15, 3.0, 1.55, WHITE, color, True)
    text_box(slide, label, x + 0.2, 3.48, 2.6, 0.35, 18, color, True, PP_ALIGN.CENTER)
    text_box(slide, sub, x + 0.2, 4.05, 2.6, 0.3, 14, MUTED, False, PP_ALIGN.CENTER)
    if x < 9:
        text_box(slide, "+", x + 3.32, 3.62, 0.45, 0.45, 25, INK, True, PP_ALIGN.CENTER)
text_box(slide, "Randomness is moved into ε, leaving a differentiable deterministic transformation.", 1.0, 5.65, 11.3, 0.42, 18, INK, True, PP_ALIGN.CENTER)
footer(slide, 29)
notes(slide, "1:45", "Direct sampling appears to interrupt the computation graph. Reparameterization writes a Gaussian sample as a deterministic function of encoder outputs and parameter-free noise. Gradients then flow through mu and sigma. This produces a low-variance pathwise gradient estimator and makes end-to-end VAE training practical.")

image_slide(
    "A latent space is useful only when validated",
    FIGURES["vae_latent"],
    "A two-dimensional code can separate structure, but orientation and scale are generally not identifiable.",
    "A clean latent visualization is suggestive, not sufficient. We should test held-out likelihood or reconstruction, stability across seeds, sensitivity to latent dimension and prior, and correspondence with external variables that were not used to train the model. Latent axes can rotate or permute without changing the observation model.",
    "Source: executed notebook 04_standard_vae.ipynb",
    "1:30",
)

section_slide("5", "LFADS", "A sequential VAE for latent neural dynamics", "LFADS extends the VAE from independent observations to entire spike-train sequences generated by a recurrent dynamical system.")

# 32
image_slide(
    "LFADS separates inferred initial condition, dynamics, and rates",
    FIGURES["lfads_arch"],
    "An encoder infers a distribution over initial generator state; recurrent dynamics produce factors and firing rates.",
    "LFADS treats a trial as a sequence generated by a dynamical system. The encoder summarizes the observed spikes and infers an initial condition. A recurrent generator evolves that state. A low-dimensional factor readout drives nonnegative firing rates, and spikes are modeled with a Poisson likelihood. More complete variants also infer time-varying inputs.",
    "Source: executed notebook 06_lfads.ipynb",
    "1:55",
)

# 33
slide = base_slide("The sequential ELBO assigns uncertainty to a trajectory")
for x, title, body, color in [
    (0.75, "Poisson reconstruction", "Do inferred rates explain\nheld-out spikes?", BLUE),
    (4.75, "Initial-state KL", "How much trial-specific\ninformation enters g0?", TEAL),
    (8.75, "Input KL", "Are inferred perturbations\nnecessary and plausible?", ORANGE),
]:
    rect(slide, x, 1.7, 3.8, 3.4, WHITE, color, True)
    text_box(slide, title, x + 0.25, 2.08, 3.3, 0.55, 18, color, True, PP_ALIGN.CENTER)
    text_box(slide, body, x + 0.35, 3.17, 3.1, 1.05, 13, MUTED, False, PP_ALIGN.CENTER)
text_box(slide, "Scientific danger: inferred inputs can absorb model mismatch.", 1.1, 5.8, 11.1, 0.4, 19, ORANGE, True, PP_ALIGN.CENTER)
footer(slide, 33)
notes(slide, "1:50", "The LFADS objective is a sequential ELBO. The likelihood evaluates whether inferred rates explain spike counts. KL terms constrain trial-specific initial states and, when present, inferred controller inputs. Inputs are powerful but scientifically delicate: they may represent real unobserved perturbations, or they may compensate for inadequate autonomous dynamics. Their interpretation requires perturbation timing and held-out validation.")

image_slide(
    "Learned dynamics must generalize beyond training trajectories",
    FIGURES["lfads_dynamics"],
    "Compare learned and true vector fields, not only reconstructed observations.",
    "For synthetic data, we can compare the learned vector field with the known system. For real data, we instead rely on held-out spike prediction, consistency across trials, perturbation responses, and dynamical invariants. Reconstruction alone can be achieved by a flexible encoder-decoder without recovering the correct autonomous dynamics.",
    "Source: executed notebook 06_lfads.ipynb",
    "1:40",
)

section_slide("6", "gpSLDS", "Nonlinear dynamics assembled from interpretable local regimes", "The final model retains nonlinear flexibility while exposing local linear dynamics and uncertainty.")

# 36
image_slide(
    "gpSLDS: nonlinear dynamics with local regimes",
    FIGURES["gpslds_arch"],
    "A smoothly switching kernel blends local linear dynamics within a global nonlinear field.",
    "A Gaussian-process switching linear dynamical system places a GP prior on the vector field. Local linear regimes define interpretable dynamical templates, while smooth gating functions blend them. Inducing points make GP inference scalable. The result is a nonlinear field with uncertainty and a local linear decomposition.",
    "Source: executed notebook 08_gpslds.ipynb; model: Hu et al., NeurIPS 2024",
    "1:55",
)

# 37
image_slide(
    "Local linear structure makes nonlinear dynamics inspectable",
    FIGURES["gpslds_local"],
    "Regime weights partition latent space smoothly; local Jacobians support fixed-point and stability analysis.",
    "The key scientific advantage is inspectability. At each location we can examine regime weights, local linear operators, fixed points, and uncertainty in the vector field. However, this interpretation still depends on the latent observation model and identifiability constraints. A smooth partition is not automatically a biological module.",
    "Source: executed notebook 08_gpslds.ipynb",
    "1:40",
)

# 38
slide = base_slide("Choose the simplest model that answers the scientific question")
headers = ["Model", "Latent structure", "Inference", "Best use"]
xs = [0.55, 2.1, 5.05, 8.25]
ws = [1.45, 2.85, 3.05, 4.45]
for x, w, h in zip(xs, ws, headers):
    rect(slide, x, 1.35, w, 0.55, INK, INK)
    text_box(slide, h, x + 0.08, 1.5, w - 0.16, 0.25, 12, WHITE, True, PP_ALIGN.CENTER)
rows = [
    ("GMM", "discrete class", "exact EM", "heterogeneous populations"),
    ("HMM", "discrete sequence", "forward-backward", "recurring temporal regimes"),
    ("VAE", "continuous code", "amortized VI", "nonlinear representation"),
    ("LFADS", "recurrent trajectory", "sequential VI", "population spike dynamics"),
    ("gpSLDS", "continuous + local regimes", "variational GP", "interpretable nonlinear flow"),
]
for i, row in enumerate(rows):
    y = 1.98 + i * 0.91
    fill = WHITE if i % 2 == 0 else RGBColor(235, 240, 238)
    for x, w, val in zip(xs, ws, row):
        rect(slide, x, y, w, 0.78, fill, fill)
        text_box(slide, val, x + 0.1, y + 0.2, w - 0.2, 0.33, 13, INK, i == 0 and x == xs[0], PP_ALIGN.CENTER)
footer(slide, 38)
notes(slide, "2:00", "Model choice should follow the scientific question, not a hierarchy of sophistication. Use a mixture when temporal order is irrelevant. Use an HMM for discrete recurring regimes. Use a VAE for nonlinear continuous representations. Use LFADS when trial-level recurrent dynamics are central. Use gpSLDS when nonlinear flow and local dynamical interpretation are both required. More flexible models demand stronger validation.")

# 39
slide = base_slide("A practical validation ladder")
steps = [
    ("1", "Implementation", "synthetic recovery, invariants, numerical stability", BLUE),
    ("2", "Predictive", "held-out likelihood, spikes, trajectories", TEAL),
    ("3", "Robustness", "seeds, hyperparameters, latent dimension", ORANGE),
    ("4", "Scientific", "external variables, perturbations, falsifiable predictions", YELLOW),
]
for i, (num, title, body, color) in enumerate(steps):
    y = 1.35 + i * 1.28
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(y), Inches(0.72), Inches(0.72))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.color.rgb = color
    text_box(slide, num, 1.04, y + 0.18, 0.44, 0.25, 15, WHITE, True, PP_ALIGN.CENTER)
    text_box(slide, title, 2.0, y + 0.03, 2.4, 0.35, 18, INK, True)
    text_box(slide, body, 4.55, y + 0.03, 7.4, 0.48, 16, MUTED)
footer(slide, 39)
notes(slide, "1:45", "Validation should proceed in layers. First verify the implementation using synthetic data and mathematical invariants. Then test predictive generalization. Next test robustness across seeds and modeling choices. Finally ask whether latent variables correspond to external measurements or produce falsifiable perturbation predictions. A compelling visualization without these layers is exploratory evidence, not a mechanistic result.")

# 40
slide = base_slide("Take-home message", dark=True)
messages = [
    ("1", "Inference is the computation that connects a generative hypothesis to data."),
    ("2", "EM, forward-backward, and VAEs share the same posterior-learning logic."),
    ("3", "Temporal structure changes both the algorithm and the scientific meaning of a latent state."),
    ("4", "Interpretability comes from constrained assumptions plus validation, not from low dimension alone."),
]
for i, (num, body) in enumerate(messages):
    y = 1.45 + i * 1.18
    text_box(slide, num, 0.9, y, 0.6, 0.48, 22, YELLOW, True, PP_ALIGN.CENTER)
    text_box(slide, body, 1.75, y, 10.4, 0.58, 17, WHITE, i == 3)
text_box(slide, "Variational Neural Inference", 0.9, 6.55, 5.2, 0.35, 13, RGBColor(173, 190, 191), True)
footer(slide, 40, "Tutorial notebooks and references are available in this repository.", dark=True)
notes(slide, "1:35", "The unifying idea is that a generative model makes assumptions explicit, and inference connects those assumptions to observations. EM, HMM inference, VAEs, and sequential latent models are not isolated tricks; they are increasingly expressive solutions to the same posterior-learning problem. The value of the latent state comes from what the model constrains and what independent evidence validates.")


for index, slide in enumerate(prs.slides, start=1):
    if index not in {1} and not any(shape.has_text_frame and shape.text == f"{index:02d}" for shape in slide.shapes):
        # Section slides intentionally omit a footer during construction; add a subtle number.
        if slide.background.fill.fore_color.rgb == INK:
            footer(slide, index, dark=True)

output_pptx = OUT / "Variational_Neural_Inference_Group_Talk.pptx"
prs.save(output_pptx)

notes_path = OUT / "Variational_Neural_Inference_Speaker_Notes.md"


def slide_title(slide, fallback):
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame or not shape.text.strip():
            continue
        sizes = [
            size.pt
            for paragraph in shape.text_frame.paragraphs
            for size in [paragraph.font.size, *(run.font.size for run in paragraph.runs)]
            if size is not None
        ]
        candidates.append((max(sizes, default=0), -shape.top, shape.text.strip()))
    return max(candidates, default=(0, 0, fallback))[2]


with notes_path.open("w", encoding="utf-8") as handle:
    handle.write("# Variational Neural Inference: Speaker Notes\n\n")
    handle.write("Target duration: 60 minutes. All timing cues are approximate.\n\n")
    for index, slide in enumerate(prs.slides, start=1):
        title = slide_title(slide, f"Slide {index}")
        note_text = slide.notes_slide.notes_text_frame.text
        handle.write(f"## {index}. {title}\n\n{note_text}\n\n")

print(f"Built {len(prs.slides)} slides: {output_pptx}")
print(f"Speaker notes: {notes_path}")
