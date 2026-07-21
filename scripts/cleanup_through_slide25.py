from pathlib import Path
import shutil

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_parts2_3 import (
    title, footer, rect, text,
    GREEN, GREEN_DARK, GREEN_LIGHT, BLUE, BLUE_LIGHT, ORANGE,
    ORANGE_LIGHT, CHARCOAL, MID, LIGHT, WHITE,
)


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_working.pptx"
OUTPUT = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_working_11.pptx"


def delete_slide(prs, index):
    slide_id = prs.slides._sldIdLst[index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    prs.slides._sldIdLst.remove(slide_id)


def replace_text(slide, old_start, new):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip().startswith(old_start):
            tf = shape.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = new
                for run in tf.paragraphs[0].runs[1:]:
                    run.text = ""
            else:
                tf.text = new
            return True
    return False


def add_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    new_id = prs.slides._sldIdLst[-1]
    rect(slide, 0, 0, 13.333, 7.5, WHITE, None, False)
    title(slide, "Five models, five assumptions about the hidden state", 25, "SECTION SUMMARY")

    headers = ["MODEL", "LATENT STRUCTURE", "OBSERVATION MODEL", "SCIENTIFIC QUESTION"]
    x = [.55, 2.35, 6.18, 9.12]
    w = [1.72, 3.75, 2.86, 3.63]
    for xi, wi, value in zip(x, w, headers):
        rect(slide, xi, 1.20, wi, .56, GREEN, None, False)
        text(slide, value, xi+.08, 1.31, wi-.16, .29, 12, WHITE, True, PP_ALIGN.CENTER)

    rows = [
        ("Mixture", "independent discrete assignment", "feature distribution", "Which source generated it?", BLUE),
        ("HMM", "temporally connected discrete state", "state-dependent likelihood", "Which regime are we in?", ORANGE),
        ("LDS", "continuous linear trajectory", "Gaussian", "Where is the system moving?", GREEN),
        ("SLDS", "discrete regime + continuous trajectory", "flexible", "Which local dynamics apply?", ORANGE),
        ("PLDS", "continuous linear trajectory", "Poisson spikes", "What dynamics generated the spikes?", BLUE),
    ]
    for i, (model, latent, obs, question, accent) in enumerate(rows):
        y = 1.83 + i*.78
        fill = WHITE if i % 2 == 0 else LIGHT
        for xi, wi in zip(x, w):
            rect(slide, xi, y, wi, .68, fill, None, False)
        rect(slide, .55, y, .07, .68, accent, None, False)
        text(slide, model, .72, y+.14, 1.42, .32, 16, accent, True, PP_ALIGN.CENTER)
        text(slide, latent, 2.52, y+.10, 3.42, .42, 15, CHARCOAL, False, PP_ALIGN.CENTER)
        text(slide, obs, 6.35, y+.10, 2.52, .42, 15, CHARCOAL, False, PP_ALIGN.CENTER)
        text(slide, question, 9.28, y+.10, 3.30, .42, 15, CHARCOAL, False, PP_ALIGN.CENTER)

    rect(slide, 1.05, 6.15, 11.25, .62, GREEN_LIGHT, GREEN, True, .9)
    text(slide, "Choosing the latent structure is only half the problem—we still need to infer it from data.",
         1.30, 6.27, 10.75, .34, 18, GREEN_DARK, True, PP_ALIGN.CENTER)
    footer(slide, 25)
    return new_id


def renumber(prs):
    for page, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            value = shape.text.strip()
            if not value.isdigit():
                continue
            left = shape.left / Inches(1)
            top = shape.top / Inches(1)
            if (left < .85 and top < .55) or (left > 12.0 and top > 6.75):
                if shape.text_frame.paragraphs[0].runs:
                    shape.text_frame.paragraphs[0].runs[0].text = str(page)
                else:
                    shape.text = str(page)
                if shape.width < Inches(.42):
                    shape.width = Inches(.42)


def build():
    shutil.copy2(SOURCE, OUTPUT)
    prs = Presentation(OUTPUT)

    # Remove the duplicate PLDS slide and the two blank slides (physical 25–27).
    for index in (26, 25, 24):
        delete_slide(prs, index)

    # The retained PLDS slide is physical slide 24.
    plds = prs.slides[23]
    replace_text(plds, "λₜ = exp", "λₜₙ = exp(cₙᵀ zₜ + dₙ)")
    replace_text(
        plds,
        "PLDS replaces",
        "PLDS models population spikes as noisy observations of a continuous low-dimensional trajectory.",
    )

    new_id = add_summary(prs)
    slide_ids = prs.slides._sldIdLst
    slide_ids.remove(new_id)
    slide_ids.insert(24, new_id)
    renumber(prs)
    # Ensure the two-digit number renders fully on the newly inserted summary slide.
    summary = prs.slides[24]
    for shape in summary.shapes:
        if not getattr(shape, "has_text_frame", False) or shape.text.strip() != "25":
            continue
        shape.width = Inches(.72)
        if shape.top / Inches(1) < .6:
            shape.left = Inches(.18)
        elif shape.top / Inches(1) > 6.7:
            shape.left = Inches(12.22)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
