from pathlib import Path
import shutil
import base64
import json

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_parts2_3 import (
    title, footer, rect, text, formula, arrow, image,
    GREEN, GREEN_DARK, GREEN_LIGHT, BLUE, BLUE_LIGHT, ORANGE,
    ORANGE_LIGHT, CHARCOAL, MID, LIGHT, WHITE,
)


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_working.pptx"
OUTPUT = ROOT / "Presentation" / "Latent_Variable_Models_review_26_28.pptx"
ASSETS = ROOT / "Presentation" / "review_26_28_assets"
NBFIG = ASSETS / "03a_hmm_foundations_cell8.png"
ASSETS.mkdir(parents=True, exist_ok=True)


def delete_slide(prs, index):
    sid = prs.slides._sldIdLst[index]
    prs.part.drop_rel(sid.rId)
    prs.slides._sldIdLst.remove(sid)


def new_slide(prs, page, heading, section):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    rect(s, 0, 0, 13.333, 7.5, WHITE, None, False)
    title(s, heading, page, section)
    for shape in s.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip() == f"{page:02d}":
            shape.left = Inches(.18); shape.width = Inches(.68)
    return s


def make_poisson_posterior():
    notebook = ROOT / "VAE-Tutorial" / "Part2_Dynamics" / "03a_hmm_foundations.ipynb"
    nb = json.loads(notebook.read_text(encoding="utf-8"))
    for output in nb["cells"][8].get("outputs", []):
        png = output.get("data", {}).get("image/png")
        if png:
            NBFIG.write_bytes(base64.b64decode("".join(png) if isinstance(png, list) else png))
            break
    z = np.linspace(-3.2, 2.3, 700)
    prior = np.exp(-.5*(z/.95)**2); prior /= np.trapz(prior, z)
    y = 5
    rate = np.exp(z + .55)
    log_lik = y*np.log(rate) - rate
    likelihood = np.exp(log_lik-log_lik.max()); likelihood /= np.trapz(likelihood,z)
    posterior = prior*likelihood; posterior /= np.trapz(posterior,z)
    fig, ax = plt.subplots(figsize=(7.6, 4.2))
    ax.plot(z,prior,color="#3F78B5",lw=2.4,label="Gaussian prediction")
    ax.plot(z,likelihood,color="#EE772F",lw=2.4,label="Poisson likelihood")
    ax.fill_between(z,posterior,color="#B8DCCB",alpha=.65)
    ax.plot(z,posterior,color="#008A45",lw=3.0,label="posterior")
    ax.axvline(z[np.argmax(posterior)],color="#008A45",ls="--",lw=1.2)
    ax.set_xlabel("latent state  z"); ax.set_yticks([]); ax.set_title("Gaussian prediction × Poisson evidence",loc="left",weight="bold")
    ax.legend(frameon=False,ncol=3,loc="upper left"); ax.spines[["top","right","left"]].set_visible(False); ax.grid(axis="x",alpha=.10)
    fig.tight_layout(); fig.savefig(ASSETS/"poisson_posterior.png",dpi=230,facecolor="white"); plt.close(fig)


def build():
    make_poisson_posterior()
    shutil.copy2(SOURCE,OUTPUT); prs=Presentation(OUTPUT)
    while len(prs.slides)>25: delete_slide(prs,25)

    # 26: posterior uncertainty shown with a real notebook output.
    s=new_slide(prs,26,"Inference is a distribution over possible histories","INFERENCE BRIDGE")
    text(s,"The model does not simply return a colored state sequence.",.72,1.13,5.12,.45,20,CHARCOAL,True)
    formula(s,"p(zₜ | y₁:T)",.72,1.82,3.18,.72,25,GREEN_LIGHT,GREEN)
    text(s,"At every time point, it assigns probability\nto multiple latent explanations.",.82,2.82,4.52,.88,19,CHARCOAL,False)
    rect(s,.82,4.12,4.20,1.28,WHITE,BLUE,True,1.0)
    text(s,"Filtering",1.05,4.30,1.22,.30,16,BLUE,True)
    text(s,"uses evidence through the current time",2.24,4.29,2.52,.36,15,CHARCOAL,False)
    rect(s,.82,5.60,4.20,1.00,WHITE,ORANGE,True,1.0)
    text(s,"Smoothing",1.05,5.80,1.38,.30,16,ORANGE,True)
    text(s,"uses the complete sequence",2.45,5.79,2.28,.34,15,CHARCOAL,False)
    image(s,NBFIG,5.28,1.18,7.55,5.18)
    text(s,"Future observations can revise what we believe happened earlier.",6.12,6.38,5.82,.34,17,GREEN_DARK,True,PP_ALIGN.CENTER)
    footer(s,26,"Notebook figure: Part2_Dynamics/03a_hmm_foundations.ipynb · cell 8")

    # 27: two reasons exact inference closes algebraically.
    s=new_slide(prs,27,"Why exact inference works: two closure properties","INFERENCE BRIDGE")
    text(s,"FINITE-STATE HMM",.72,1.18,5.55,.34,17,BLUE,True,PP_ALIGN.CENTER)
    text(s,"LINEAR–GAUSSIAN LDS",7.05,1.18,5.55,.34,17,GREEN_DARK,True,PP_ALIGN.CENTER)
    formula(s,"αₜ(k) ∝ ℓₜ(k) Σⱼ Pⱼₖ αₜ₋₁(j)",.82,1.80,5.35,.82,22,BLUE_LIGHT,BLUE)
    text(s,"K possibilities",.90,3.04,1.68,.36,17,CHARCOAL,True,PP_ALIGN.CENTER)
    arrow(s,2.68,3.22,3.52,3.22,BLUE,2.0)
    text(s,"finite sum",3.62,3.04,1.48,.36,17,BLUE,True,PP_ALIGN.CENTER)
    arrow(s,5.13,3.22,5.72,3.22,BLUE,2.0)
    text(s,"new belief",4.96,3.62,1.12,.34,15,CHARCOAL,True,PP_ALIGN.CENTER)
    text(s,"Dynamic programming reuses the previous message;\nit never enumerates all Kᵀ sequences.",1.05,4.48,4.90,.90,18,CHARCOAL,False,PP_ALIGN.CENTER)
    formula(s,"Gaussian",7.33,1.80,2.00,.72,22,BLUE_LIGHT,BLUE)
    arrow(s,9.48,2.16,10.08,2.16,GREEN,2.0)
    formula(s,"linear map",10.22,1.80,2.00,.72,21,GREEN_LIGHT,GREEN)
    arrow(s,11.22,2.70,11.22,3.12,GREEN,2.0)
    formula(s,"Gaussian",10.22,3.27,2.00,.72,22,BLUE_LIGHT,BLUE)
    text(s,"Prediction and correction remain Gaussian,\nso mean and covariance are sufficient.",7.45,4.48,4.82,.90,18,CHARCOAL,False,PP_ALIGN.CENTER)
    line=rect(s,6.60,1.28,.025,4.62,LIGHT,None,False)
    rect(s,1.08,6.12,11.15,.58,LIGHT,MID,True,.7)
    text(s,"Exactness is a consequence of finite sums or Gaussian closure—not a generic property of latent models.",1.33,6.24,10.65,.32,17,CHARCOAL,True,PP_ALIGN.CENTER)
    footer(s,27,"Forward recursion follows Part2_Dynamics/03a_hmm_foundations.ipynb")

    # 28: a concrete likelihood change that destroys Kalman closure.
    s=new_slide(prs,28,"Poisson spikes break Kalman's Gaussian closure","INFERENCE BRIDGE")
    formula(s,"prediction:  p(zₜ | y₁:ₜ₋₁) = Gaussian",.62,1.20,5.18,.72,21,BLUE_LIGHT,BLUE)
    formula(s,"spikes:  yₜₙ ~ Poisson(exp(cₙᵀzₜ+dₙ))",.62,2.13,5.18,.82,18,ORANGE_LIGHT,ORANGE)
    text(s,"Bayes update",1.72,3.43,2.60,.34,17,CHARCOAL,True,PP_ALIGN.CENTER)
    formula(s,"posterior ∝ prediction × likelihood",.82,3.95,4.78,.72,20,LIGHT,None)
    text(s,"The product is generally not Gaussian.\nA mean and covariance no longer describe it exactly.",.82,5.05,4.78,.82,18,CHARCOAL,False,PP_ALIGN.CENTER)
    image(s,ASSETS/"poisson_posterior.png",6.02,1.15,6.52,4.55)
    rect(s,6.47,5.78,5.62,.58,ORANGE_LIGHT,ORANGE,True,.8)
    text(s,"Non-Gaussian does not mean impossible—only no closed-form Kalman update.",6.73,5.90,5.10,.32,16,ORANGE,True,PP_ALIGN.CENTER)
    rect(s,1.10,6.45,11.10,.42,GREEN_LIGHT,GREEN,True,.8)
    text(s,"This is the opening for approximate inference.",1.35,6.51,10.60,.26,17,GREEN_DARK,True,PP_ALIGN.CENTER)
    footer(s,28,"PLDS observation model recalled from slide 24")

    prs.save(OUTPUT); print(OUTPUT)


if __name__=="__main__": build()
