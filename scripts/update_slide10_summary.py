from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_working.pptx"
OUTPUT = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_working_03.pptx"

GREEN = RGBColor(0x00, 0x8A, 0x45)
GREEN_DARK = RGBColor(0x00, 0x68, 0x36)
GREEN_LIGHT = RGBColor(0xE8, 0xF4, 0xEE)
BLUE = RGBColor(0x3F, 0x78, 0xB5)
CHARCOAL = RGBColor(0x33, 0x35, 0x37)
MID = RGBColor(0x68, 0x6D, 0x71)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def text_box(slide, text, x, y, w, h, size, color, bold=False,
             align=PP_ALIGN.CENTER, font="Arial"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def formula_card(slide, label, formula, x, accent):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.96), Inches(4.45), Inches(0.84)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = GREEN_LIGHT
    card.line.color.rgb = accent
    card.line.width = Pt(1.0)
    try:
        card.adjustments[0] = 0.06
    except Exception:
        pass
    text_box(slide, label, x + 0.18, 6.03, 1.45, 0.20, 8.5, accent, True, PP_ALIGN.LEFT)
    text_box(slide, formula, x + 1.50, 6.04, 2.70, 0.54, 20, CHARCOAL, False,
             PP_ALIGN.CENTER, "Cambria Math")


def main():
    shutil.copy2(SOURCE, OUTPUT)
    prs = Presentation(OUTPUT)
    slide = prs.slides[9]

    # Remove the previous bottom takeaway elements; they otherwise remain visible
    # beneath the new formula cards because the source slide contains two layers.
    for shape in list(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        value = shape.text.strip().lower()
        if value.startswith("neural dynamics") or value.startswith("different signals"):
            shape._element.getparent().remove(shape._element)

    # Cover the old slogan bar before inserting the common state-space formulation.
    cover = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.95), Inches(5.88), Inches(9.55), Inches(1.00)
    )
    cover.fill.solid()
    cover.fill.fore_color.rgb = WHITE
    cover.line.fill.background()

    formula_card(slide, "OBSERVATION MODEL", "yₜ ∼ p(yₜ | zₜ)", 2.03, BLUE)
    formula_card(slide, "DYNAMICS MODEL", "zₜ ∼ p(zₜ | zₜ₋₁)", 6.86, GREEN)

    # Correct the copied slide number without relying on template-specific shape names.
    number_cover = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(12.38), Inches(7.00), Inches(0.48), Inches(0.30)
    )
    number_cover.fill.solid()
    number_cover.fill.fore_color.rgb = WHITE
    number_cover.line.fill.background()
    text_box(slide, "10", 12.45, 7.04, 0.30, 0.18, 8, MID, False, PP_ALIGN.RIGHT)

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
