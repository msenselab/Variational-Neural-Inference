"""Build the LMU-styled latent neural dynamics presentation from the supplied template."""

from __future__ import annotations

import math
import shutil
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from matplotlib.colors import LinearSegmentedColormap
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = Path(r"C:\Users\10993\Desktop\Latent Variable Models for Neural Data and Dynamics .pptx")
OUTPUT = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_working_02.pptx"
ASSETS = ROOT / "Presentation" / "assets_v2"

GREEN = RGBColor(0x00, 0x8A, 0x45)
GREEN_DARK = RGBColor(0x00, 0x61, 0x32)
GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xEE)
LIME = RGBColor(0x6D, 0xAA, 0x2C)
CHARCOAL = RGBColor(0x31, 0x33, 0x36)
MID = RGBColor(0x66, 0x6A, 0x70)
LIGHT = RGBColor(0xF4, 0xF6, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xE8, 0x7B, 0x35)
BLUE = RGBColor(0x3F, 0x78, 0xB5)
RED = RGBColor(0xC9, 0x45, 0x45)

FONT = "Arial"
MATH_FONT = "Cambria Math"
SLIDE_W = 13.333
SLIDE_H = 7.5


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def clear_slide(slide, keep_top_line: bool = True) -> None:
    for shape in list(slide.shapes):
        if keep_top_line and shape.shape_type == 9 and shape.top < Inches(0.9):
            continue
        remove_shape(shape)


def set_run(run, size=20, bold=False, color=CHARCOAL, font=FONT, italic=False) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=20,
    bold=False,
    color=CHARCOAL,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.04,
    font=FONT,
    italic=False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = text
    set_run(r, size, bold, color, font, italic)
    return box


def add_rich_lines(slide, lines, x, y, w, h, size=18, color=CHARCOAL, gap=7):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, bold, item_color = item, False, color
        else:
            text, bold, item_color = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        r = p.add_run()
        r.text = text
        set_run(r, size, bold, item_color)
    return box


def add_bullets(slide, bullets, x, y, w, h, size=19, color=CHARCOAL, gap=9, level0_color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, bullet in enumerate(bullets):
        if isinstance(bullet, tuple):
            text, level = bullet
        else:
            text, level = bullet, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # A literal glyph is portable across python-pptx versions and remains editable.
        p.text = f"{'   ' * level}• {text}"
        p.level = level
        p.space_after = Pt(gap if level == 0 else max(3, gap - 4))
        p.line_spacing = 1.05
        p.font.name = FONT
        p.font.size = Pt(size if level == 0 else size - 2)
        p.font.color.rgb = level0_color if level == 0 and level0_color else color
    return box


def add_rect(slide, x, y, w, h, fill=LIGHT, line=GREEN, radius=True, width=1.2):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(width)
    if radius:
        try:
            sh.adjustments[0] = 0.08
        except Exception:
            pass
    return sh


def add_card(slide, title, body, x, y, w, h, accent=GREEN, icon=None, body_size=16):
    add_rect(slide, x, y, w, h, WHITE, RGBColor(0xD7, 0xDD, 0xDA), True, 0.9)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    if icon:
        add_text(slide, icon, x + 0.2, y + 0.16, 0.5, 0.45, 22, True, accent, PP_ALIGN.CENTER)
        tx = x + 0.75
    else:
        tx = x + 0.22
    add_text(slide, title, tx, y + 0.14, w - (tx - x) - 0.15, 0.42, 18, True, CHARCOAL)
    add_text(slide, body, x + 0.22, y + 0.66, w - 0.4, h - 0.78, body_size, False, MID)


def add_title(slide, title, number=None, section=None, subtitle=None):
    if number is not None:
        add_text(slide, f"{number:02d}", 0.32, 0.13, 0.45, 0.35, 11, True, GREEN, PP_ALIGN.CENTER)
    add_text(slide, title, 0.85 if number is not None else 0.35, 0.10, 11.8, 0.55, 27, True, CHARCOAL)
    if section:
        add_text(slide, section.upper(), 10.9, 0.20, 1.95, 0.24, 8.5, True, GREEN, PP_ALIGN.RIGHT)
    if subtitle:
        add_text(slide, subtitle, 0.42, 0.82, 12.0, 0.44, 14.5, False, MID)


def add_footer(slide, number, source=None):
    add_text(slide, f"{number}", 12.55, 7.12, 0.35, 0.2, 8, False, MID, PP_ALIGN.RIGHT)
    if source:
        add_text(slide, source, 0.42, 7.08, 11.8, 0.24, 7.5, False, MID)


def add_arrow(slide, x1, y1, x2, y2, color=GREEN, width=2.2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_pill(slide, text, x, y, w, fill=GREEN_LIGHT, color=GREEN_DARK, size=13):
    sh = add_rect(slide, x, y, w, 0.42, fill, fill, True, 0)
    add_text(slide, text, x + 0.05, y + 0.07, w - 0.1, 0.24, size, True, color, PP_ALIGN.CENTER)
    return sh


def add_formula(slide, text, x, y, w, h, size=22, color=CHARCOAL, fill=LIGHT):
    add_rect(slide, x, y, w, h, fill, RGBColor(0xD8, 0xDE, 0xDB), True, 0.8)
    add_text(slide, text, x + 0.12, y + 0.08, w - 0.24, h - 0.16, size, False, color,
             PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 0.02, MATH_FONT)


def add_image(slide, path, x, y, w, h=None):
    kwargs = {"left": Inches(x), "top": Inches(y), "width": Inches(w)}
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), **kwargs)


def mpl_style():
    plt.rcParams.update({
        "font.family": "DejaVu Sans",
        "axes.edgecolor": "#666a70",
        "axes.labelcolor": "#313336",
        "xtick.color": "#666a70",
        "ytick.color": "#666a70",
        "axes.titleweight": "bold",
    })


def generate_assets() -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    mpl_style()
    rng = np.random.default_rng(7)
    assets = {}

    # Neural observations versus latent trajectory.
    t = np.linspace(0, 4, 240)
    z1 = np.sin(1.4 * t) * np.exp(-0.05 * t)
    z2 = np.cos(0.8 * t + 0.4)
    rates = np.exp(0.8 * rng.normal(size=(28, 2)) @ np.vstack([z1, z2]) - 2.3)
    spikes = rng.poisson(rates * 0.35)
    fig, axes = plt.subplots(2, 1, figsize=(9, 5.2), sharex=True, gridspec_kw={"height_ratios": [1.6, 1]})
    rr, cc = np.where(spikes > 0)
    axes[0].scatter(t[cc], rr, s=6, color="#313336", alpha=0.8)
    axes[0].set_ylabel("neuron")
    axes[0].set_title("Observed spikes: sparse and variable", loc="left")
    axes[0].spines[["top", "right"]].set_visible(False)
    axes[1].plot(t, z1, color="#008A45", lw=3, label="latent 1")
    axes[1].plot(t, z2, color="#3F78B5", lw=3, label="latent 2")
    axes[1].set_ylabel("latent state"); axes[1].set_xlabel("time (s)")
    axes[1].legend(frameon=False, ncol=2, loc="upper right")
    axes[1].spines[["top", "right"]].set_visible(False)
    fig.tight_layout(); path = ASSETS / "phenomenon.png"; fig.savefig(path, dpi=200, transparent=True); plt.close(fig)
    assets["phenomenon"] = path

    # Mixture with uncertainty.
    a = rng.multivariate_normal([-1.5, 0], [[0.7, 0.35], [0.35, 0.8]], 220)
    b = rng.multivariate_normal([1.5, 0.7], [[0.8, -0.25], [-0.25, 0.6]], 220)
    xy = np.vstack([a, b])
    d0 = np.sum((xy - np.array([-1.5, 0])) ** 2, axis=1)
    d1 = np.sum((xy - np.array([1.5, 0.7])) ** 2, axis=1)
    resp = 1 / (1 + np.exp(np.clip(d1 - d0, -20, 20)))
    colors = np.c_[resp * 0.0 + (1-resp)*0.25, resp*0.54 + (1-resp)*0.47, resp*0.27 + (1-resp)*0.71]
    fig, ax = plt.subplots(figsize=(6.5, 4.6))
    ax.scatter(xy[:, 0], xy[:, 1], c=colors, s=22, alpha=0.75, edgecolor="none")
    ax.scatter([-1.5, 1.5], [0, 0.7], c=["#3F78B5", "#008A45"], s=180, marker="X", edgecolor="white", lw=2)
    ax.set_title("Soft assignments preserve uncertainty", loc="left")
    ax.set_xlabel("population feature 1"); ax.set_ylabel("population feature 2")
    ax.spines[["top", "right"]].set_visible(False); fig.tight_layout()
    path = ASSETS / "mixture.png"; fig.savefig(path, dpi=200, transparent=True); plt.close(fig); assets["mixture"] = path

    # HMM posterior heatmap.
    states = np.repeat([0, 1, 2, 1, 0], [35, 30, 42, 28, 45])
    K, T = 3, len(states)
    post = np.full((K, T), 0.06)
    for i, s in enumerate(states):
        post[s, i] = 0.88
    post += rng.normal(0, 0.035, post.shape); post = np.clip(post, 0, 1); post /= post.sum(0)
    cmap = LinearSegmentedColormap.from_list("lmu", ["#F4F6F5", "#8BC8A6", "#008A45"])
    fig, ax = plt.subplots(figsize=(9, 3.6))
    im = ax.imshow(post, aspect="auto", cmap=cmap, vmin=0, vmax=1)
    ax.set_yticks(range(K), ["rest", "plan", "move"]); ax.set_xlabel("time")
    ax.set_title("Posterior probability over discrete states", loc="left")
    ax.spines[:].set_visible(False); fig.colorbar(im, ax=ax, fraction=0.025, pad=0.02, label="probability")
    fig.tight_layout(); path = ASSETS / "hmm_posterior.png"; fig.savefig(path, dpi=200, transparent=True); plt.close(fig); assets["hmm"] = path

    # LDS vector field.
    xx, yy = np.meshgrid(np.linspace(-2.6, 2.6, 17), np.linspace(-2.2, 2.2, 15))
    u = -0.22 * xx - 0.95 * yy; v = 0.82 * xx - 0.28 * yy
    traj = np.zeros((150, 2)); traj[0] = [2.4, 0.4]
    dt = 0.05
    for i in range(1, len(traj)):
        x, y = traj[i-1]; traj[i] = traj[i-1] + dt * np.array([-0.22*x-0.95*y, 0.82*x-0.28*y])
    fig, ax = plt.subplots(figsize=(6.2, 5.0))
    ax.streamplot(xx, yy, u, v, color="#A7CDB8", density=1.1, linewidth=1.1, arrowsize=1)
    ax.plot(traj[:, 0], traj[:, 1], color="#008A45", lw=3)
    ax.scatter(*traj[0], color="#E87B35", s=80, zorder=4, label="initial condition")
    ax.scatter(0, 0, color="#313336", marker="X", s=100, zorder=4, label="fixed point")
    ax.set_title("Continuous latent flow", loc="left"); ax.set_xlabel("latent 1"); ax.set_ylabel("latent 2")
    ax.legend(frameon=False, loc="upper right"); ax.set_aspect("equal"); ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout(); path = ASSETS / "lds_flow.png"; fig.savefig(path, dpi=200, transparent=True); plt.close(fig); assets["lds"] = path

    # LFADS-like denoising illustration.
    rate_true = 2.5 + 2*np.sin(2*np.pi*t/2.8) + 1.3*np.sin(2*np.pi*t/1.05 + .5)
    rate_true = np.clip(rate_true, .1, None)
    counts = rng.poisson(rate_true)
    kernel = np.ones(13) / 13
    smooth = np.convolve(counts, kernel, mode="same")
    fig, ax = plt.subplots(figsize=(9, 4.5))
    ax.vlines(t, 0, counts, color="#9CA1A4", lw=1, alpha=.8, label="observed counts")
    ax.plot(t, rate_true, color="#008A45", lw=3, label="underlying rate")
    ax.plot(t, smooth, color="#3F78B5", lw=2.2, ls="--", label="inferred rate")
    ax.set_title("A dynamical prior pools information across time", loc="left")
    ax.set_xlabel("time (s)"); ax.set_ylabel("count / rate"); ax.legend(frameon=False, ncol=3)
    ax.spines[["top", "right"]].set_visible(False); fig.tight_layout()
    path = ASSETS / "lfads_denoise.png"; fig.savefig(path, dpi=200, transparent=True); plt.close(fig); assets["lfads"] = path

    # gpSLDS smoothly switching field.
    xx, yy = np.meshgrid(np.linspace(-2.7, 2.7, 19), np.linspace(-2.2, 2.2, 17))
    w = 1 / (1 + np.exp(-2.5 * xx))
    u1, v1 = -0.2*xx - 1.0*yy, 0.9*xx - 0.2*yy
    u2, v2 = -0.25*xx + 0.85*yy, -0.9*xx - 0.25*yy
    u, v = (1-w)*u1 + w*u2, (1-w)*v1 + w*v2
    mag = np.hypot(u, v)
    fig, ax = plt.subplots(figsize=(7.2, 5.0))
    ax.contourf(xx, yy, w, levels=np.linspace(0,1,15), cmap="PRGn", alpha=.22)
    ax.streamplot(xx, yy, u, v, color=mag, cmap="Greens", density=1.15, linewidth=1.2, arrowsize=1)
    ax.axvline(0, color="#666A70", lw=1.2, ls="--")
    ax.text(-2.45,1.9,"local regime A",color="#3F78B5",weight="bold")
    ax.text(1.25,1.9,"local regime B",color="#008A45",weight="bold")
    ax.set_title("Smooth interpolation between local linear regimes", loc="left")
    ax.set_xlabel("latent 1"); ax.set_ylabel("latent 2"); ax.set_aspect("equal")
    ax.spines[["top", "right"]].set_visible(False); fig.tight_layout()
    path = ASSETS / "gpslds_flow.png"; fig.savefig(path, dpi=200, transparent=True); plt.close(fig); assets["gpslds"] = path

    # Continuous-time stability plane.
    fig, ax = plt.subplots(figsize=(6.4, 4.5))
    ax.axvline(0, color="#313336", lw=1.4); ax.axhline(0, color="#999999", lw=.8)
    ax.axvspan(-2.2, 0, color="#E8F5EE"); ax.axvspan(0, 2.2, color="#FBEDEA")
    stable = np.array([[-1.4,.7],[-.8,-1.0],[-.35,.2]])
    unstable = np.array([[.35,.8],[1.0,-.6],[1.45,.15]])
    ax.scatter(stable[:,0],stable[:,1],s=90,color="#008A45",label="stable")
    ax.scatter(unstable[:,0],unstable[:,1],s=90,color="#C94545",label="unstable")
    ax.scatter([0],[1.25],s=100,color="#E87B35",marker="D",label="marginal")
    ax.set_xlim(-2.2,2.2); ax.set_ylim(-1.6,1.7)
    ax.set_xlabel("Re(λ)"); ax.set_ylabel("Im(λ)"); ax.set_title("Continuous-time local stability",loc="left")
    ax.legend(frameon=False,ncol=3,loc="lower center"); ax.spines[["top","right"]].set_visible(False)
    fig.tight_layout(); path=ASSETS/"stability.png"; fig.savefig(path,dpi=200,transparent=True); plt.close(fig); assets["stability"]=path
    return assets


def fill_title_slide(slide):
    # Preserve the template artwork; correct and rebalance text only.
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text.strip()
        if "Latent Variable Models" in txt:
            sh.text_frame.clear(); p=sh.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
            for line in ["Latent Variable Models for", "Neural Data and Dynamics"]:
                if p.text:
                    p = sh.text_frame.add_paragraph(); p.alignment=PP_ALIGN.CENTER
                r=p.add_run(); r.text=line; set_run(r,25,True,WHITE)
        elif "Chunyu Qu" in txt or "22.07.2026" in txt:
            sh.text="Chunyu Qu\n22.07.2026"
            for p in sh.text_frame.paragraphs:
                p.alignment=PP_ALIGN.CENTER
                for r in p.runs: set_run(r,19,True,WHITE)
        elif "Probabilistic" in txt:
            sh.text="From Probabilistic State-Space Models\nto Deep Generative Models"
            for p in sh.text_frame.paragraphs:
                p.alignment=PP_ALIGN.CENTER
                for r in p.runs: set_run(r,15,False,WHITE,italic=True)


def build_deck():
    if not TEMPLATE.exists():
        raise FileNotFoundError(TEMPLATE)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    assets = generate_assets()
    shutil.copy2(TEMPLATE, OUTPUT)
    prs = Presentation(OUTPUT)
    if len(prs.slides) < 36:
        raise RuntimeError("Template must contain at least 36 slides")

    fill_title_slide(prs.slides[0])

    # 2. Two motivating phenomena, one inference problem.
    s=prs.slides[1]; clear_slide(s)
    add_rect(s,0,.78,13.333,6.72,RGBColor(0xF7,0xF9,0xF8),RGBColor(0xF7,0xF9,0xF8),False,0)
    add_title(s,"Two phenomena, one hidden-state problem",section="MOTIVATION")

    # Left: neural spikes are discrete observations of a latent rate/state.
    add_rect(s,.58,1.18,5.95,4.70,WHITE,RGBColor(0xD5,0xDD,0xD9),True,.9)
    add_pill(s,"01",.88,1.48,.52,BLUE,WHITE,12)
    add_text(s,"NEURAL ACTIVITY",1.58,1.44,2.55,.38,15,True,BLUE)
    add_text(s,"We observe spikes",.90,1.98,4.90,.48,24,True,CHARCOAL)
    rng=np.random.default_rng(23)
    for row in range(8):
        for xdot in rng.uniform(1.02,3.55,5+row%3):
            dot=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(xdot),Inches(2.72+row*.16),Inches(.038),Inches(.038))
            dot.fill.solid(); dot.fill.fore_color.rgb=CHARCOAL; dot.line.fill.background()
    add_text(s,"neurons",.72,3.18,.55,.28,10,False,MID,PP_ALIGN.CENTER)
    add_arrow(s,3.86,3.32,4.38,3.32,RGBColor(0xA8,0xB4,0xAE),1.5)
    add_formula(s,"xₜ ~ Poisson(λₜ)",4.42,2.80,1.63,.78,17,CHARCOAL,GREEN_LIGHT)
    add_text(s,"What latent rate or state generated them?",.92,4.42,4.98,.62,17,True,GREEN_DARK,PP_ALIGN.CENTER)
    add_text(s,"spikes  →  firing rate  →  neural dynamics",1.05,5.20,4.70,.35,13,False,MID,PP_ALIGN.CENTER)

    # Right: sequential perceptual reports reveal evolving beliefs and precision.
    add_rect(s,6.80,1.18,5.95,4.70,WHITE,RGBColor(0xD5,0xDD,0xD9),True,.9)
    add_pill(s,"02",7.10,1.48,.52,ORANGE,WHITE,12)
    add_text(s,"PERCEPTUAL DYNAMICS",7.80,1.44,3.15,.38,15,True,ORANGE)
    add_text(s,"We observe trial-by-trial bias",7.12,1.98,5.05,.48,24,True,CHARCOAL)
    stim=[.20,.62,.35,.78,.54,.70,.42]
    for i,v in enumerate(stim):
        x=7.35+i*.68; y=3.63-v*1.08
        stem=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x),Inches(3.72),Inches(x),Inches(y))
        stem.line.color.rgb=RGBColor(0xC3,0xCB,0xC7); stem.line.width=Pt(1.2)
        dot=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x-.055),Inches(y-.055),Inches(.11),Inches(.11))
        dot.fill.solid(); dot.fill.fore_color.rgb=ORANGE; dot.line.fill.background()
    add_text(s,"trial",9.13,3.78,1.20,.28,10,False,MID,PP_ALIGN.CENTER)
    add_formula(s,"beliefₜ  →  responseₜ",8.00,4.18,3.55,.68,18,CHARCOAL,RGBColor(0xFD,0xF3,0xED))
    add_text(s,"What hidden belief determines how history is weighted?",7.16,4.93,5.18,.60,16.5,True,GREEN_DARK,PP_ALIGN.CENTER)

    add_formula(s,"noisy observations  →  infer an evolving latent cause",2.10,6.18,9.15,.62,20,GREEN_DARK,GREEN_LIGHT)
    add_footer(s,2,"Linderman, ML4ND · Glasauer & Shi, Sci Rep (2022) · Qu & Shi, bioRxiv (2026)")

    # 3.
    s=prs.slides[2]; add_title(s,"We observe activity—not the computation itself",3,"PHENOMENON")
    add_image(s,assets["phenomenon"],.55,1.18,7.25,4.9)
    add_card(s,"What changes across trials?","Spikes vary even when task conditions appear identical.",8.15,1.42,4.55,1.25,GREEN,"?")
    add_card(s,"What is shared?","Population-level trajectories may remain structured and reproducible.",8.15,2.96,4.55,1.25,BLUE,"≈")
    add_card(s,"The scientific question","Which variation reflects latent computation, observation noise, or unobserved input?",8.15,4.50,4.55,1.55,ORANGE,"→")
    add_footer(s,3)

    # 4.
    s=prs.slides[3]; add_title(s,"The measurement problem has three layers",4,"PHENOMENON")
    cards=[("Partial","We record only a subset of a much larger circuit.","◌",BLUE),("Stochastic","Spiking remains variable conditional on latent state.","⋮",ORANGE),("Indirect","Behavior and events provide incomplete external context.","↗",GREEN)]
    for i,(t,b,ic,c) in enumerate(cards):add_card(s,t,b,.65+i*4.22,1.55,3.75,2.15,c,ic,17)
    add_formula(s,"observed spikes  =  structured population signal  +  residual variability",1.2,4.35,10.9,1.0,23,GREEN_DARK,GREEN_LIGHT)
    add_text(s,"A latent-variable model formalizes this decomposition; it does not guarantee a unique biological explanation.",1.1,5.75,11.1,.62,17,False,MID,PP_ALIGN.CENTER)
    add_footer(s,4)

    # 5.
    s=prs.slides[4]; add_title(s,"The latent dynamics hypothesis",5,"HYPOTHESIS", "A low-dimensional hidden state generates high-dimensional neural observations")
    add_rect(s,.6,1.55,3.1,3.9,WHITE,RGBColor(0xD9,0xDF,0xDC)); add_text(s,"LATENT STATE",.85,1.83,2.6,.35,13,True,GREEN,PP_ALIGN.CENTER)
    add_text(s,"zₜ ∈ ℝᵈ",1.22,2.45,1.9,.55,30,True,CHARCOAL,PP_ALIGN.CENTER,font=MATH_FONT)
    add_text(s,"low-dimensional\ncontinuous trajectory",1.0,3.28,2.3,.85,16,False,MID,PP_ALIGN.CENTER)
    add_arrow(s,3.78,3.35,5.05,3.35,GREEN,3)
    add_rect(s,5.1,1.55,3.1,3.9,GREEN_LIGHT,GREEN); add_text(s,"OBSERVATION MODEL",5.33,1.83,2.65,.35,13,True,GREEN_DARK,PP_ALIGN.CENTER)
    add_text(s,"λₜ = fθ(zₜ)",5.53,2.42,2.25,.55,26,True,CHARCOAL,PP_ALIGN.CENTER,font=MATH_FONT)
    add_text(s,"yₜ,n ~ Poisson(λₜ,n)",5.35,3.34,2.6,.55,19,False,CHARCOAL,PP_ALIGN.CENTER,font=MATH_FONT)
    add_arrow(s,8.28,3.35,9.55,3.35,GREEN,3)
    add_rect(s,9.6,1.55,3.1,3.9,WHITE,RGBColor(0xD9,0xDF,0xDC)); add_text(s,"RECORDED DATA",9.86,1.83,2.55,.35,13,True,GREEN,PP_ALIGN.CENTER)
    for i in range(6):
        x=10.15+(i%3)*.72; y=2.55+(i//3)*1.0
        circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(.34),Inches(.34)); circ.fill.solid();circ.fill.fore_color.rgb=CHARCOAL;circ.line.fill.background()
    add_text(s,"N neurons × T bins",9.94,4.25,2.45,.4,15,False,MID,PP_ALIGN.CENTER)
    add_text(s,"Hypothesis ≠ proof: latent coordinates can rotate, rescale, or change across equally predictive models.",.85,5.85,11.7,.5,17,True,ORANGE,PP_ALIGN.CENTER)
    add_footer(s,5)

    # 6.
    s=prs.slides[5]; add_title(s,"Five questions prevent model-driven storytelling",6,"HYPOTHESIS")
    qs=[("1","What is latent?"),("2","How does it evolve?"),("3","How does it generate observations?"),("4","Is inference exact or approximate?"),("5","What interpretation is justified?")]
    for i,(n,q) in enumerate(qs):
        y=1.15+i*1.02; add_pill(s,n,.75,y,.55,GREEN,WHITE,15); add_text(s,q,1.55,y-.01,10.8,.5,21,True if i==4 else False,CHARCOAL)
    add_text(s,"Every model in this talk gives a different answer.",1.55,6.25,10.7,.48,19,True,GREEN_DARK)
    add_footer(s,6)

    # 7.
    s=prs.slides[6]; add_title(s,"Mixture models introduce a latent assignment",7,"EXACT MODELS")
    add_image(s,assets["mixture"],.45,1.18,6.65,4.95)
    add_formula(s,"zᵢ ~ Categorical(π)",7.55,1.48,4.85,.72,21)
    add_formula(s,"yᵢ | zᵢ=k ~ 𝒩(μₖ, Σₖ)",7.55,2.48,4.85,.72,20)
    add_card(s,"What it explains","Heterogeneity: observations arise from different latent regimes.",7.55,3.58,4.85,1.15,GREEN)
    add_card(s,"What it ignores","Time: assignments are independent across observations.",7.55,4.98,4.85,1.15,ORANGE)
    add_footer(s,7)

    # 8.
    s=prs.slides[7]; add_title(s,"EM alternates uncertainty and parameter learning",8,"EXACT MODELS")
    add_card(s,"E-step","Compute soft responsibilities\nγᵢₖ = p(zᵢ=k | yᵢ, θold)",.7,1.45,4.9,2.15,BLUE,"E",19)
    add_arrow(s,5.78,2.52,7.42,2.52,GREEN,3)
    add_card(s,"M-step","Update parameters using responsibilities as fractional counts.",7.6,1.45,4.9,2.15,GREEN,"M",19)
    add_arrow(s,10.0,3.82,3.18,3.82,MID,1.8)
    add_text(s,"repeat until convergence",5.07,3.67,3.2,.32,13,True,MID,PP_ALIGN.CENTER)
    add_formula(s,"log evidence = ELBO + KL(q(z) ‖ p(z | y, θ))",1.25,4.55,10.8,.9,22,GREEN_DARK,GREEN_LIGHT)
    add_text(s,"When the exact posterior is available, the E-step can close the ELBO gap.",1.15,5.80,11.0,.45,18,False,MID,PP_ALIGN.CENTER)
    add_footer(s,8)

    # 9.
    s=prs.slides[8]; add_title(s,"HMM: connect latent assignments through time",9,"EXACT MODELS")
    labels=["rest","plan","move","move","return"]
    colors=[BLUE,ORANGE,GREEN,GREEN,RGBColor(0x79,0x63,0xA6)]
    for i,(lab,c) in enumerate(zip(labels,colors)):
        x=.72+i*2.43; add_rect(s,x,1.6,1.72,.8,c,c,True,0); add_text(s,lab,x+.08,1.82,1.56,.28,17,True,WHITE,PP_ALIGN.CENTER)
        if i<4:add_arrow(s,x+1.75,2.0,x+2.30,2.0,MID,2)
        add_arrow(s,x+.86,2.45,x+.86,3.15,c,1.8)
        add_rect(s,x,3.25,1.72,1.05,WHITE,c,True,1.2); add_text(s,"yₜ",x+.12,3.57,1.48,.32,21,True,c,PP_ALIGN.CENTER,font=MATH_FONT)
    add_formula(s,"p(z₁:T, y₁:T) = p(z₁) ∏ₜ p(zₜ | zₜ₋₁) ∏ₜ p(yₜ | zₜ)",1.15,4.92,11.05,.95,22)
    add_text(s,"The state label means only what the transition and emission assumptions make it mean.",1.05,6.14,11.2,.45,18,True,ORANGE,PP_ALIGN.CENTER)
    add_footer(s,9)

    # 10.
    s=prs.slides[9]; add_title(s,"Forward–backward makes sequence inference exact",10,"EXACT MODELS")
    add_image(s,assets["hmm"],.62,1.25,7.25,3.65)
    add_card(s,"Filtering","p(zₜ | y₁:ₜ)\nonline; past + present",8.25,1.32,4.25,1.25,BLUE,"→")
    add_card(s,"Smoothing","p(zₜ | y₁:T)\noffline; complete sequence",8.25,2.87,4.25,1.25,GREEN,"↔")
    add_card(s,"Viterbi","Most probable joint path—not pointwise argmax.",8.25,4.42,4.25,1.25,ORANGE,"V")
    add_formula(s,"O(TK²) instead of enumerating Kᵀ paths",2.15,5.55,6.9,.7,20,GREEN_DARK,GREEN_LIGHT)
    add_footer(s,10)

    # 11.
    s=prs.slides[10]; add_title(s,"A standard HMM is interpretable—but coarse",11,"EXACT MODELS")
    items=[("Discrete trajectory","Continuous neural flow is compressed into labels.",BLUE),("Within-state dynamics","Gaussian emissions describe distributions, not flow.",GREEN),("Duration assumption","Homogeneous HMMs imply geometric state durations.",ORANGE),("Identifiability","Labels permute; K and priors change the partition.",RGBColor(0x79,0x63,0xA6))]
    for i,(t,b,c) in enumerate(items):add_card(s,t,b,.65+(i%2)*6.12,1.3+(i//2)*2.25,5.75,1.72,c,str(i+1),17)
    add_text(s,"Extensions: HSMM · AR-HMM · LDS · switching LDS",1.05,6.05,11.2,.5,21,True,GREEN_DARK,PP_ALIGN.CENTER)
    add_footer(s,11)

    # 12.
    s=prs.slides[11]; add_title(s,"LDS: move the dynamics into a continuous latent space",12,"EXACT MODELS")
    add_image(s,assets["lds"],.45,1.18,6.35,5.45)
    add_formula(s,"zₜ = A zₜ₋₁ + εₜ",7.25,1.55,4.95,.86,25)
    add_formula(s,"yₜ = C zₜ + d + ηₜ",7.25,2.75,4.95,.86,25)
    add_card(s,"Dynamics matrix A","Defines local rotation, contraction, expansion, and timescales.",7.25,4.00,4.95,1.12,GREEN)
    add_card(s,"Observation matrix C","Maps the hidden trajectory to measured population activity.",7.25,5.35,4.95,1.12,BLUE)
    add_footer(s,12)

    # 13.
    s=prs.slides[12]; add_title(s,"Kalman inference alternates prediction and correction",13,"EXACT MODELS")
    xs=[.75,3.55,6.35,9.15]; titles=["Previous posterior","Predict","Observe","Update"]
    bodies=["p(zₜ₋₁ | y₁:ₜ₋₁)","p(zₜ | y₁:ₜ₋₁)","likelihood p(yₜ | zₜ)","p(zₜ | y₁:ₜ)"]
    cols=[MID,BLUE,ORANGE,GREEN]
    for i,(t,b,c) in enumerate(zip(titles,bodies,cols)):
        add_rect(s,xs[i],1.72,2.15,2.25,WHITE,c,True,1.4); add_text(s,t,xs[i]+.12,1.98,1.9,.45,17,True,c,PP_ALIGN.CENTER); add_text(s,b,xs[i]+.13,2.72,1.88,.55,16,False,CHARCOAL,PP_ALIGN.CENTER,font=MATH_FONT)
        if i<3:add_arrow(s,xs[i]+2.18,2.85,xs[i]+2.70,2.85,GREEN,2.2)
    add_formula(s,"Linear + Gaussian  ⇒  exact Gaussian filtering and smoothing",1.25,4.65,10.85,.86,23,GREEN_DARK,GREEN_LIGHT)
    add_text(s,"Continuous latents do not automatically require variational inference.",1.2,5.85,10.95,.48,20,True,ORANGE,PP_ALIGN.CENTER)
    add_footer(s,13)

    # 14.
    s=prs.slides[13]; add_title(s,"Where exact inference ends",14,"INFERENCE BOUNDARY")
    add_rect(s,.7,1.30,5.75,4.85,GREEN_LIGHT,GREEN,True,1.3); add_text(s,"EXACT / ANALYTIC",.95,1.63,5.25,.45,18,True,GREEN_DARK,PP_ALIGN.CENTER)
    add_bullets(s,["Finite-state HMM","Linear-Gaussian LDS","Conjugate exponential-family models"],1.25,2.40,4.75,2.1,20,CHARCOAL,14)
    add_text(s,"Dynamic programming · Kalman · exact E-step",1.15,5.20,4.95,.45,15,True,GREEN_DARK,PP_ALIGN.CENTER)
    add_rect(s,6.88,1.30,5.75,4.85,RGBColor(0xFD,0xF3,0xED),ORANGE,True,1.3); add_text(s,"APPROXIMATE",7.13,1.63,5.25,.45,18,True,ORANGE,PP_ALIGN.CENTER)
    add_bullets(s,["Nonlinear decoder","Poisson / nonconjugate observations","Flexible neural or GP dynamics"],7.43,2.40,4.75,2.1,20,CHARCOAL,14)
    add_text(s,"Variational inference · sampling · numerical approximation",7.28,5.20,4.95,.45,15,True,ORANGE,PP_ALIGN.CENTER)
    add_footer(s,14)

    # 15.
    s=prs.slides[14]; add_title(s,"The hard object is a posterior over a trajectory",15,"VARIATIONAL INFERENCE")
    add_formula(s,"pθ(y₁:T) = ∫ pθ(y₁:T, z₁:T) dz₁:T",1.65,1.55,10.1,1.05,27)
    add_text(s,"For T time bins and d latent dimensions, the integral spans Td coupled variables.",1.2,2.95,10.95,.48,20,False,MID,PP_ALIGN.CENTER)
    add_arrow(s,6.66,3.65,6.66,4.28,ORANGE,3)
    add_formula(s,"Replace pθ(z₁:T | y₁:T) with a tractable qφ(z₁:T | y₁:T)",1.35,4.48,10.65,1.05,24,GREEN_DARK,GREEN_LIGHT)
    add_text(s,"Approximation is now part of the model's uncertainty story.",1.25,5.95,10.85,.45,18,True,ORANGE,PP_ALIGN.CENTER)
    add_footer(s,15)

    # 16.
    s=prs.slides[15]; add_title(s,"The ELBO balances fit and posterior regularization",16,"VARIATIONAL INFERENCE")
    add_rect(s,.7,1.45,5.65,3.65,RGBColor(0xED,0xF3,0xFA),BLUE,True,1.2); add_text(s,"EXPLAIN THE DATA",1.0,1.78,5.05,.42,18,True,BLUE,PP_ALIGN.CENTER)
    add_text(s,"𝔼q [ log pθ(y | z) ]",1.12,2.55,4.8,.65,25,True,CHARCOAL,PP_ALIGN.CENTER,font=MATH_FONT)
    add_text(s,"Does the sampled latent state reproduce the observation?",1.18,3.65,4.7,.85,17,False,MID,PP_ALIGN.CENTER)
    add_rect(s,6.98,1.45,5.65,3.65,GREEN_LIGHT,GREEN,True,1.2); add_text(s,"CONSTRAIN THE CODE",7.28,1.78,5.05,.42,18,True,GREEN_DARK,PP_ALIGN.CENTER)
    add_text(s,"− KL(qφ(z | y) ‖ pθ(z))",7.40,2.55,4.8,.65,24,True,CHARCOAL,PP_ALIGN.CENTER,font=MATH_FONT)
    add_text(s,"Does the posterior respect the structure encoded by the prior?",7.47,3.65,4.7,.85,17,False,MID,PP_ALIGN.CENTER)
    add_formula(s,"log pθ(y) = ELBO + KL(qφ(z|y) ‖ pθ(z|y))",1.45,5.45,10.45,.82,21)
    add_footer(s,16)

    # 17.
    s=prs.slides[16]; add_title(s,"A VAE learns a reusable inference rule",17,"VARIATIONAL INFERENCE")
    nodes=[("DATA","y",.6,2.15,BLUE),("ENCODER","μφ(y), σφ(y)",3.1,1.75,GREEN),("SAMPLE","z = μ + σ⊙ε",5.95,2.15,ORANGE),("DECODER","pθ(y | z)",8.65,1.75,GREEN),("RECONSTRUCTION","ŷ",11.25,2.15,BLUE)]
    for i,(head,body,x,y,c) in enumerate(nodes):
        w=1.55 if i in [0,4] else 2.15; add_rect(s,x,y,w,1.65,WHITE,c,True,1.3); add_text(s,head,x+.08,y+.20,w-.16,.3,11,True,c,PP_ALIGN.CENTER); add_text(s,body,x+.08,y+.78,w-.16,.45,18,True,CHARCOAL,PP_ALIGN.CENTER,font=MATH_FONT)
        if i<len(nodes)-1:
            nx=nodes[i+1][2]; add_arrow(s,x+w,y+.82,nx-.12,nodes[i+1][3]+.82,GREEN,2.2)
    add_card(s,"Amortization","One encoder predicts variational parameters for new observations in a forward pass.",1.15,4.65,5.25,1.2,GREEN)
    add_card(s,"Trade-off","The shared encoder can introduce an amortization gap and does not solve identifiability.",6.92,4.65,5.25,1.2,ORANGE)
    add_footer(s,17)

    # 18.
    s=prs.slides[17]; add_title(s,"Reparameterization separates randomness from parameters",18,"VARIATIONAL INFERENCE")
    add_rect(s,.85,1.45,3.15,3.9,WHITE,RGBColor(0xD9,0xDF,0xDC)); add_text(s,"ENCODER OUTPUT",1.1,1.82,2.65,.35,14,True,GREEN,PP_ALIGN.CENTER); add_formula(s,"μ, σ",1.43,2.55,2.0,.8,28,CHARCOAL,WHITE)
    add_rect(s,5.08,1.45,3.15,3.9,RGBColor(0xFD,0xF3,0xED),ORANGE); add_text(s,"PARAMETER-FREE NOISE",5.33,1.82,2.65,.35,14,True,ORANGE,PP_ALIGN.CENTER); add_formula(s,"ε ~ 𝒩(0, I)",5.48,2.55,2.35,.8,23,CHARCOAL,WHITE)
    add_rect(s,9.30,1.45,3.15,3.9,GREEN_LIGHT,GREEN); add_text(s,"LATENT SAMPLE",9.55,1.82,2.65,.35,14,True,GREEN_DARK,PP_ALIGN.CENTER); add_formula(s,"z = μ + σ ⊙ ε",9.62,2.55,2.5,.8,21,CHARCOAL,WHITE)
    add_arrow(s,4.12,3.35,4.92,3.35,GREEN,3); add_arrow(s,8.35,3.35,9.14,3.35,GREEN,3)
    add_text(s,"Gradients flow through μ and σ while ε carries the stochasticity.",1.15,5.90,11.0,.5,20,True,GREEN_DARK,PP_ALIGN.CENTER)
    add_footer(s,18)

    # 19.
    s=prs.slides[18]; add_title(s,"A standard VAE is not yet a dynamical model",19,"VARIATIONAL INFERENCE")
    add_card(s,"Standard VAE","One observation ↔ one latent code\nNo temporal prior is implied.",.75,1.55,3.65,2.0,BLUE,"1")
    add_arrow(s,4.65,2.55,5.35,2.55,GREEN,3)
    add_card(s,"Sequential VAE","Posterior and prior are defined over an entire trajectory.",5.6,1.55,3.65,2.0,GREEN,"T")
    add_arrow(s,9.50,2.55,10.20,2.55,GREEN,3)
    add_card(s,"Neural dynamics","Latent evolution generates time-varying rates and spikes.",10.45,1.55,2.25,2.0,ORANGE,"λ",15)
    add_formula(s,"p(z₁:T) = p(z₁) ∏ₜ p(zₜ | zₜ₋₁, uₜ)",1.45,4.25,10.45,.9,24)
    add_text(s,"RNNs are one parameterization—not the definition of sequential variational inference.",1.1,5.55,11.1,.55,18,True,ORANGE,PP_ALIGN.CENTER)
    add_footer(s,19)

    # 20.
    s=prs.slides[19]; add_title(s,"Sequential inference can look backward, forward—or both",20,"VARIATIONAL INFERENCE")
    rows=[("Filtering posterior","q(zₜ | y₁:ₜ)","online decoding",BLUE),("Smoothing posterior","q(zₜ | y₁:T)","offline trajectory recovery",GREEN),("Factorized posterior","∏ₜ q(zₜ | y)","simple, but drops dependencies",ORANGE),("Structured posterior","q(z₁:T | y₁:T)","richer temporal uncertainty",RGBColor(0x79,0x63,0xA6))]
    for i,(a,b,c,col) in enumerate(rows):
        y=1.15+i*1.32; add_rect(s,.65,y,12.0,1.0,WHITE,RGBColor(0xD8,0xDE,0xDB),True,.8); add_pill(s,a,.85,y+.27,2.35,col,WHITE,12); add_text(s,b,3.55,y+.27,3.25,.35,18,True,CHARCOAL,font=MATH_FONT); add_text(s,c,7.15,y+.27,4.95,.35,17,False,MID)
    add_footer(s,20)

    # 21.
    s=prs.slides[20]; add_title(s,"LFADS: infer a dynamical generator from spike trains",21,"LFADS")
    xvals=[.45,2.35,4.65,7.05,9.25,11.35]; labels=[("SPIKES","y₁:T",BLUE),("ENCODER","q(g₀|y)",GREEN),("INIT","g₀",ORANGE),("GENERATOR","g₁:T",GREEN),("FACTORS","f₁:T",BLUE),("RATES","λ₁:T",ORANGE)]
    for i,(x,(h,b,c)) in enumerate(zip(xvals,labels)):
        w=1.55 if i not in [3] else 1.75; add_rect(s,x,1.75,w,1.65,WHITE,c,True,1.2); add_text(s,h,x+.05,1.98,w-.1,.3,10.5,True,c,PP_ALIGN.CENTER); add_text(s,b,x+.08,2.62,w-.16,.35,17,True,CHARCOAL,PP_ALIGN.CENTER,font=MATH_FONT)
        if i<5:add_arrow(s,x+w,2.58,xvals[i+1]-.1,2.58,GREEN,2)
    add_formula(s,"yₜ,n ~ Poisson(λₜ,n)",4.18,4.15,4.95,.75,22)
    add_card(s,"Optional controller","Infers time-varying inputs uₜ when autonomous dynamics are insufficient.",.75,5.15,5.55,1.25,ORANGE,None,12)
    add_card(s,"Key distinction","The generator trajectory is deterministic conditional on g₀ and u₁:T.",6.78,5.15,5.55,1.25,GREEN,None,12)
    add_footer(s,21,"Pandarinath et al., Nature Methods (2018)")

    # 22.
    s=prs.slides[21]; add_title(s,"LFADS separates stochastic inference from deterministic evolution",22,"LFADS")
    add_rect(s,.65,1.35,5.8,4.65,RGBColor(0xED,0xF3,0xFA),BLUE,True,1.2); add_text(s,"STOCHASTIC INFERENCE",.95,1.72,5.2,.4,18,True,BLUE,PP_ALIGN.CENTER)
    add_formula(s,"g₀ ~ qφ(g₀ | y₁:T)",1.35,2.42,4.4,.72,21,CHARCOAL,WHITE)
    add_formula(s,"uₜ ~ qφ(uₜ | y₁:T)",1.35,3.42,4.4,.72,21,CHARCOAL,WHITE)
    add_text(s,"What initial condition and input best explain this trial?",1.13,4.75,4.95,.6,17,False,MID,PP_ALIGN.CENTER)
    add_rect(s,6.88,1.35,5.8,4.65,GREEN_LIGHT,GREEN,True,1.2); add_text(s,"DETERMINISTIC GENERATOR",7.18,1.72,5.2,.4,18,True,GREEN_DARK,PP_ALIGN.CENTER)
    add_formula(s,"gₜ = RNNθ(gₜ₋₁, uₜ)",7.58,2.42,4.4,.72,21,CHARCOAL,WHITE)
    add_formula(s,"λₜ = exp(W fₜ + b)",7.58,3.42,4.4,.72,21,CHARCOAL,WHITE)
    add_text(s,"How does the inferred system evolve and generate rates?",7.35,4.75,4.95,.6,17,False,MID,PP_ALIGN.CENTER)
    add_footer(s,22,"Pandarinath et al., Nature Methods (2018)")

    # 23.
    s=prs.slides[22]; add_title(s,"The LFADS objective has three jobs",23,"LFADS")
    jobs=[("Explain spikes","𝔼q[log p(y₁:T | g₀,u₁:T)]",BLUE),("Constrain initial state","− KL(q(g₀|y) ‖ p(g₀))",GREEN),("Constrain inferred inputs","− KL(q(u₁:T|y) ‖ p(u₁:T))",ORANGE)]
    for i,(t,f,c) in enumerate(jobs):
        x=.65+i*4.22; add_card(s,t,f,x,1.45,3.75,2.4,c,str(i+1),18)
    add_formula(s,"ELBO = spike likelihood − initial-condition KL − input KL",1.25,4.65,10.85,.9,23,GREEN_DARK,GREEN_LIGHT)
    add_text(s,"A high ELBO validates optimization—not the biological interpretation of g₀ or uₜ.",1.0,5.95,11.3,.46,18,True,ORANGE,PP_ALIGN.CENTER)
    add_footer(s,23,"The input KL is omitted in an autonomous LFADS model.")

    # 24.
    s=prs.slides[23]; add_title(s,"A dynamical prior pools evidence across time",24,"LFADS")
    add_image(s,assets["lfads"],.55,1.15,8.0,4.95)
    add_card(s,"Denoising","The model estimates a rate trajectory that explains sparse counts over time.",8.9,1.45,3.75,1.35,GREEN)
    add_card(s,"Single-trial structure","Initial conditions and inputs can vary between trials.",8.9,3.15,3.75,1.35,BLUE)
    add_card(s,"Not guaranteed","Forward simulation is possible; accurate extrapolation must be tested.",8.9,4.85,3.75,1.35,ORANGE)
    add_footer(s,24,"Illustration uses synthetic data; no empirical performance claim is implied.")

    # 25.
    s=prs.slides[24]; add_title(s,"What LFADS demonstrated—and what remains conditional",25,"LFADS")
    add_rect(s,.65,1.25,5.85,4.95,GREEN_LIGHT,GREEN,True,1.2); add_text(s,"SUPPORTED IN THE ORIGINAL STUDY",.95,1.58,5.25,.38,16,True,GREEN_DARK,PP_ALIGN.CENTER)
    add_bullets(s,["Recovery of known synthetic dynamics","Improved single-trial firing-rate estimates","Prediction of behavioral variables","Inferred perturbations and multisession alignment"],1.08,2.20,4.95,3.2,16,CHARCOAL,11)
    add_rect(s,6.85,1.25,5.85,4.95,RGBColor(0xFD,0xF3,0xED),ORANGE,True,1.2); add_text(s,"STILL REQUIRES VALIDATION",7.15,1.58,5.25,.38,16,True,ORANGE,PP_ALIGN.CENTER)
    add_bullets(s,["Generalization to new conditions","Autonomous vs. input-driven interpretation","Fixed points and local Jacobians","Robustness across seeds and model choices"],7.28,2.20,4.95,3.2,16,CHARCOAL,11)
    add_footer(s,25,"Pandarinath et al., Nature Methods (2018) · doi:10.1038/s41592-018-0109-9")

    # 26.
    s=prs.slides[25]; add_title(s,"Flexibility creates a new scientific question",26,"INTERPRETABILITY")
    add_text(s,"LFADS can fit nonlinear recurrent dynamics.",.8,1.35,5.35,.55,24,True,GREEN_DARK)
    add_text(s,"But which aspects of the learned vector field are stable, local, and supported by data?",.8,2.15,11.7,1.0,29,True,CHARCOAL)
    qs=[("Fixed points","Where does flow slow or stop?"),("Stability","Which directions contract or expand?"),("Regimes","Does local flow change across state space?"),("Uncertainty","Where are dynamics weakly determined?")]
    for i,(t,b) in enumerate(qs):add_card(s,t,b,.7+i*3.12,4.15,2.85,1.35,[GREEN,BLUE,ORANGE,RGBColor(0x79,0x63,0xA6)][i],str(i+1),15)
    add_footer(s,26)

    # 27.
    s=prs.slides[26]; add_title(s,"gpSLDS places a structured GP prior on continuous-time flow",27,"gpSLDS")
    add_formula(s,"dzₜ = f(zₜ) dt + Σ¹ᐟ² dWₜ",.85,1.38,5.55,.85,25)
    add_formula(s,"fₖ(·) ~ GP(0, κSSL(·,·))",.85,2.58,5.55,.85,24)
    add_bullets(s,["GP-SDE models nonlinear vector fields","SSL kernel encourages local linearity","Smooth boundaries replace hard switching artifacts","Posterior uncertainty is defined over dynamics"],.92,3.78,5.55,2.25,17,CHARCOAL,9)
    add_image(s,assets["gpslds"],6.78,1.15,5.85,5.4)
    add_footer(s,27,"Hu et al., NeurIPS (2024)")

    # 28.
    s=prs.slides[27]; add_title(s,"Local linearity turns a nonlinear field into inspectable pieces",28,"gpSLDS")
    add_image(s,assets["gpslds"],.55,1.15,6.55,5.5)
    add_formula(s,"ż ≈ Aⱼ z + bⱼ",7.55,1.45,4.7,.8,26)
    add_card(s,"Local operator Aⱼ","Summarizes contraction, expansion, rotation, and timescale near a region.",7.55,2.58,4.7,1.3,GREEN)
    add_card(s,"Smooth regime weight","Describes where a local linear template is relevant—without a standard HMM state posterior.",7.55,4.18,4.7,1.55,BLUE)
    add_footer(s,28,"The GP models the vector field; the SSL kernel encodes smoothly switching linear structure.")

    # 29.
    s=prs.slides[28]; add_title(s,"Use continuous-time stability criteria",29,"gpSLDS")
    add_image(s,assets["stability"],.55,1.25,6.6,4.95)
    add_card(s,"Re(λ) < 0","Locally stable direction",7.65,1.45,4.65,1.15,GREEN,"−")
    add_card(s,"Re(λ) > 0","Locally unstable direction",7.65,2.92,4.65,1.15,RED,"+")
    add_card(s,"Re(λ) = 0","Marginal direction; inspect nonlinear terms and noise.",7.65,4.39,4.65,1.35,ORANGE,"0")
    add_footer(s,29,"The familiar |λ|<1 criterion applies to discrete-time maps, not continuous-time flow.")

    # 30.
    s=prs.slides[29]; add_title(s,"gpSLDS uses sparse variational EM—not VAE amortization",30,"gpSLDS")
    add_formula(s,"q(z,f,u) = q(z) ∏ₖ p(fₖ | uₖ, Θ) q(uₖ)",1.05,1.35,11.15,.9,24)
    steps=[("1","Latent path","Update q(z)"),("2","Inducing variables","Update q(u)"),("3","Hyperparameters","Optimize SSL kernel Θ"),("4","Dynamics posterior","Evaluate q(f(z*))")]
    for i,(n,t,b) in enumerate(steps):
        x=.55+i*3.15; add_rect(s,x,2.85,2.85,2.25,WHITE,[BLUE,GREEN,ORANGE,RGBColor(0x79,0x63,0xA6)][i],True,1.1); add_pill(s,n,x+.18,3.05,.48,[BLUE,GREEN,ORANGE,RGBColor(0x79,0x63,0xA6)][i],WHITE,13); add_text(s,t,x+.78,3.04,1.85,.4,16,True,CHARCOAL); add_text(s,b,x+.22,3.82,2.4,.45,16,False,MID,PP_ALIGN.CENTER)
        if i<3:add_arrow(s,x+2.88,4.0,x+3.08,4.0,GREEN,1.6)
    add_text(s,"Sparse inducing points make GP inference tractable; quadrature remains the high-dimensional bottleneck.",.8,5.75,11.7,.55,17,True,ORANGE,PP_ALIGN.CENTER)
    add_footer(s,30,"Hu et al., NeurIPS (2024)")

    # 31.
    s=prs.slides[30]; add_title(s,"What the gpSLDS paper actually evaluated",31,"EVIDENCE")
    exp=[("Synthetic flow","30 trials · 50 output dimensions","Recover trajectories, boundaries,\nfixed points, and uncertainty",BLUE),("Mouse VMH aggression","104 neurons · calcium imaging","Approximate line-attractor region\nwith posterior support",GREEN),("Monkey LIP decisions","58 neurons · 50 trials","2D dynamics · 2 local regimes\nknown pulse input",ORANGE)]
    for i,(t,meta,res,c) in enumerate(exp):
        x=.55+i*4.22; add_rect(s,x,1.30,3.8,4.65,WHITE,c,True,1.2); add_text(s,t,x+.2,1.62,3.4,.55,20,True,c,PP_ALIGN.CENTER); add_pill(s,meta,x+.28,2.48,3.24,LIGHT,CHARCOAL,11.5); add_text(s,res,x+.34,3.28,3.12,1.45,12.5,False,CHARCOAL,PP_ALIGN.CENTER); add_text(s,"paper-supported",x+.55,5.18,2.7,.35,12,True,GREEN_DARK,PP_ALIGN.CENTER)
    add_footer(s,31,"Hu et al., NeurIPS (2024) · no LFADS comparison or 8D motor-reaching result is claimed.")

    # 32.
    s=prs.slides[31]; add_title(s,"Interpretability comes with computational limits",32,"gpSLDS")
    limits=[("Latent dimension","Quadrature memory scales poorly; >≈3 dimensions was difficult in the paper.",ORANGE),("Time discretization","Small steps may be required for stable continuous-time inference.",BLUE),("Variational bias","KL-based VI can underestimate posterior variance.",RGBColor(0x79,0x63,0xA6)),("Causality","Local regimes describe flow; they do not establish circuit connectivity.",RED)]
    for i,(t,b,c) in enumerate(limits):add_card(s,t,b,.65+(i%2)*6.12,1.25+(i//2)*2.25,5.75,1.72,c,str(i+1),16.5)
    add_text(s,"A structured prior makes questions inspectable—not automatically true.",1.05,5.95,11.2,.5,21,True,GREEN_DARK,PP_ALIGN.CENTER)
    add_footer(s,32,"Hu et al., NeurIPS (2024), Discussion")

    # 33.
    s=prs.slides[32]; add_title(s,"The model family is a ladder of assumptions—not a leaderboard",33,"MODEL CHOICE")
    chain=[("Mixture","groups"),("HMM","discrete time"),("LDS","linear flow"),("VAE","nonlinear code"),("LFADS","RNN dynamics"),("gpSLDS","structured GP flow")]
    for i,(m,r) in enumerate(chain):
        x=.35+i*2.16; c=[MID,BLUE,GREEN,ORANGE,RGBColor(0x79,0x63,0xA6),GREEN_DARK][i]; add_rect(s,x,1.55,1.72,2.0,WHITE,c,True,1.1); add_text(s,m,x+.08,1.88,1.56,.45,18,True,c,PP_ALIGN.CENTER); add_text(s,r,x+.12,2.65,1.48,.48,13.5,False,MID,PP_ALIGN.CENTER)
        if i<5:add_arrow(s,x+1.75,2.55,x+2.08,2.55,GREEN,1.8)
    add_rect(s,.65,4.25,12.0,1.45,GREEN_LIGHT,GREEN,True,1.0)
    add_text(s,"Choose the simplest model whose assumptions answer the scientific question.",.95,4.58,11.4,.45,23,True,GREEN_DARK,PP_ALIGN.CENTER)
    add_text(s,"More flexibility → more ways to fit the data → stronger validation required",1.0,5.24,11.3,.35,15,False,MID,PP_ALIGN.CENTER)
    add_footer(s,33)

    # 34.
    s=prs.slides[33]; add_title(s,"A validation ladder: from code to scientific claim",34,"VALIDATION")
    levels=[("1","Implementation","Synthetic recovery · finite objectives"),("2","Generalization","Held-out likelihood · co-smoothing · decoding"),("3","Robustness","Seeds · latent dimension · priors"),("4","Interpretation","Fixed points · regimes · inferred inputs"),("5","Falsification","Independent perturbation or prediction")]
    widths=[10.8,9.5,8.2,6.9,5.6]
    for i,(n,t,b) in enumerate(levels):
        w=widths[i]; x=(13.33-w)/2; y=1.05+i*1.08; c=[RGBColor(0xD9,0xEA,0xE1),RGBColor(0xB9,0xDA,0xC7),RGBColor(0x8C,0xC5,0xA6),RGBColor(0x4E,0xA6,0x78),GREEN][i]
        sh=add_rect(s,x,y,w,.82,c,c,True,0); add_text(s,n,x+.15,y+.20,.45,.32,13,True,WHITE if i>2 else GREEN_DARK,PP_ALIGN.CENTER); add_text(s,t,x+.72,y+.15,2.05,.34,16,True,WHITE if i>2 else CHARCOAL); add_text(s,b,x+2.85,y+.15,w-3.1,.34,14.5,False,WHITE if i>2 else CHARCOAL)
    add_text(s,"A compelling latent plot is exploratory evidence—not a mechanism by itself.",1.0,6.47,11.3,.42,18,True,ORANGE,PP_ALIGN.CENTER)
    add_footer(s,34)

    # 35 takeaways; clear old images.
    s=prs.slides[34]; clear_slide(s); add_title(s,"Take-home: the latent state is only as useful as its tests",35,"SUMMARY")
    takes=[("01","Observations motivate—but do not prove—a latent dynamical description."),("02","HMMs and LDS models show when exact inference is possible."),("03","VI extends inference to nonlinear and nonconjugate systems."),("04","LFADS prioritizes flexible recurrent dynamics; gpSLDS exposes local structure and uncertainty."),("05","Interpretation requires prediction, robustness, and independent experimental evidence.")]
    for i,(n,t) in enumerate(takes):
        y=1.12+i*1.02; add_pill(s,n,.72,y,.62,GREEN,WHITE,12); add_text(s,t,1.58,y-.02,10.95,.52,18.5,True if i==4 else False,CHARCOAL if i<4 else GREEN_DARK)
    add_footer(s,35)

    # 36.
    s=prs.slides[35]; clear_slide(s,False)
    add_text(s,"THANK YOU",1.0,1.45,11.3,.75,36,True,GREEN,PP_ALIGN.CENTER)
    add_text(s,"Which assumptions are necessary for the scientific claim—\nand how would we test whether they are wrong?",1.15,2.75,11.0,1.35,24,True,CHARCOAL,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    add_text(s,"Questions?",1.0,5.25,11.3,.6,28,True,GREEN_DARK,PP_ALIGN.CENTER)
    add_text(s,"Chunyu Qu · Latent Variable Models for Neural Data and Dynamics",1.0,6.65,11.3,.3,11,False,MID,PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f"Built {len(prs.slides)} slides: {OUTPUT}")
    return OUTPUT


if __name__ == "__main__":
    build_deck()
