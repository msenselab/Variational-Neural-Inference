from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_working_03.pptx"
OUTPUT = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_working_04.pptx"
ASSETS = ROOT / "Presentation" / "tutorial_assets"

GREEN = RGBColor(0x00, 0x8A, 0x45)
GREEN_DARK = RGBColor(0x00, 0x68, 0x36)
GREEN_LIGHT = RGBColor(0xE8, 0xF4, 0xEE)
BLUE = RGBColor(0x3F, 0x78, 0xB5)
BLUE_LIGHT = RGBColor(0xED, 0xF3, 0xFA)
ORANGE = RGBColor(0xEE, 0x77, 0x2F)
ORANGE_LIGHT = RGBColor(0xFD, 0xF3, 0xED)
CHARCOAL = RGBColor(0x33, 0x35, 0x37)
MID = RGBColor(0x68, 0x6D, 0x71)
LIGHT = RGBColor(0xF3, 0xF5, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC9, 0x43, 0x43)
FONT = "Arial"
MATH = "Cambria Math"


def clear(slide):
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def rect(slide, x, y, w, h, fill=WHITE, line=None, rounded=True, width=1.0):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(width)
    if rounded:
        try: sh.adjustments[0] = 0.06
        except Exception: pass
    return sh


def text(slide, value, x, y, w, h, size=18, color=CHARCOAL, bold=False,
         align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(.03)
    tf.margin_top = tf.margin_bottom = Inches(.01)
    p = tf.paragraphs[0]; p.text = value; p.alignment = align
    p.font.name = font; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    return box


def formula(slide, value, x, y, w, h, size=22, fill=LIGHT, line=None, color=CHARCOAL):
    rect(slide, x, y, w, h, fill, line, True, .8)
    return text(slide, value, x+.12, y+.06, w-.24, h-.12, size, color, False, PP_ALIGN.CENTER, MATH)


def title(slide, value, number, section):
    text(slide, f"{number:02d}", .30, .13, .48, .30, 11, GREEN, True, PP_ALIGN.CENTER)
    text(slide, value, .86, .08, 11.6, .58, 27, CHARCOAL, True)
    text(slide, section.upper(), 10.85, .19, 2.00, .23, 8.5, GREEN, True, PP_ALIGN.RIGHT)
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(.76), Inches(13.333), Inches(.76))
    ln.line.color.rgb = GREEN; ln.line.width = Pt(1.7)


def footer(slide, number, source=""):
    if source: text(slide, source, .42, 7.08, 11.7, .20, 7.2, MID)
    text(slide, str(number), 12.48, 7.08, .32, .20, 8, MID, False, PP_ALIGN.RIGHT)


def arrow(slide, x1, y1, x2, y2, color=GREEN, width=2.0):
    sh = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    sh.line.color.rgb = color; sh.line.width = Pt(width); sh.line.end_arrowhead = True
    return sh


def image(slide, path, x, y, w, h=None):
    if h is None: return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def card(slide, heading, body, x, y, w, h, accent=GREEN, body_size=15):
    rect(slide, x, y, w, h, WHITE, RGBColor(0xD7,0xDD,0xDA), True, .8)
    bar = rect(slide, x, y, .08, h, accent, None, False)
    text(slide, heading, x+.22, y+.13, w-.4, .38, 17, CHARCOAL, True)
    text(slide, body, x+.22, y+.61, w-.42, h-.72, body_size, MID)


def build():
    shutil.copy2(SOURCE, OUTPUT)
    prs = Presentation(OUTPUT)

    # 11 — the posterior is the computational target.
    s=prs.slides[10]; clear(s); title(s,"From a generative story to a posterior",11,"EXACT INFERENCE")
    text(s,"The model specifies how hidden states generate data:",.75,1.12,5.8,.42,20,CHARCOAL,True)
    formula(s,"p(z₁:T, y₁:T) = p(z₁) ∏ₜ₌₂ᵀ p(zₜ | zₜ₋₁) ∏ₜ₌₁ᵀ p(yₜ | zₜ)",.72,1.68,11.9,.88,23,GREEN_LIGHT,GREEN)
    arrow(s,6.67,2.74,6.67,3.20,GREEN,2.2)
    text(s,"Inference reverses that story:",.75,3.18,5.8,.42,20,CHARCOAL,True)
    formula(s,"p(z₁:T | y₁:T) = p(z₁:T, y₁:T) / p(y₁:T)",1.55,3.74,10.2,.85,25,BLUE_LIGHT,BLUE)
    formula(s,"p(y₁:T) = ∫ p(z₁:T, y₁:T) dz₁:T",3.25,4.91,6.85,.72,21,ORANGE_LIGHT,ORANGE)
    text(s,"The hard part is the denominator: summing or integrating over every possible latent trajectory.",1.05,5.95,11.2,.54,18,ORANGE,True,PP_ALIGN.CENTER)
    footer(s,11)

    # 12 — exact HMM recursion.
    s=prs.slides[11]; clear(s); title(s,"HMM: exact inference without enumerating Kᵀ sequences",12,"EXACT INFERENCE")
    formula(s,"αₜ(k) = p(y₁:ₜ, zₜ = k)",.62,1.26,4.00,.78,24,BLUE_LIGHT,BLUE)
    formula(s,"αₜ(k) = p(yₜ | zₜ=k) ∑ⱼ p(zₜ=k | zₜ₋₁=j) αₜ₋₁(j)",4.90,1.26,7.80,.78,20,GREEN_LIGHT,GREEN)
    steps=[("PAST MESSAGE","αₜ₋₁ summarizes all earlier evidence",BLUE),
           ("TRANSITION","mix probability mass across states",GREEN),
           ("EMISSION","weight states by the new observation",ORANGE)]
    for i,(h,b,c) in enumerate(steps):
        x=.62+i*4.18; card(s,h,b,x,2.55,3.75,1.45,c,14)
        if i<2: arrow(s,x+3.78,3.28,x+4.05,3.28,RGBColor(0xA8,0xB4,0xAE),1.5)
    rect(s,1.02,4.48,11.25,1.30,GREEN_LIGHT,GREEN,True,1.0)
    text(s,"Naive enumeration",1.30,4.73,2.80,.35,16,MID,True,PP_ALIGN.CENTER)
    text(s,"O(Kᵀ)",1.30,5.15,2.80,.38,23,RED,True,PP_ALIGN.CENTER,MATH)
    text(s,"dynamic programming",5.20,4.73,2.80,.35,16,MID,True,PP_ALIGN.CENTER)
    text(s,"O(TK²)",5.20,5.15,2.80,.38,23,GREEN_DARK,True,PP_ALIGN.CENTER,MATH)
    text(s,"Exact because a finite state space can be summed efficiently.",8.20,4.85,3.65,.55,17,CHARCOAL,True,PP_ALIGN.CENTER)
    footer(s,12,"Tutorial: Part2_Dynamics/03a_hmm_foundations.ipynb")

    # 13 — Kalman update as precision-weighted prediction error.
    s=prs.slides[12]; clear(s); title(s,"Kalman filtering is precision-weighted belief updating",13,"EXACT INFERENCE")
    card(s,"1  PREDICT","ẑₜ|ₜ₋₁ = A ẑₜ₋₁|ₜ₋₁\nPₜ|ₜ₋₁ = A Pₜ₋₁|ₜ₋₁ Aᵀ + Q",.55,1.30,3.75,2.20,BLUE,18)
    card(s,"2  COMPARE","prediction error\neₜ = yₜ − C ẑₜ|ₜ₋₁",4.78,1.30,3.75,2.20,ORANGE,18)
    card(s,"3  UPDATE","ẑₜ|ₜ = ẑₜ|ₜ₋₁ + Kₜ eₜ",9.00,1.30,3.75,2.20,GREEN,18)
    arrow(s,4.33,2.40,4.70,2.40,GREEN,1.8); arrow(s,8.56,2.40,8.92,2.40,GREEN,1.8)
    formula(s,"Kₜ = Pₜ|ₜ₋₁ Cᵀ (C Pₜ|ₜ₋₁ Cᵀ + R)⁻¹",2.10,4.05,9.15,.78,23,LIGHT,RGBColor(0xD0,0xD6,0xD3))
    text(s,"new belief  =  prediction  +  confidence in the evidence × prediction error",1.00,5.15,11.35,.58,22,GREEN_DARK,True,PP_ALIGN.CENTER)
    text(s,"The Kalman gain is not an arbitrary learning rate; it follows from relative uncertainty.",1.20,5.95,10.95,.45,16,MID,False,PP_ALIGN.CENTER)
    footer(s,13,"Connection to the group work: precision-weighted updating across trials")

    # 14 — boundary of analytic inference.
    s=prs.slides[13]; clear(s); title(s,"Exact inference depends on restrictive closure properties",14,"EXACT INFERENCE")
    rect(s,.62,1.25,5.85,4.65,GREEN_LIGHT,GREEN,True,1.0)
    text(s,"EXACT / ANALYTIC",.95,1.58,5.20,.36,17,GREEN_DARK,True,PP_ALIGN.CENTER)
    text(s,"HMM\nfinite discrete states\n\nLDS / Kalman\nlinear maps + Gaussian noise\n\nConjugate exponential-family models",1.18,2.18,4.72,2.72,18,CHARCOAL,False,PP_ALIGN.CENTER)
    rect(s,6.86,1.25,5.85,4.65,ORANGE_LIGHT,ORANGE,True,1.0)
    text(s,"TYPICALLY INTRACTABLE",7.18,1.58,5.20,.36,17,ORANGE,True,PP_ALIGN.CENTER)
    text(s,"nonlinear decoder\n\nnonlinear dynamics\n\nPoisson / nonconjugate likelihood\n\nneural-network parameterization",7.40,2.18,4.75,2.72,18,CHARCOAL,False,PP_ALIGN.CENTER)
    text(s,"If the posterior cannot be computed, learn a tractable approximation to it.",1.08,6.12,11.15,.44,20,GREEN_DARK,True,PP_ALIGN.CENTER)
    footer(s,14)

    # 15 — VAE section divider.
    s=prs.slides[14]; clear(s); rect(s,0,0,13.333,7.5,CHARCOAL,None,False)
    rect(s,.72,1.10,.15,4.95,GREEN,None,False)
    text(s,"PART III",1.15,1.22,2.0,.40,15,GREEN,True)
    text(s,"VAE Speedrun",1.12,2.00,10.8,.90,39,WHITE,True)
    text(s,"From an intractable posterior to an amortized inference network",1.15,3.10,10.3,.65,23,RGBColor(0xD7,0xDD,0xDA),False)
    formula(s,"generative model  pθ(x | z)     ⇄     inference model  qφ(z | x)",1.15,4.38,10.65,.85,22,RGBColor(0x3D,0x42,0x43),GREEN,WHITE)
    text(s,"Eight slides: derive it → implement it → inspect what it learned",1.16,5.62,10.4,.42,16,RGBColor(0xBFC,0xC8,0xC3) if False else RGBColor(0xBF,0xC8,0xC3),False)
    text(s,"15",12.48,7.08,.32,.20,8,RGBColor(0xBF,0xC8,0xC3),False,PP_ALIGN.RIGHT)

    # 16 — variational approximation.
    s=prs.slides[15]; clear(s); title(s,"The posterior is easy to write—and often impossible to compute",16,"VAE SPEEDRUN")
    formula(s,"pθ(z | x) = pθ(x,z) / ∫ pθ(x,z) dz",1.05,1.32,11.20,1.00,27,ORANGE_LIGHT,ORANGE)
    text(s,"Choose a tractable family",.85,2.80,3.50,.38,18,CHARCOAL,True,PP_ALIGN.CENTER)
    formula(s,"qφ(z | x)",1.25,3.34,2.70,.85,26,BLUE_LIGHT,BLUE)
    arrow(s,4.25,3.78,5.20,3.78,GREEN,2.2)
    text(s,"optimize φ",4.25,3.20,.95,.35,13,MID,True,PP_ALIGN.CENTER)
    formula(s,"qφ(z | x) ≈ pθ(z | x)",5.35,3.34,3.60,.85,25,GREEN_LIGHT,GREEN)
    arrow(s,9.10,3.78,10.05,3.78,GREEN,2.2)
    text(s,"reuse",9.08,3.20,.95,.35,13,MID,True,PP_ALIGN.CENTER)
    rect(s,10.20,3.20,2.15,1.20,WHITE,BLUE,True,1.1)
    text(s,"new data x\n→ posterior",10.37,3.36,1.80,.78,16,BLUE,True,PP_ALIGN.CENTER)
    text(s,"Amortization replaces a fresh optimization for every observation with one learned inference rule.",1.10,5.45,11.05,.60,20,GREEN_DARK,True,PP_ALIGN.CENTER)
    footer(s,16)

    # 17 — ELBO derivation.
    s=prs.slides[16]; clear(s); title(s,"The ELBO follows from one insertion and Jensen's inequality",17,"VAE SPEEDRUN")
    formula(s,"log pθ(x) = log ∫ pθ(x,z) dz",.70,1.12,5.80,.66,20,LIGHT,None)
    formula(s,"= log ∫ qφ(z|x)  pθ(x,z) / qφ(z|x) dz",6.82,1.12,5.80,.66,19,LIGHT,None)
    arrow(s,6.66,1.45,6.76,1.45,GREEN,1.5)
    formula(s,"≥  Eₙ [ log pθ(x,z) − log qφ(z|x) ]",1.25,2.20,10.85,.78,23,BLUE_LIGHT,BLUE)
    text(s,"Jensen",10.80,1.94,1.10,.26,11,BLUE,True,PP_ALIGN.CENTER)
    formula(s,"ELBO = Eₙ[log pθ(x|z)] − KL(qφ(z|x) || p(z))",1.02,3.38,11.30,.92,25,GREEN_LIGHT,GREEN)
    card(s,"RECONSTRUCTION","Can the sampled latent state explain x?",1.05,4.70,5.35,1.20,BLUE,15)
    card(s,"REGULARIZATION","Does the posterior remain compatible with the prior?",6.93,4.70,5.35,1.20,ORANGE,15)
    text(s,"log pθ(x) = ELBO + KL(qφ(z|x) || pθ(z|x))",2.35,6.15,8.65,.38,16,MID,False,PP_ALIGN.CENTER,MATH)
    footer(s,17,"Derivation anchor: Part3_Variational/05_variational_em.ipynb")

    # 18 — architecture from the executable tutorial.
    s=prs.slides[17]; clear(s); title(s,"A VAE learns an inverse map and a generative map together",18,"VAE SPEEDRUN")
    image(s,ASSETS/"vae_cell4_0.png",.50,1.08,12.35,4.78)
    rect(s,.95,5.98,11.42,.63,GREEN_LIGHT,GREEN,True,.8)
    text(s,"encoder qφ(z|x)  →  stochastic latent code z  →  decoder pθ(x|z)",1.12,6.08,11.05,.36,19,GREEN_DARK,True,PP_ALIGN.CENTER,MATH)
    footer(s,18,"Executable figure: Part3_Variational/04_standard_vae.ipynb")

    # 19 — encoder distribution.
    s=prs.slides[18]; clear(s); title(s,"The encoder predicts a distribution—not a single embedding",19,"VAE SPEEDRUN")
    formula(s,"qφ(z | x) = N(μφ(x), diag σφ²(x))",.78,1.28,6.05,.86,25,BLUE_LIGHT,BLUE)
    card(s,"MEAN  μφ(x)","Where the posterior is centered in latent space.",.82,2.55,2.82,1.40,BLUE,15)
    card(s,"SCALE  σφ(x)","How uncertain the encoder is about this observation.",3.98,2.55,2.82,1.40,ORANGE,15)
    image(s,ASSETS/"vae_cell24_0.png",7.25,1.20,5.15,4.82)
    text(s,"Uncertainty is part of the representation; it is not added after training.",.92,5.25,5.70,.65,20,GREEN_DARK,True,PP_ALIGN.CENTER)
    footer(s,19,"Tutorial latent space: color tracks the original x-coordinate")

    # 20 — reparameterization.
    s=prs.slides[19]; clear(s); title(s,"Reparameterization moves randomness outside the network",20,"VAE SPEEDRUN")
    rect(s,.62,1.30,4.15,3.85,ORANGE_LIGHT,ORANGE,True,1.0)
    text(s,"DIRECT SAMPLING",.92,1.64,3.55,.35,16,ORANGE,True,PP_ALIGN.CENTER)
    formula(s,"z ∼ N(μφ(x), σφ²(x))",1.02,2.38,3.35,.82,22,WHITE,None)
    text(s,"The sampling operation appears to block a simple pathwise gradient.",1.05,3.66,3.28,.82,16,MID,False,PP_ALIGN.CENTER)
    arrow(s,4.95,3.18,5.65,3.18,GREEN,2.3)
    rect(s,5.82,1.30,6.88,3.85,GREEN_LIGHT,GREEN,True,1.0)
    text(s,"REPARAMETERIZED",6.15,1.64,6.20,.35,16,GREEN_DARK,True,PP_ALIGN.CENTER)
    formula(s,"ε ∼ N(0,I)",6.32,2.25,2.05,.72,20,WHITE,None)
    formula(s,"z = μφ(x) + σφ(x) ⊙ ε",8.62,2.25,3.58,.72,22,WHITE,None)
    text(s,"Randomness lives in ε; gradients flow through μφ and σφ.",6.40,3.62,5.70,.74,18,GREEN_DARK,True,PP_ALIGN.CENTER)
    text(s,"One algebraic rewrite makes stochastic latent-variable learning compatible with backpropagation.",.95,5.70,11.45,.62,20,CHARCOAL,True,PP_ALIGN.CENTER)
    footer(s,20)

    # 21 — likelihood choice.
    s=prs.slides[20]; clear(s); title(s,"The reconstruction loss is a likelihood assumption",21,"VAE SPEEDRUN")
    rows=[("CONTINUOUS DATA","Gaussian","negative log likelihood ∝ MSE",BLUE),
          ("BINARY DATA","Bernoulli","negative log likelihood = BCE",ORANGE),
          ("SPIKE COUNTS","Poisson","yₜ,ₙ ∼ Poisson(λₜ,ₙ)",GREEN)]
    for i,(a,b,c,col) in enumerate(rows):
        y=1.32+i*1.48; rect(s,.72,y,11.92,1.06,WHITE,RGBColor(0xD7,0xDD,0xDA),True,.8)
        rect(s,.72,y,.10,1.06,col,None,False)
        text(s,a,1.02,y+.23,2.45,.38,15,col,True)
        text(s,b,3.75,y+.23,2.20,.38,19,CHARCOAL,True)
        text(s,c,6.08,y+.21,5.90,.43,18,MID,False,PP_ALIGN.CENTER,MATH)
    formula(s,"decoder output → parameters of pθ(x | z), not merely a reconstructed point",1.20,5.92,10.95,.68,20,GREEN_LIGHT,GREEN,GREEN_DARK)
    footer(s,21)

    # 22 — executable tutorial results.
    s=prs.slides[21]; clear(s); title(s,"VAE speedrun: train, reconstruct, inspect, sample",22,"VAE SPEEDRUN")
    image(s,ASSETS/"vae_cell20_0.png",.42,1.14,4.00,2.18)
    image(s,ASSETS/"vae_cell22_0.png",4.66,1.14,4.00,2.18)
    image(s,ASSETS/"vae_cell26_0.png",8.90,1.14,4.00,2.18)
    caps=[("1","OPTIMIZE","ELBO converges"),("2","RECONSTRUCT","clusters are recovered"),("3","GENERATE","sample z ∼ N(0,I)")]
    for i,(n,h,b) in enumerate(caps):
        x=.55+i*4.22; rect(s,x,3.58,3.75,1.12,WHITE,RGBColor(0xD7,0xDD,0xDA),True,.8)
        rect(s,x+.18,3.82,.45,.45,[BLUE,GREEN,ORANGE][i],None,True)
        text(s,n,x+.18,3.84,.45,.30,12,WHITE,True,PP_ALIGN.CENTER)
        text(s,h,x+.78,3.70,2.60,.30,14,[BLUE,GREEN,ORANGE][i],True)
        text(s,b,x+.78,4.08,2.65,.30,13,MID)
    rect(s,1.05,5.24,11.22,1.08,GREEN_LIGHT,GREEN,True,.9)
    text(s,"A trained VAE is simultaneously",1.30,5.47,2.75,.35,16,MID,True,PP_ALIGN.CENTER)
    text(s,"a compressor",4.08,5.44,2.05,.40,20,BLUE,True,PP_ALIGN.CENTER)
    text(s,"a density model",6.35,5.44,2.25,.40,20,GREEN_DARK,True,PP_ALIGN.CENTER)
    text(s,"a generator",8.88,5.44,2.05,.40,20,ORANGE,True,PP_ALIGN.CENTER)
    footer(s,22,"Results extracted from Part3_Variational/04_standard_vae.ipynb")

    # 23 — sequential bridge.
    s=prs.slides[22]; clear(s); title(s,"A standard VAE learns a code—not yet a dynamical system",23,"SEQUENTIAL VAE")
    rect(s,.62,1.22,4.05,2.10,BLUE_LIGHT,BLUE,True,1.0)
    text(s,"STANDARD VAE",.92,1.54,3.45,.34,16,BLUE,True,PP_ALIGN.CENTER)
    formula(s,"qφ(z | x)",1.12,2.17,3.05,.72,24,WHITE,None)
    arrow(s,4.88,2.27,5.42,2.27,GREEN,2.0)
    rect(s,5.62,1.22,7.08,2.10,GREEN_LIGHT,GREEN,True,1.0)
    text(s,"SEQUENTIAL VAE",5.95,1.54,6.42,.34,16,GREEN_DARK,True,PP_ALIGN.CENTER)
    formula(s,"qφ(z₁:T | y₁:T)",6.02,2.17,2.65,.72,23,WHITE,None)
    formula(s,"p(z₁) ∏ₜ p(zₜ|zₜ₋₁) ∏ₜ p(yₜ|zₜ)",8.90,2.17,3.35,.72,18,WHITE,None)
    formula(s,"L = Eₙ[∑ₜ log pθ(yₜ|zₜ)] − KL(qφ(z₁:T|y₁:T) || pθ(z₁:T))",1.02,3.78,11.30,.82,21,LIGHT,RGBColor(0xD0,0xD6,0xD3))
    text(s,"Replace an independent latent code with a latent trajectory and a dynamical prior.",1.08,4.93,11.18,.48,20,GREEN_DARK,True,PP_ALIGN.CENTER)
    rect(s,2.05,5.70,9.25,.72,ORANGE_LIGHT,ORANGE,True,.8)
    text(s,"Next: LFADS makes that dynamical prior a recurrent neural network.",2.28,5.84,8.80,.36,19,ORANGE,True,PP_ALIGN.CENTER)
    footer(s,23,"Next notebook: Part3_Variational/06_lfads.ipynb")

    prs.save(OUTPUT)
    print(f"Built slides 11–23: {OUTPUT}")


if __name__ == "__main__":
    build()
