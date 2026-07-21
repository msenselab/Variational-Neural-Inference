#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generate professional presentation using template color scheme and structure.
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Template color scheme extracted from template
COLORS = {
    'dark_blue': RGBColor(15, 25, 135),      # #0F1987
    'bright_blue': RGBColor(0, 159, 227),    # #009FE3
    'green': RGBColor(0, 136, 58),           # #00883A
    'purple': RGBColor(140, 64, 145),        # #8C4091
    'orange': RGBColor(241, 135, 0),         # #F18700
    'red': RGBColor(215, 25, 25),            # #D71919
    'dark_gray': RGBColor(35, 35, 35),       # #232323
    'medium_gray': RGBColor(98, 100, 104),   # #626468
    'light_gray': RGBColor(192, 193, 195),   # #C0C1C3
    'very_light_gray': RGBColor(230, 230, 231),  # #E6E6E7
    'almost_white': RGBColor(245, 245, 245), # #F5F5F5
    'white': RGBColor(255, 255, 255),
    'black': RGBColor(0, 0, 0),
}

class SlideBuilder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.slide_num = 0

    def add_title_slide(self, title, subtitle, author, date):
        """Title slide with blue background."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        self.slide_num += 1

        # Dark blue background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['dark_blue']

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = COLORS['white']

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(20)
        subtitle_para.font.color.rgb = COLORS['bright_blue']

        # Metadata
        meta_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        meta_frame = meta_box.text_frame
        meta_para = meta_frame.paragraphs[0]
        meta_para.text = f"{author}  •  {date}"
        meta_para.font.size = Pt(14)
        meta_para.font.color.rgb = COLORS['light_gray']
        meta_para.alignment = PP_ALIGN.CENTER

    def add_content_slide(self, title, content_blocks, bg_color=None):
        """Standard content slide with title bar and multi-block layout."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        self.slide_num += 1

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color or COLORS['almost_white']

        # Title bar (dark blue with bright blue line)
        title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_bar_fill = title_bar.fill
        title_bar_fill.solid()
        title_bar_fill.fore_color.rgb = COLORS['dark_blue']
        title_bar.line.color.rgb = COLORS['bright_blue']
        title_bar.line.width = Pt(3)

        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = COLORS['white']
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Content blocks
        block_height = 1.5  # Height per block
        top_pos = 1.0

        for block in content_blocks:
            self._add_content_block(slide, top_pos, block)
            top_pos += block_height + 0.3

    def _add_content_block(self, slide, top, block):
        """Add a single content block (title + content)."""
        block_title = block.get('title', '')
        block_content = block.get('content', [])
        left = Inches(block.get('left', 0.5))
        width = Inches(block.get('width', 4.5))

        # Block title
        if block_title:
            title_box = slide.shapes.add_textbox(left, Inches(top), width, Inches(0.35))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = block_title
            title_para.font.size = Pt(18)
            title_para.font.bold = True
            title_para.font.color.rgb = COLORS['dark_blue']
            top += 0.35

        # Block content
        content_box = slide.shapes.add_textbox(left, Inches(top), width, Inches(1.2))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, line in enumerate(block_content):
            if i > 0:
                text_frame.add_paragraph()

            para = text_frame.paragraphs[i]
            para.text = line
            para.font.size = Pt(14)
            para.font.color.rgb = COLORS['dark_gray']
            para.space_before = Pt(3)
            para.space_after = Pt(3)

            # Highlight key terms
            if '→' in line or '=' in line:
                para.font.color.rgb = COLORS['purple']

    def add_section_slide(self, section_title):
        """Section divider slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        # Gradient-like effect: dark to purple
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['purple']

        # Large section title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = section_title
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = COLORS['white']
        title_para.alignment = PP_ALIGN.CENTER

    def save(self, output_path):
        """Save presentation."""
        self.prs.save(output_path)
        print(f"✓ Presentation saved: {output_path}")
        print(f"✓ Total slides: {self.slide_num}")

def build_presentation():
    """Build complete presentation."""
    builder = SlideBuilder()

    # === TITLE SLIDE ===
    builder.add_title_slide(
        "Latent Variable Models for Neural Data and Dynamics",
        "From Probabilistic Modeling to Interpretable Neural Dynamics",
        "Chunyu Qu",
        "22.07.2026"
    )

    # === PART 1: OPENING ===
    builder.add_section_slide("Part 1: Understanding Neural Data")

    builder.add_content_slide(
        "The Fundamental Challenge",
        [
            {
                'title': 'What we observe',
                'content': [
                    '• Spike trains: high-dimensional, noisy',
                    '• Variable trial-to-trial',
                    '• Confounded with behavior & environment'
                ]
            },
            {
                'title': 'What we want to understand',
                'content': [
                    '• Hidden population codes',
                    '• Latent dynamical structure',
                    '• Circuit computations'
                ]
            }
        ]
    )

    builder.add_content_slide(
        "The Inferential Gap",
        [
            {
                'title': 'The problem',
                'content': [
                    'Observation ≠ Mechanism',
                    '',
                    'Need to infer hidden states z',
                    'from noisy observations y'
                ]
            },
            {
                'title': 'Statistical solution',
                'content': [
                    'Generative models:',
                    'p(y, z | θ) = p(y | z, θ) · p(z | θ)',
                    '',
                    'Bayesian inference: p(z | y, θ)'
                ]
            }
        ]
    )

    # === PART 2: PROBABILISTIC FRAMEWORK ===
    builder.add_section_slide("Part 2: Probabilistic Modeling")

    builder.add_content_slide(
        "Bayes Rule for Neural Inference",
        [
            {
                'title': 'Generative model',
                'content': [
                    'Likelihood: p(y | z, θ)',
                    '  How does latent z produce y?',
                    '',
                    'Prior: p(z | θ)',
                    '  What structure do latents have?'
                ]
            },
            {
                'title': 'Posterior inference',
                'content': [
                    'p(z | y, θ) ∝ p(y | z) · p(z)',
                    '',
                    'Goal: Learn θ, infer z'
                ]
            }
        ]
    )

    # === PART 3: DISCRETE LATENTS ===
    builder.add_section_slide("Part 3: From Static to Dynamic")

    builder.add_content_slide(
        "Mixture Models: Heterogeneity",
        [
            {
                'title': 'Model',
                'content': [
                    'p(y) = Σ_k π_k N(y | μ_k, Σ_k)',
                    '',
                    'Discrete latent z ∈ {1,...,K}'
                ]
            },
            {
                'title': 'What it captures',
                'content': [
                    '✓ Data heterogeneity',
                    '✓ Multiple regimes',
                    '✗ No temporal structure'
                ]
            }
        ]
    )

    builder.add_content_slide(
        "The EM Algorithm & ELBO",
        [
            {
                'title': 'Challenge',
                'content': [
                    'Exact posterior p(z | y) is',
                    'intractable (sum over K^N terms)'
                ]
            },
            {
                'title': 'Solution: EM Algorithm',
                'content': [
                    'E-step: γ_ik = p(z_i = k | y_i)',
                    'M-step: Update μ_k, π_k',
                    '',
                    'Theoretical foundation: ELBO'
                ]
            }
        ]
    )

    builder.add_content_slide(
        "HMM: Adding Temporal Dependence",
        [
            {
                'title': 'Key innovation',
                'content': [
                    'Markov transition:',
                    'p(z_t | z_{t-1})',
                    '',
                    'Explains sequential patterns'
                ]
            },
            {
                'title': 'Why limited?',
                'content': [
                    '✗ Discrete states too coarse',
                    '✗ Fixed within-state dynamics',
                    '✗ Need continuous latents for',
                    '  rich neural population codes'
                ]
            }
        ]
    )

    # === PART 4: CONTINUOUS LATENTS ===
    builder.add_section_slide("Part 4: Continuous Latent Dynamics")

    builder.add_content_slide(
        "Why Variational Inference?",
        [
            {
                'title': 'The jump to continuous z',
                'content': [
                    'z ∈ ℝ^d (not discrete)',
                    '',
                    'Posterior p(z | y) is doubly',
                    'intractable: continuous + nonlinear'
                ]
            },
            {
                'title': 'Approximate with ELBO',
                'content': [
                    'log p(y) ≥ E_q[log p(y,z)] - KL(q||p)',
                    '',
                    'Find tractable q(z|y)'
                ]
            }
        ]
    )

    builder.add_content_slide(
        "VAE: Amortized Variational Inference",
        [
            {
                'title': 'Architecture',
                'content': [
                    'Encoder: q_φ(z | y) → μ, σ',
                    'Decoder: p_θ(y | z)',
                    '',
                    'Reparameterization trick'
                ]
            },
            {
                'title': 'Key innovation: Amortization',
                'content': [
                    '✓ One encoder for all data',
                    '✓ No per-datum optimization',
                    '✓ Continuous latent space',
                    '✓ Generative model'
                ]
            }
        ]
    )

    builder.add_content_slide(
        "VAE Limitation for Sequences",
        [
            {
                'title': 'Problem',
                'content': [
                    'VAE assumes i.i.d. observations',
                    'Ignores temporal structure',
                    'Each y_t independent (given z_t)'
                ]
            },
            {
                'title': 'Solution needed',
                'content': [
                    'Sequential VAE with dynamics prior:',
                    'p(z_1:T) = p(z_1) ∏ p(z_t|z_{t-1})',
                    '',
                    '→ LFADS'
                ]
            }
        ]
    )

    # === PART 5: LFADS ===
    builder.add_section_slide("Part 5: LFADS - Neural Sequences")

    builder.add_content_slide(
        "LFADS: The Leap to Spike Trains",
        [
            {
                'title': 'What LFADS infers',
                'content': [
                    '• z_0: initial population code',
                    '• z_{1:T}: latent trajectory',
                    '• Generator RNN: dynamics',
                    '• Firing rates: λ_t = exp(Wz_t)'
                ]
            },
            {
                'title': 'Why it works',
                'content': [
                    '✓ Continuous population codes',
                    '✓ Learned, flexible dynamics',
                    '✓ Amortized inference (fast)',
                    '✓ Poisson observation model'
                ]
            }
        ]
    )

    builder.add_content_slide(
        "LFADS Architecture (Simplified)",
        [
            {
                'title': 'Data encoder',
                'content': [
                    'BiRNN(y_{1:T}) → context'
                ]
            },
            {
                'title': 'Posterior inference',
                'content': [
                    'RNN(context) → q(z_{1:T}|y)'
                ]
            },
            {
                'title': 'Generator & decoder',
                'content': [
                    'RNN(z_0) → z_{1:T}',
                    'Linear(z_t) → λ_t',
                    'Poisson(λ_t) → y_t'
                ]
            }
        ]
    )

    builder.add_content_slide(
        "Why LFADS Matters",
        [
            {
                'title': 'Scientific insights',
                'content': [
                    '• Denoised spike rasters',
                    '• Trial-to-trial variability',
                    '• Extrapolate missing data'
                ]
            },
            {
                'title': 'Neural dynamics',
                'content': [
                    '• Learn generator RNN dynamics',
                    '• Analyze latent trajectories',
                    '• (But generator is still black-box)'
                ]
            }
        ]
    )

    # === PART 6: gpSLDS ===
    builder.add_section_slide("Part 6: Interpretable Dynamics")

    builder.add_content_slide(
        "The Interpretability Problem",
        [
            {
                'title': 'LFADS strength',
                'content': [
                    '✓ Flexible, accurate',
                    '✓ Learns complex dynamics'
                ]
            },
            {
                'title': 'LFADS weakness',
                'content': [
                    '✗ Black-box RNN generator',
                    '✗ Hard to reverse-engineer',
                    '✗ Why this solution?'
                ]
            }
        ]
    )

    builder.add_content_slide(
        "gpSLDS: Piecewise Linear Dynamics",
        [
            {
                'title': 'Key idea',
                'content': [
                    'Replace black-box RNN with',
                    'explicit piecewise-linear dynamics:',
                    'z_t = f(z_{t-1}) + noise'
                ]
            },
            {
                'title': 'Each "piece"',
                'content': [
                    '• Local linear regime',
                    '• Simple: z_t = A_k z_{t-1}',
                    '• GP-SDE smoothly interpolates'
                ]
            }
        ]
    )

    builder.add_content_slide(
        "Why Local Linearity?",
        [
            {
                'title': 'Interpretability',
                'content': [
                    '✓ Eigenvalues = stability',
                    '✓ Eigenvectors = directions',
                    '✓ Matches circuit motifs',
                    '✓ Mechanistic insight'
                ]
            },
            {
                'title': 'Expressiveness',
                'content': [
                    '✓ Piecewise linear ≈ nonlinear',
                    '✓ Can represent complex behavior',
                    '✓ Not as flexible as RNN,',
                    '  but more understandable'
                ]
            }
        ]
    )

    builder.add_content_slide(
        "Uncertainty Quantification",
        [
            {
                'title': 'Beyond point estimates',
                'content': [
                    'gpSLDS provides:',
                    '• p(regimes | data)',
                    '• p(dynamics | data)',
                    '• Which are well-determined?'
                ]
            },
            {
                'title': 'Benefits',
                'content': [
                    '✓ Know what we know',
                    '✓ Identify ambiguous regions',
                    '✓ Transfer to new conditions'
                ]
            }
        ]
    )

    # === SUMMARY ===
    builder.add_section_slide("Part 7: Key Takeaways")

    builder.add_content_slide(
        "The Modeling Chain",
        [
            {
                'title': 'Progression',
                'content': [
                    '1. Mixture: discrete, static',
                    '2. HMM: discrete, temporal',
                    '3. VAE: continuous, static',
                    '4. LFADS: continuous, RNN dynamics',
                    '5. gpSLDS: continuous, interpretable'
                ]
            },
            {
                'title': 'Principle',
                'content': [
                    'Increase expressiveness,',
                    'maintain tractability via',
                    'variational inference'
                ]
            }
        ]
    )

    builder.add_content_slide(
        "Open Questions",
        [
            {
                'title': 'Generalization',
                'content': [
                    '• Transfer across animals/conditions?',
                    '• Which parameters stable?'
                ]
            },
            {
                'title': 'Causality & control',
                'content': [
                    '• Infer causal structure?',
                    '• Optimal intervention?',
                    '• Neural prosthetics?'
                ]
            }
        ]
    )

    return builder

if __name__ == '__main__':
    if sys.platform == 'win32':
        import io
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

    print("Building presentation...")
    builder = build_presentation()

    output_path = Path(r"D:\Variational-Neural-Inference\presentation_final.pptx")
    builder.save(output_path)
