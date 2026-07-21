#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Convert markdown presentation to PPTX using a template as style base.
"""

import re
import shutil
import sys
from pathlib import Path

# Fix Windows encoding
if sys.platform == 'win32':
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def parse_markdown(md_path):
    """Parse Slidev markdown into slides."""
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by --- (slide separator)
    raw_slides = content.split('\n---\n')

    slides = []
    for raw_slide in raw_slides:
        # Skip empty slides
        if not raw_slide.strip():
            continue

        # Parse YAML frontmatter if present
        yaml_match = re.match(r'^---\n(.*?)\n---', raw_slide, re.DOTALL)
        body = raw_slide

        if yaml_match:
            yaml_block = yaml_match.group(1)
            body = raw_slide[yaml_match.end():].lstrip('\n')

        # Extract title (first # heading)
        title_match = re.search(r'^#\s+(.+?)$', body, re.MULTILINE)
        title = title_match.group(1) if title_match else ""

        # Extract subtitle (second # heading or ###)
        subtitle_match = re.search(r'^###\s+(.+?)$', body, re.MULTILINE)
        subtitle = subtitle_match.group(1) if subtitle_match else ""

        # Clean body for content
        content_text = body
        if title_match:
            content_text = content_text.replace(f"# {title}", "", 1).lstrip()
        if subtitle_match:
            content_text = content_text.replace(f"### {subtitle}", "", 1).lstrip()

        content_text = content_text.strip()

        slides.append({
            'title': title,
            'subtitle': subtitle,
            'content': content_text,
            'layout': 'title' if subtitle else ('content' if title else 'blank')
        })

    return slides

def create_pptx_from_template(template_path, slides, output_path):
    """
    Create PPTX using template as style base, populate with parsed slides.
    """
    # Load template
    prs = Presentation(template_path)

    # Use template's color scheme and fonts
    template_slide = prs.slides[0] if prs.slides else None

    def add_title_slide(prs, title, subtitle):
        """Add title slide (mimics template slide 1)."""
        blank_slide_layout = prs.slide_layouts[0]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add background or use template's style
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(40, 40, 50)  # Dark background

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Add subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RGBColor(200, 200, 200)

        # Add metadata at bottom
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = "Chunyu Qu  •  22.07.2026"
        date_para.font.size = Pt(12)
        date_para.font.color.rgb = RGBColor(150, 150, 150)
        date_para.alignment = PP_ALIGN.CENTER

        return slide

    def add_content_slide(prs, title, content):
        """Add content slide."""
        blank_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(blank_slide_layout)

        # White background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(40, 40, 50)

        # Separator line
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(0))
        line.line.color.rgb = RGBColor(100, 150, 200)
        line.line.width = Pt(2)

        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        # Parse content markdown (basic)
        lines = content.split('\n')
        for i, line in enumerate(lines):
            if i > 0:
                text_frame.add_paragraph()

            para = text_frame.paragraphs[i]

            # Check for list markers
            if line.strip().startswith('- '):
                para.text = line.strip()[2:]
                para.level = 0
                para.font.size = Pt(16)
            elif line.strip().startswith('  - '):
                para.text = line.strip()[4:]
                para.level = 1
                para.font.size = Pt(14)
            elif re.match(r'^\d+\.\s', line.strip()):
                # Numbered list
                para.text = re.sub(r'^\d+\.\s', '', line.strip())
                para.level = 0
                para.font.size = Pt(16)
            elif '|' in line and '-' in line:
                # Skip table lines for now
                continue
            elif line.strip().startswith('###'):
                para.text = line.replace('###', '').strip()
                para.font.size = Pt(20)
                para.font.bold = True
            elif line.strip().startswith('**'):
                para.text = line.replace('**', '').strip()
                para.font.bold = True
                para.font.size = Pt(16)
            elif '$' in line:
                # Keep math as is (for now)
                para.text = line
                para.font.size = Pt(12)
                para.font.name = 'Courier New'
            else:
                para.text = line
                para.font.size = Pt(16)

            para.font.color.rgb = RGBColor(60, 60, 60)
            para.space_before = Pt(4)
            para.space_after = Pt(4)

        return slide

    # Clear existing slides (except keep layout reference)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # Add slides from parsed content
    for slide_dict in slides:
        if slide_dict['layout'] == 'title':
            add_title_slide(prs, slide_dict['title'], slide_dict['subtitle'])
        elif slide_dict['layout'] == 'content':
            add_content_slide(prs, slide_dict['title'], slide_dict['content'])

    # Save
    prs.save(output_path)
    print(f"✓ PPTX saved to: {output_path}")
    print(f"✓ Total slides: {len(prs.slides)}")

if __name__ == '__main__':
    md_path = Path(__file__).parent.parent / 'presentation.md'
    template_path = Path(r'C:\Users\10993\Desktop\Latent Variable Models for Neural Data and Dynamics .pptx')
    output_path = Path(__file__).parent.parent / 'presentation_generated.pptx'

    print(f"Reading markdown: {md_path}")
    slides = parse_markdown(md_path)
    print(f"Parsed {len(slides)} slides")

    print(f"Using template: {template_path}")
    print(f"Generating PPTX...")
    create_pptx_from_template(template_path, slides, output_path)
