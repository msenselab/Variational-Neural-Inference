from pathlib import Path
import shutil

import matplotlib.pyplot as plt
import numpy as np
from matplotlib.patches import Patch
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from build_parts2_3 import (
    clear, title, footer, rect, text, formula, arrow, image,
    GREEN, GREEN_DARK, GREEN_LIGHT, BLUE, BLUE_LIGHT, ORANGE,
    ORANGE_LIGHT, CHARCOAL, MID, LIGHT, WHITE,
)


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_working_06.pptx"
OUTPUT = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_working_07.pptx"
ASSETS = ROOT / "Presentation" / "season_example"
ASSETS.mkdir(parents=True, exist_ok=True)

SEASON_NAMES = ["Winter", "Spring", "Summer", "Autumn"]
SEASON_COLORS = ["#3F78B5", "#45A778", "#EE772F", "#C39B2E"]


def make_figures():
    rng = np.random.default_rng(27)

    # Finite mixture: shuffled days in temperature-humidity space.
    means = np.array([[2, 78], [12, 66], [25, 48], [14, 72]])
    covs = [np.array([[8, -3], [-3, 35]]), np.array([[10, -4], [-4, 30]]),
            np.array([[9, -5], [-5, 28]]), np.array([[8, 2], [2, 28]])]
    fig, ax = plt.subplots(figsize=(7.2, 4.6))
    for k, (mu, cov, name, col) in enumerate(zip(means, covs, SEASON_NAMES, SEASON_COLORS)):
        pts = rng.multivariate_normal(mu, cov, 130)
        ax.scatter(pts[:, 0], pts[:, 1], s=18, alpha=.64, c=col, label=name, edgecolors="none")
        ax.scatter(*mu, s=150, c=col, marker="X", edgecolors="white", linewidths=1.6)
    ax.set_xlabel("temperature (°C)"); ax.set_ylabel("humidity (%)")
    ax.set_title("The calendar is hidden: infer a component for each day", loc="left", weight="bold")
    ax.legend(frameon=False, ncol=2, loc="lower left")
    ax.spines[["top", "right"]].set_visible(False); ax.grid(alpha=.14)
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")
    fig.tight_layout(); fig.savefig(ASSETS / "season_mixture.png", dpi=210, transparent=False, facecolor="white"); plt.close(fig)

    # HMM: restore temporal ordering and show persistent discrete regimes.
    T = 240
    states = np.repeat(np.arange(4), 60)
    temp_means = np.array([2, 12, 25, 14])
    temps = temp_means[states] + rng.normal(0, 2.2, T)
    temps += 1.4 * np.sin(np.linspace(0, 12*np.pi, T))
    fig, ax = plt.subplots(figsize=(9.0, 4.4))
    for k in range(4):
        mask = states == k
        idx = np.where(mask)[0]
        ax.axvspan(idx[0], idx[-1]+1, color=SEASON_COLORS[k], alpha=.18, lw=0)
    ax.plot(temps, color="#333537", lw=1.15, alpha=.65, label="observed temperature")
    ax.plot(temp_means[states], color="#008A45", lw=2.4, label="state-dependent mean")
    ax.set_xlabel("day"); ax.set_ylabel("temperature (°C)")
    ax.set_title("Temporal order turns independent assignments into persistent regimes", loc="left", weight="bold")
    handles=[Patch(facecolor=c, alpha=.35, label=n) for n,c in zip(SEASON_NAMES,SEASON_COLORS)]
    ax.legend(handles=handles, frameon=False, ncol=4, loc="upper center")
    ax.spines[["top", "right"]].set_visible(False); ax.grid(axis="y", alpha=.14)
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")
    fig.tight_layout(); fig.savefig(ASSETS / "season_hmm.png", dpi=210, transparent=False, facecolor="white"); plt.close(fig)

    # LDS: local level/trend state and noisy observations, inferred with a Kalman filter.
    T = 150
    F = np.array([[1., 1.], [0., 1.]])
    Q = np.diag([.05, .008]); H = np.array([[1., 0.]]); R = np.array([[2.2]])
    latent = np.zeros((T,2)); latent[0] = [8., .07]
    for t in range(1,T):
        latent[t] = F @ latent[t-1] + rng.multivariate_normal(np.zeros(2), Q)
    obs = latent[:,0] + rng.normal(0, np.sqrt(R[0,0]), T)
    m=np.array([7.,0.]); P=np.diag([5.,.3]); means_f=[]; vars_f=[]
    for y in obs:
        mp=F@m; Pp=F@P@F.T+Q
        S=H@Pp@H.T+R; K=Pp@H.T@np.linalg.inv(S)
        m=mp+(K[:,0]*(y-(H@mp)[0])); P=(np.eye(2)-K@H)@Pp
        means_f.append(m.copy()); vars_f.append(P.copy())
    means_f=np.asarray(means_f); vars_f=np.asarray(vars_f); sd=np.sqrt(vars_f[:,0,0])
    fig, axes = plt.subplots(2,1,figsize=(8.6,5.4),sharex=True,gridspec_kw={"height_ratios":[2.2,1]})
    tt=np.arange(T)
    axes[0].scatter(tt,obs,s=12,c="#A9B2AE",alpha=.62,label="noisy thermometer")
    axes[0].plot(tt,latent[:,0],c="#EE772F",lw=1.8,label="true latent level")
    axes[0].plot(tt,means_f[:,0],c="#008A45",lw=2.3,label="Kalman estimate")
    axes[0].fill_between(tt,means_f[:,0]-2*sd,means_f[:,0]+2*sd,color="#008A45",alpha=.13,label="posterior ±2 SD")
    axes[0].set_ylabel("temperature"); axes[0].legend(frameon=False,ncol=2,loc="upper left")
    axes[0].spines[["top","right"]].set_visible(False); axes[0].grid(alpha=.13)
    axes[1].plot(tt,latent[:,1],c="#EE772F",lw=1.5,label="true trend")
    axes[1].plot(tt,means_f[:,1],c="#3F78B5",lw=2,label="inferred trend")
    axes[1].axhline(0,c="#A9B2AE",lw=.8); axes[1].set_ylabel("trend"); axes[1].set_xlabel("day")
    axes[1].legend(frameon=False,ncol=2,loc="upper left"); axes[1].spines[["top","right"]].set_visible(False); axes[1].grid(alpha=.13)
    fig.suptitle("A continuous hidden state separates weather level from trend",x=.08,ha="left",weight="bold")
    fig.patch.set_facecolor("white")
    for ax in axes: ax.set_facecolor("white")
    fig.tight_layout(); fig.savefig(ASSETS / "season_lds.png", dpi=210, transparent=False, facecolor="white"); plt.close(fig)


def add_cycle(slide):
    xs=[.98,2.12,3.26,4.40]
    for i,(name,col,x) in enumerate(zip(SEASON_NAMES,SEASON_COLORS,xs)):
        sh=rect(slide,x,1.42,.92,.66,WHITE,BLUE if i==0 else GREEN,True,1.0)
        # Match the intuitive seasonal palette in the label rather than the border.
        from pptx.dml.color import RGBColor
        c=RGBColor.from_string(col.lstrip("#"))
        sh.line.color.rgb=c
        text(slide,name,x+.04,1.55,.84,.28,12,c,True,PP_ALIGN.CENTER)
        if i<3: arrow(slide,x+.94,1.75,xs[i+1]-.05,1.75,GREEN,1.4)
    arrow(slide,5.34,1.74,5.62,1.74,GREEN,1.4)
    text(slide,"…",5.65,1.49,.40,.30,18,MID,True,PP_ALIGN.CENTER)


def renumber(slides):
    for idx, slide in enumerate(slides, 1):
        if idx < 12:
            continue
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            value=shape.text.strip()
            if not value.isdigit():
                continue
            x=shape.left/914400; y=shape.top/914400
            if (x < .85 and y < .60) or (x > 12.0 and y > 6.80):
                shape.text=str(idx)
                for p in shape.text_frame.paragraphs:
                    p.font.name="Arial"
                    if x < .85:
                        p.font.size=Pt(11); p.font.bold=True; p.font.color.rgb=GREEN
                        p.alignment=PP_ALIGN.CENTER
                        shape.width=Inches(.62)
                    else:
                        p.font.size=Pt(8); p.font.bold=False; p.font.color.rgb=MID
                        p.alignment=PP_ALIGN.RIGHT
                        shape.left=Inches(12.40); shape.width=Inches(.38)


def build():
    make_figures()
    shutil.copy2(SOURCE, OUTPUT)
    prs=Presentation(OUTPUT)
    new=[]
    for _ in range(3):
        sl=prs.slides.add_slide(prs.slide_layouts[0]); clear(sl); new.append(sl)

    # Move the new slides to immediately after slide 11.
    ids=prs.slides._sldIdLst
    new_ids=list(ids)[-3:]
    for item in new_ids: ids.remove(item)
    for offset,item in enumerate(new_ids): ids.insert(11+offset,item)

    # Slide 12: finite mixture.
    s=prs.slides[11]; rect(s,0,0,13.333,7.5,WHITE,None,False); title(s,"Finite mixture: four seasons, but no calendar",12,"INTUITIVE EXAMPLE")
    text(s,"Imagine a bag of shuffled daily records.",.62,1.10,4.55,.42,19,CHARCOAL,True)
    formula(s,"zₙ ∈ {winter, spring, summer, autumn}",.62,1.68,4.45,.72,19,BLUE_LIGHT,BLUE)
    formula(s,"xₙ = [temperature, humidity]ᵀ",.62,2.65,4.45,.72,20,GREEN_LIGHT,GREEN)
    formula(s,"xₙ | zₙ=k ∼ N(μₖ, Σₖ)",.62,3.62,4.45,.72,21,LIGHT,None)
    text(s,"Each day gets a component probability,\nbut yesterday tells us nothing about today.",.82,4.74,4.02,.86,18,GREEN_DARK,True,PP_ALIGN.CENTER)
    image(s,ASSETS/"season_mixture.png",5.30,1.18,7.48,4.92)
    rect(s,1.15,6.14,11.05,.56,GREEN_LIGHT,GREEN,True,.8)
    text(s,"Mixture groups observations; it does not model their temporal order.",1.40,6.25,10.55,.32,18,GREEN_DARK,True,PP_ALIGN.CENTER)
    footer(s,12,"Conceptual example: observations are continuous; component identity is discrete")

    # Slide 13: HMM.
    s=prs.slides[12]; rect(s,0,0,13.333,7.5,WHITE,None,False); title(s,"HMM: put the calendar back",13,"INTUITIVE EXAMPLE")
    add_cycle(s)
    formula(s,"p(zₜ | zₜ₋₁)",.72,2.46,2.55,.68,22,ORANGE_LIGHT,ORANGE)
    text(s,"The hidden season tends to persist,\nand some transitions are more plausible than others.",3.52,2.42,3.30,.78,17,CHARCOAL,True)
    formula(s,"xₜ | zₜ=k ∼ N(μₖ,Σₖ)",7.10,2.46,4.75,.68,21,BLUE_LIGHT,BLUE)
    image(s,ASSETS/"season_hmm.png",.72,3.38,11.92,2.77)
    rect(s,1.22,6.22,10.88,.50,GREEN_LIGHT,GREEN,True,.8)
    text(s,"The state is still discrete; the difference is temporal dependence.",1.48,6.30,10.38,.30,17,GREEN_DARK,True,PP_ALIGN.CENTER)
    footer(s,13,"A standard HMM does not know the seasons are cyclic unless the transition structure encodes it")

    # Slide 14: LDS.
    s=prs.slides[13]; rect(s,0,0,13.333,7.5,WHITE,None,False); title(s,"LDS: replace season labels with a continuous weather state",14,"INTUITIVE EXAMPLE")
    formula(s,"zₜ = [true temperature, trend]ᵀ",.62,1.28,4.30,.70,20,GREEN_LIGHT,GREEN)
    formula(s,"zₜ = A zₜ₋₁ + εₜ",.62,2.23,4.30,.70,22,LIGHT,None)
    formula(s,"yₜ = [1,0] zₜ + ηₜ",.62,3.18,4.30,.70,22,BLUE_LIGHT,BLUE)
    text(s,"The latent state is no longer a label.\nIt has a location and direction of change.",.82,4.36,3.88,.78,18,GREEN_DARK,True,PP_ALIGN.CENTER)
    image(s,ASSETS/"season_lds.png",5.18,1.15,7.52,4.95)
    rect(s,1.05,6.15,11.25,.56,ORANGE_LIGHT,ORANGE,True,.8)
    text(s,"Kalman inference combines a smooth dynamical prediction with noisy thermometer readings.",1.32,6.26,10.72,.32,17,ORANGE,True,PP_ALIGN.CENTER)
    footer(s,14,"Conceptual illustration; real seasonal weather is nonlinear and cyclic")

    renumber(prs.slides)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
