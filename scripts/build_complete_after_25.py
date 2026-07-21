from pathlib import Path
import shutil

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from build_parts2_3 import (
    title, footer, rect, text, formula, arrow, image, card,
    GREEN, GREEN_DARK, GREEN_LIGHT, BLUE, BLUE_LIGHT, ORANGE,
    ORANGE_LIGHT, CHARCOAL, MID, LIGHT, WHITE,
)


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_working_11.pptx"
OUTPUT = ROOT / "Presentation" / "Latent_Variable_Models_Neural_Data_Dynamics_complete_v2.pptx"
ASSETS = ROOT / "Presentation" / "complete_assets"
LFADS = ROOT / "external" / "computation-thru-dynamics" / "images" / "lfads_architecture_w_inferred_inputs_3.png"
GPSLDS = ROOT / "external" / "gpslds" / "figs" / "fig1.png"
ASSETS.mkdir(parents=True, exist_ok=True)


def delete_slide(prs, index):
    slide_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(slide_id.rId)
    prs.slides._sldIdLst.remove(slide_id)


def new_slide(prs, page, heading, section):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    rect(s, 0, 0, 13.333, 7.5, WHITE, None, False)
    title(s, heading, page, section)
    return s


def make_assets():
    rng = np.random.default_rng(5)
    t = np.linspace(0, 10, 150)
    truth = .55*np.sin(.72*t) + .12*t
    samples = truth[None, :] + rng.normal(0, .20, (45, len(t))) * (.5 + .5*np.sin(.3*t)**2)
    lo, hi = np.quantile(samples, [.08, .92], axis=0)
    fig, ax = plt.subplots(figsize=(8.2, 3.8))
    ax.fill_between(t, lo, hi, color="#B8DCCB", alpha=.75, label="posterior uncertainty")
    ax.plot(t, samples.mean(0), color="#008A45", lw=2.8, label="posterior mean")
    ax.scatter(t[::9], truth[::9] + rng.normal(0,.25,len(t[::9])), s=18, color="#3F78B5", alpha=.65, label="observations")
    ax.set_xlabel("time"); ax.set_ylabel("latent state"); ax.legend(frameon=False, ncol=3, loc="upper left")
    ax.spines[["top","right"]].set_visible(False); ax.grid(alpha=.12)
    fig.tight_layout(); fig.savefig(ASSETS/"posterior_trajectory.png", dpi=220, facecolor="white"); plt.close(fig)

    x = np.linspace(-4, 4, 500)
    p = .56*np.exp(-.5*((x+.65)/.72)**2)/(.72*np.sqrt(2*np.pi)) + .44*np.exp(-.5*((x-1.15)/.48)**2)/(.48*np.sqrt(2*np.pi))
    q = np.exp(-.5*((x-.05)/1.05)**2)/(1.05*np.sqrt(2*np.pi))
    fig, ax = plt.subplots(figsize=(7.0, 3.7))
    ax.fill_between(x,p,color="#F8D9C6",alpha=.8); ax.plot(x,p,color="#EE772F",lw=2.6,label="true posterior  pθ(z|y)")
    ax.fill_between(x,q,color="#DDE9F6",alpha=.6); ax.plot(x,q,color="#3F78B5",lw=2.6,label="approximation  qφ(z|y)")
    ax.set_yticks([]); ax.set_xlabel("latent z"); ax.legend(frameon=False); ax.spines[["top","right","left"]].set_visible(False)
    fig.tight_layout(); fig.savefig(ASSETS/"variational_approx.png",dpi=220,facecolor="white"); plt.close(fig)

    grid=np.linspace(-2.5,2.5,19); xx,yy=np.meshgrid(grid,grid)
    u=-.35*xx-1.05*yy+.28*np.sin(1.4*yy); v=.95*xx-.28*yy+.22*np.sin(1.3*xx)
    fig,ax=plt.subplots(figsize=(6.2,5.0)); ax.quiver(xx,yy,u,v,color="#82BEA2",alpha=.85,angles="xy",scale_units="xy",scale=6)
    z=np.array([1.0,.45]); f=np.array([-.35*z[0]-1.05*z[1]+.28*np.sin(1.4*z[1]), .95*z[0]-.28*z[1]+.22*np.sin(1.3*z[0])])
    ax.scatter(*z,s=100,color="#EE772F",zorder=4); ax.arrow(z[0],z[1],f[0]*.55,f[1]*.55,width=.025,head_width=.16,color="#008A45",zorder=5)
    ax.set_xlabel("latent 1"); ax.set_ylabel("latent 2"); ax.set_aspect("equal"); ax.grid(alpha=.1); ax.spines[["top","right"]].set_visible(False)
    fig.tight_layout(); fig.savefig(ASSETS/"nonlinear_local.png",dpi=220,facecolor="white"); plt.close(fig)

    fig,ax=plt.subplots(figsize=(6.4,4.2)); ax.axhline(0,color="#777",lw=1); ax.axvline(0,color="#777",lw=1)
    pts=[(-.9,0,"stable"),(.8,0,"unstable"),(-.35,1.05,"damped rotation"),(.35,1.0,"growing rotation"),(-.05,.05,"slow mode")]
    colors=[GREEN,ORANGE,BLUE,ORANGE,RGBColor(100,105,108)]
    for (a,b,l),c in zip(pts,colors):
        col='#%02x%02x%02x'%tuple(c); ax.scatter(a,b,s=90,color=col); ax.text(a+.08,b+.08,l,color=col,fontsize=11,weight='bold')
        if b: ax.scatter(a,-b,s=55,color=col,alpha=.55)
    ax.set_xlim(-1.4,1.4); ax.set_ylim(-1.55,1.55); ax.set_xlabel("Re(λ): contraction  ←  →  expansion"); ax.set_ylabel("Im(λ): oscillation frequency")
    ax.spines[["top","right"]].set_visible(False); ax.grid(alpha=.12)
    fig.tight_layout(); fig.savefig(ASSETS/"eigenvalue_plane.png",dpi=220,facecolor="white"); plt.close(fig)

    grid=np.linspace(-2.5,2.5,15); xx,yy=np.meshgrid(grid,grid); left=xx<0
    u=np.where(left,-.65*(xx+1.15),-.65*(xx-1.15)); v=-.70*yy
    fig,ax=plt.subplots(figsize=(7.2,4.5)); ax.axvspan(-2.7,0,color="#EAF2FB"); ax.axvspan(0,2.7,color="#FFF0E5")
    ax.quiver(xx,yy,u,v,color=np.where(left.ravel(),"#3F78B5","#EE772F"),alpha=.92,angles="xy",scale_units="xy",scale=5)
    ax.axvline(0,color="#777",lw=1.5,ls="--"); ax.scatter([-1.15,1.15],[0,0],marker="X",s=90,color=["#3F78B5","#EE772F"])
    ax.text(-2.35,2.15,"regime 1:  A₁, b₁",color="#3F78B5",weight="bold",fontsize=13); ax.text(.35,2.15,"regime 2:  A₂, b₂",color="#C45C1B",weight="bold",fontsize=13)
    ax.set_xlabel("latent dimension 1"); ax.set_ylabel("latent dimension 2"); ax.set_aspect("equal"); ax.set_xlim(-2.6,2.6); ax.set_ylim(-2.6,2.6); ax.grid(alpha=.1); ax.spines[["top","right"]].set_visible(False)
    fig.tight_layout(); fig.savefig(ASSETS/"piecewise.png",dpi=220,facecolor="white"); plt.close(fig)


def add_part3(prs):
    s=new_slide(prs,26,"Inference recovers a posterior over latent trajectories","PART III · INFERENCE")
    formula(s,"p(z₁:T | y₁:T)",.62,1.20,3.55,.78,27,GREEN_LIGHT,GREEN)
    text(s,"not only one best path",.82,2.24,3.15,.36,18,CHARCOAL,True,PP_ALIGN.CENTER)
    formula(s,"mean + uncertainty",.92,2.84,2.95,.64,20,BLUE_LIGHT,BLUE)
    image(s,ASSETS/"posterior_trajectory.png",4.45,1.18,8.15,4.48)
    formula(s,"filtering  p(zₜ|y₁:ₜ)     ·     smoothing  p(zₜ|y₁:T)",1.10,5.85,11.10,.62,20,LIGHT,None)
    footer(s,26,"Online inference uses the past; offline reconstruction can use the complete sequence")

    s=new_slide(prs,27,"Special structure makes exact inference possible","PART III · INFERENCE")
    rect(s,.60,1.18,5.95,4.70,BLUE_LIGHT,BLUE,True,1)
    text(s,"FINITE-STATE HMM",.90,1.48,5.35,.34,17,BLUE,True,PP_ALIGN.CENTER)
    formula(s,"αₜ(k) ∝ p(yₜ|zₜ=k) Σⱼ Pⱼₖ αₜ₋₁(j)",.92,2.05,5.30,.86,20,WHITE,BLUE)
    text(s,"finite sums\n+ dynamic programming\n\nforward–backward",1.35,3.30,4.42,1.48,20,CHARCOAL,True,PP_ALIGN.CENTER)
    rect(s,6.78,1.18,5.95,4.70,GREEN_LIGHT,GREEN,True,1)
    text(s,"LINEAR–GAUSSIAN LDS",7.08,1.48,5.35,.34,17,GREEN_DARK,True,PP_ALIGN.CENTER)
    formula(s,"ẑₜ|ₜ = ẑₜ|ₜ₋₁ + Kₜ(yₜ − Cẑₜ|ₜ₋₁)",7.10,2.05,5.30,.86,19,WHITE,GREEN)
    text(s,"linear maps\n+ Gaussian closure\n\nKalman filter / smoother",7.55,3.30,4.42,1.48,20,CHARCOAL,True,PP_ALIGN.CENTER)
    rect(s,1.15,6.15,11.02,.55,LIGHT,MID,True,.7)
    text(s,"Exactness comes from closure—not because latent-state inference is generally easy.",1.40,6.25,10.52,.32,17,CHARCOAL,True,PP_ALIGN.CENTER); footer(s,27)

    s=new_slide(prs,28,"Realistic neural models break analytic inference","PART III · INFERENCE")
    formula(s,"Gaussian observation\nyₜ = Czₜ + ηₜ",.62,1.25,3.50,1.20,21,BLUE_LIGHT,BLUE)
    arrow(s,4.30,1.86,5.08,1.86,ORANGE,2.2)
    formula(s,"Poisson spikes\nyₜₙ ~ Poisson(exp(cₙᵀzₜ+dₙ))",5.28,1.25,4.35,1.20,18,ORANGE_LIGHT,ORANGE)
    arrow(s,9.82,1.86,10.55,1.86,ORANGE,2.2)
    formula(s,"non-Gaussian\nposterior",10.72,1.25,2.00,1.20,18,LIGHT,MID)
    formula(s,"linear dynamics  Azₜ₋₁",.90,3.15,4.10,.72,21,GREEN_LIGHT,GREEN)
    arrow(s,5.20,3.51,6.00,3.51,ORANGE,2.2)
    formula(s,"nonlinear dynamics  fθ(zₜ₋₁)",6.20,3.15,5.10,.72,21,ORANGE_LIGHT,ORANGE)
    formula(s,"pθ(z₁:T | y₁:T)  is generally intractable",2.15,4.55,9.05,.82,25,LIGHT,None)
    rect(s,1.15,5.90,11.05,.70,GREEN_LIGHT,GREEN,True,.9)
    text(s,"If the posterior cannot be computed exactly, approximate it—and quantify the approximation.",1.42,6.05,10.52,.38,18,GREEN_DARK,True,PP_ALIGN.CENTER); footer(s,28)

    s=new_slide(prs,29,"Variational inference turns inference into optimization","PART III · DEEP GENERATIVE MODELS")
    formula(s,"qφ(z|y)  ≈  pθ(z|y)",.65,1.30,4.20,.82,28,GREEN_LIGHT,GREEN)
    text(s,"choose a tractable family\nand optimize its parameters φ",.82,2.55,3.85,.90,20,CHARCOAL,True,PP_ALIGN.CENTER)
    image(s,ASSETS/"variational_approx.png",5.15,1.13,7.40,4.50)
    rect(s,.92,4.02,3.70,1.28,BLUE_LIGHT,BLUE,True,.9)
    text(s,"Amortization",1.20,4.22,3.10,.32,18,BLUE,True,PP_ALIGN.CENTER)
    text(s,"one encoder predicts posterior parameters for new data",1.15,4.62,3.20,.50,15,CHARCOAL,False,PP_ALIGN.CENTER)
    rect(s,1.10,6.05,11.10,.62,LIGHT,MID,True,.7)
    text(s,"θ learns the generative model; φ learns how to infer its hidden variables.",1.35,6.18,10.60,.34,18,CHARCOAL,True,PP_ALIGN.CENTER); footer(s,29)

    s=new_slide(prs,30,"Standard VAE: inference model + generative model","PART III · STANDARD VAE")
    labels=[("DATA","y",BLUE), ("ENCODER","μ(y), σ(y)",GREEN), ("LATENT","z",ORANGE), ("DECODER","pθ(y|z)",GREEN), ("RECONSTRUCTION","ŷ",BLUE)]
    xs=[.55,3.00,5.63,7.85,10.45]; widths=[1.65,1.95,1.35,1.95,2.25]
    for i,((head,body,col),x,w) in enumerate(zip(labels,xs,widths)):
        rect(s,x,2.05,w,1.48,WHITE,col,True,1.1); text(s,head,x+.10,2.25,w-.20,.28,13,col,True,PP_ALIGN.CENTER); text(s,body,x+.10,2.73,w-.20,.42,22,CHARCOAL,True,PP_ALIGN.CENTER)
        if i<4: arrow(s,x+w+.08,2.80,x+w+.65,2.80,GREEN,2)
    formula(s,"qφ(z|y)=N(μφ(y), diag(σ²φ(y)))",1.15,4.35,5.25,.72,21,BLUE_LIGHT,BLUE)
    formula(s,"pθ(y|z)",7.10,4.35,3.95,.72,23,GREEN_LIGHT,GREEN)
    rect(s,1.20,5.75,10.95,.70,ORANGE_LIGHT,ORANGE,True,.9)
    text(s,"The encoder approximates inference; the decoder specifies how latent causes generate observations.",1.47,5.90,10.42,.38,18,ORANGE,True,PP_ALIGN.CENTER); footer(s,30)

    s=new_slide(prs,31,"One objective trains reconstruction and representation","PART III · STANDARD VAE")
    text(s,"REPARAMETERIZE",.72,1.18,3.85,.32,15,ORANGE,True,PP_ALIGN.CENTER)
    formula(s,"ε ~ N(0,I)\nz = μ + σ ⊙ ε",.72,1.65,3.85,1.28,24,ORANGE_LIGHT,ORANGE)
    text(s,"randomness is isolated in ε,\nso gradients reach μ and σ",.85,3.20,3.60,.78,17,CHARCOAL,False,PP_ALIGN.CENTER)
    text(s,"OPTIMIZE THE ELBO",5.05,1.18,7.45,.32,15,GREEN_DARK,True,PP_ALIGN.CENTER)
    formula(s,"L = E_q[log pθ(y|z)] − KL[qφ(z|y) || p(z)]",5.05,1.65,7.45,1.28,23,GREEN_LIGHT,GREEN)
    card(s,"EXPLAIN THE DATA","reconstruction / likelihood",5.25,3.28,3.15,1.40,BLUE,16)
    card(s,"CONTROL COMPLEXITY","stay compatible with the prior",8.85,3.28,3.15,1.40,ORANGE,16)
    rect(s,1.18,5.78,11.02,.72,LIGHT,MID,True,.8)
    text(s,"The ELBO is not merely a reconstruction loss: it is a lower bound on log evidence.",1.45,5.94,10.48,.38,18,CHARCOAL,True,PP_ALIGN.CENTER); footer(s,31)

    s=new_slide(prs,32,"A standard VAE is not yet a dynamical model","PART III · FROM VAE TO DYNAMICS")
    text(s,"STANDARD VAE",.90,1.20,4.85,.34,17,BLUE,True,PP_ALIGN.CENTER)
    formula(s,"yₜ → qφ(zₜ|yₜ) → zₜ → ŷₜ",.82,1.75,5.00,.82,23,BLUE_LIGHT,BLUE)
    text(s,"Each observation may be encoded independently.",1.15,2.92,4.32,.50,18,CHARCOAL,True,PP_ALIGN.CENTER)
    text(s,"DYNAMICAL LATENT MODEL",7.15,1.20,5.20,.34,17,GREEN_DARK,True,PP_ALIGN.CENTER)
    formula(s,"z₁ → z₂ → ··· → z_T",7.20,1.75,5.05,.82,25,GREEN_LIGHT,GREEN)
    text(s,"The generative model specifies how states evolve through time.",7.52,2.92,4.42,.60,18,CHARCOAL,True,PP_ALIGN.CENTER)
    arrow(s,5.92,2.17,6.95,2.17,ORANGE,2.5)
    formula(s,"p(z₁:T)=p(z₁) ∏ₜ₌₂ᵀ p(zₜ|zₜ₋₁)",2.15,4.25,9.05,.78,24,LIGHT,None)
    rect(s,1.15,5.85,11.05,.68,ORANGE_LIGHT,ORANGE,True,.9)
    text(s,"For neural dynamics, the decoder must generate a trajectory—not only reconstruct a sample.",1.42,6.00,10.52,.38,18,ORANGE,True,PP_ALIGN.CENTER); footer(s,32)

    s=new_slide(prs,33,"LFADS makes the decoder a nonlinear dynamical generator","PART III · LFADS")
    image(s,LFADS,.45,1.12,12.42,3.60)
    formula(s,"g₀ ~ qφ(g₀|y₁:T)     →     gₜ = RNNθ(gₜ₋₁,uₜ)     →     λₜ     →     yₜ",1.05,5.05,11.20,.74,21,GREEN_LIGHT,GREEN)
    rect(s,1.22,6.08,10.90,.58,BLUE_LIGHT,BLUE,True,.8)
    text(s,"The inferred initial condition and inputs generate a single-trial rate trajectory.",1.48,6.20,10.38,.32,17,BLUE,True,PP_ALIGN.CENTER)
    footer(s,33,"Architecture adapted from the computation-through-dynamics LFADS tutorial")

    s=new_slide(prs,34,"LFADS: powerful generator, conditional explanation","PART III · LFADS")
    rect(s,.62,1.20,5.85,4.68,GREEN_LIGHT,GREEN,True,1)
    text(s,"WHAT IT ENABLES",.95,1.52,5.20,.34,17,GREEN_DARK,True,PP_ALIGN.CENTER)
    text(s,"• single-trial denoising\n\n• nonlinear recurrent trajectories\n\n• inferred inputs and initial conditions\n\n• behavioral prediction",1.20,2.12,4.72,2.92,19,CHARCOAL,False)
    rect(s,6.85,1.20,5.85,4.68,ORANGE_LIGHT,ORANGE,True,1)
    text(s,"WHAT REMAINS CONDITIONAL",7.18,1.52,5.20,.34,17,ORANGE,True,PP_ALIGN.CENTER)
    text(s,"• latent coordinates are not unique\n\n• flexible RNN fields can be hard to inspect\n\n• reconstruction does not prove mechanism\n\n• inferred inputs require validation",7.40,2.12,4.72,2.92,19,CHARCOAL,False)
    rect(s,1.02,6.15,11.30,.58,LIGHT,MID,True,.8)
    text(s,"Next: retain nonlinear flexibility while making local dynamics inspectable.",1.28,6.27,10.78,.32,18,CHARCOAL,True,PP_ALIGN.CENTER); footer(s,34)


def add_part4(prs):
    s=new_slide(prs,35,"Good spike reconstruction does not guarantee correct dynamics","PART IV · INTERPRETABILITY")
    formula(s,"high held-out likelihood",.75,1.42,3.70,.72,23,BLUE_LIGHT,BLUE)
    text(s,"≠",4.62,1.52,.60,.52,28,ORANGE,True,PP_ALIGN.CENTER)
    formula(s,"correct vector field",5.38,1.42,3.70,.72,23,ORANGE_LIGHT,ORANGE)
    text(s,"≠",9.25,1.52,.60,.52,28,ORANGE,True,PP_ALIGN.CENTER)
    formula(s,"causal mechanism",10.00,1.42,2.62,.72,22,GREEN_LIGHT,GREEN)
    card(s,"IDENTIFIABILITY","different latent coordinates can explain the same observations",.72,3.00,3.75,1.52,BLUE,15)
    card(s,"EXTRAPOLATION","the learned field may be unsupported away from observed trajectories",4.80,3.00,3.75,1.52,ORANGE,15)
    card(s,"CAUSALITY","perturbations—not visualization alone—test whether dynamics are intrinsic",8.88,3.00,3.75,1.52,GREEN,15)
    rect(s,1.15,5.82,11.03,.72,GREEN_LIGHT,GREEN,True,.9)
    text(s,"An interpretable vector field is a scientific hypothesis with explicit tests.",1.42,5.98,10.50,.38,20,GREEN_DARK,True,PP_ALIGN.CENTER); footer(s,35)

    s=new_slide(prs,36,"rSLDS lets position influence which local dynamics apply","PART IV · INTERPRETABILITY")
    formula(s,"SLDS:  p(sₜ₊₁ | sₜ)",.62,1.22,4.20,.72,23,LIGHT,MID)
    formula(s,"rSLDS:  p(sₜ₊₁ | sₜ, zₜ)",.62,2.18,4.20,.72,23,ORANGE_LIGHT,ORANGE)
    text(s,"The continuous location zₜ changes the probability of switching regimes.",.82,3.25,3.82,.90,19,CHARCOAL,True,PP_ALIGN.CENTER)
    image(s,ASSETS/"piecewise.png",5.05,1.08,7.35,4.75)
    rect(s,1.02,6.05,11.30,.62,GREEN_LIGHT,GREEN,True,.8)
    text(s,"The result is globally nonlinear behavior assembled from interpretable local linear rules.",1.28,6.18,10.78,.34,18,GREEN_DARK,True,PP_ALIGN.CENTER); footer(s,36)

    s=new_slide(prs,37,"gpSLDS: a smooth nonlinear field with uncertainty","PART IV · INTERPRETABILITY")
    formula(s,"ż = f(z) + noise",.62,1.25,3.80,.72,25,GREEN_LIGHT,GREEN)
    formula(s,"fₖ(·) ~ GP(0, κSSL)",.62,2.20,3.80,.72,22,BLUE_LIGHT,BLUE)
    text(s,"A structured kernel favors locally simple dynamics while allowing the field to bend smoothly.",.75,3.25,3.55,1.15,18,CHARCOAL,True,PP_ALIGN.CENTER)
    image(s,GPSLDS,4.72,1.08,7.92,4.88)
    rect(s,1.10,6.10,11.10,.60,ORANGE_LIGHT,ORANGE,True,.8)
    text(s,"Unlike a single fitted trajectory, a probabilistic field represents both dynamics and uncertainty.",1.35,6.22,10.60,.34,18,ORANGE,True,PP_ALIGN.CENTER); footer(s,37,"gpSLDS repository figure: structured GP dynamics and point-process observations")

    s=new_slide(prs,38,"Local linearization makes a nonlinear field inspectable","PART IV · LOCAL ANALYSIS")
    formula(s,"f(z) ≈ f(z*) + Jf(z*)(z − z*)",.62,1.22,5.30,.78,25,GREEN_LIGHT,GREEN)
    formula(s,"Jf(z*) = ∂f/∂z |z=z*",.92,2.28,4.70,.70,22,LIGHT,None)
    text(s,"Near z*, the Jacobian plays the same role as the LDS matrix A.",.78,3.32,4.95,.78,19,CHARCOAL,True,PP_ALIGN.CENTER)
    image(s,ASSETS/"nonlinear_local.png",6.20,1.05,6.20,4.90)
    rect(s,.85,4.62,4.80,1.05,ORANGE_LIGHT,ORANGE,True,.9)
    text(s,"If f(z*) = 0, z* is a fixed point.\nThe local field is approximately linear around it.",1.08,4.82,4.34,.62,17,ORANGE,True,PP_ALIGN.CENTER)
    rect(s,1.08,6.10,11.15,.60,GREEN_LIGHT,GREEN,True,.8)
    text(s,"This is how LDS intuition becomes a tool for analyzing nonlinear dynamics.",1.35,6.22,10.60,.34,18,GREEN_DARK,True,PP_ALIGN.CENTER); footer(s,38)

    s=new_slide(prs,39,"Eigenvalues summarize local stability and timescale","PART IV · LOCAL ANALYSIS")
    image(s,ASSETS/"eigenvalue_plane.png",.55,1.10,6.65,4.72)
    rows=[("Re(λ) < 0","contraction · locally stable",GREEN), ("Re(λ) > 0","expansion · locally unstable",ORANGE), ("Im(λ) ≠ 0","rotation / oscillation",BLUE), ("Re(λ) ≈ 0","slow mode · line-attractor-like",MID)]
    for i,(eq,meaning,col) in enumerate(rows):
        y=1.30+i*1.03; rect(s,7.55,y,4.95,.78,WHITE,col,True,.9); text(s,eq,7.78,y+.16,1.62,.32,17,col,True); text(s,meaning,9.40,y+.13,2.88,.38,16,CHARCOAL,False,PP_ALIGN.RIGHT)
    formula(s,"τ ≈ −1 / Re(λ)",8.30,5.55,3.45,.66,22,LIGHT,None)
    footer(s,39,"Interpret eigenvalues locally: they depend on the point and coordinate representation")

    s=new_slide(prs,40,"Validation must rise with the strength of the scientific claim","PART IV · VALIDATION")
    steps=[("1","IMPLEMENTATION","tests · simulation",BLUE), ("2","PREDICTION","held-out likelihood",BLUE), ("3","RECOVERY","known synthetic dynamics",ORANGE), ("4","RELEVANCE","behavior · cross-session",ORANGE), ("5","CAUSALITY","targeted perturbation",GREEN)]
    for i,(num,head,body,col) in enumerate(steps):
        x=.48+i*2.56; y=4.72-i*.72
        rect(s,x,y,2.25,1.15,WHITE,col,True,1); text(s,num,x+.08,y+.12,.40,.35,18,col,True,PP_ALIGN.CENTER); text(s,head,x+.48,y+.12,1.60,.28,13,col,True,PP_ALIGN.CENTER); text(s,body,x+.20,y+.58,1.86,.34,13,CHARCOAL,False,PP_ALIGN.CENTER)
        if i<4: arrow(s,x+2.28,y+.55,x+2.48,y-.10,GREEN,1.7)
    text(s,"A validation ladder",.72,1.35,4.10,.42,25,CHARCOAL,True)
    text(s,"Each step rules out a stronger alternative explanation.",.75,2.08,5.62,.54,19,MID,False)
    formula(s,"fit  →  predict  →  recover  →  generalize  →  perturb",4.75,1.35,7.80,.74,22,GREEN_LIGHT,GREEN)
    rect(s,1.10,6.22,11.10,.50,LIGHT,MID,True,.7); text(s,"A beautiful latent trajectory is the beginning of validation—not the end.",1.35,6.30,10.60,.30,17,CHARCOAL,True,PP_ALIGN.CENTER); footer(s,40)

    s=new_slide(prs,41,"Take-home: useful latent dynamics remain testable","PART IV · SUMMARY")
    items=[("01","OBSERVE","High-dimensional variability motivates hidden-state explanations."), ("02","MODEL","Mixture, HMM, LDS, SLDS and PLDS encode different assumptions."), ("03","INFER","Exact inference is special; variational inference extends the model class."), ("04","INTERPRET","Local fields, fixed points and eigenvalues translate fits into hypotheses."), ("05","VALIDATE","Prediction, recovery and perturbation determine what we may claim.")]
    for i,(n,h,b) in enumerate(items):
        y=1.12+i*1.02; text(s,n,.72,y,.48,.34,14,GREEN,True,PP_ALIGN.CENTER); rect(s,1.32,y-.04,.06,.66,GREEN,None,False); text(s,h,1.62,y,1.52,.34,16,CHARCOAL,True); text(s,b,3.15,y,9.10,.42,17,MID,False)
    rect(s,1.10,6.30,11.10,.45,GREEN_LIGHT,GREEN,True,.8); text(s,"The latent state is not discovered by the model alone—it earns meaning through successful tests.",1.35,6.37,10.60,.28,17,GREEN_DARK,True,PP_ALIGN.CENTER); footer(s,41)


def add_tutorial(prs):
    s=prs.slides.add_slide(prs.slide_layouts[0]); rect(s,0,0,13.333,7.5,CHARCOAL,None,False)
    rect(s,.72,1.05,.13,5.10,GREEN,None,False); text(s,"PRACTICAL EPILOGUE",1.15,1.18,3.10,.36,15,GREEN,True)
    text(s,"Tutorial Speedrun",1.12,1.90,10.75,.82,40,WHITE,True)
    text(s,"From PyTorch tensors to variational neural inference",1.15,2.92,10.20,.58,23,RGBColor(210,218,214),False)
    formula(s,"FOUNDATIONS  →  DYNAMICS  →  VARIATIONAL MODELS  →  GPSLDS",1.15,4.35,10.85,.78,20,RGBColor(61,66,67),GREEN,WHITE)
    text(s,"A map of the repository—not a second lecture.",1.15,5.65,9.40,.42,17,RGBColor(190,200,195),False); text(s,"42",12.45,7.05,.40,.20,8,RGBColor(190,200,195),False,PP_ALIGN.RIGHT)

    s=new_slide(prs,43,"The tutorials follow the same conceptual ladder as the talk","TUTORIAL SPEEDRUN")
    groups=[
        ("0  FOUNDATIONS",["00_pytorch_basics","01_pytorch_neuroscience_introduction"],BLUE),
        ("1–2  CLASSICAL",["02_probabilistic_modeling","02b_mixtures_em","03a_hmm_foundations · 03_hmm_lds"],ORANGE),
        ("3  VARIATIONAL",["04_standard_vae","05_variational_em","06_lfads · 07_FULL_LFADS_Tutorial"],GREEN),
        ("4  ADVANCED",["08_gpslds"],BLUE),
    ]
    for i,(head,files,col) in enumerate(groups):
        x=.45+i*3.18; rect(s,x,1.28,2.90,4.65,WHITE,col,True,1); text(s,head,x+.16,1.58,2.58,.42,15,col,True,PP_ALIGN.CENTER)
        for j,f in enumerate(files): rect(s,x+.18,2.28+j*.85,2.54,.62,LIGHT,None,True); text(s,f,x+.28,2.38+j*.85,2.34,.36,12,CHARCOAL,True,PP_ALIGN.CENTER)
    rect(s,1.12,6.22,11.08,.48,GREEN_LIGHT,GREEN,True,.8); text(s,"Run the notebooks in order only when you need the full implementation path.",1.37,6.30,10.58,.30,17,GREEN_DARK,True,PP_ALIGN.CENTER); footer(s,43)

    s=new_slide(prs,44,"Three routes through the repository","TUTORIAL SPEEDRUN")
    routes=[
        ("I WANT THE BASICS","00 → 01 → 02 → 04","PyTorch, probability, then a standard VAE",BLUE),
        ("I WANT DYNAMICS","03a → 03 → 06 → 07","HMM/LDS foundations, then LFADS",ORANGE),
        ("I WANT INTERPRETABILITY","03 → 08","Linear dynamics, vector fields and gpSLDS",GREEN),
    ]
    for i,(head,path,desc,col) in enumerate(routes):
        y=1.25+i*1.55; rect(s,.75,y,11.85,1.20,WHITE,col,True,1.0); text(s,head,1.05,y+.18,2.35,.32,15,col,True); formula(s,path,3.52,y+.16,3.38,.68,19,LIGHT,None); text(s,desc,7.18,y+.18,5.02,.55,17,CHARCOAL,False,PP_ALIGN.CENTER)
    rect(s,1.15,6.15,11.05,.60,ORANGE_LIGHT,ORANGE,True,.8); text(s,"Best practice: change one assumption, inspect one diagnostic, and explain one failure mode.",1.42,6.27,10.52,.34,18,ORANGE,True,PP_ALIGN.CENTER); footer(s,44,"Repository root: VAE-Tutorial/")

    s=prs.slides.add_slide(prs.slide_layouts[0]); rect(s,0,0,13.333,7.5,CHARCOAL,None,False)
    text(s,"THANK YOU",.85,.70,3.00,.40,16,GREEN,True); text(s,"Questions?",.82,2.15,8.60,1.05,46,WHITE,True)
    text(s,"Latent Variable Models for Neural Data and Dynamics",.88,3.55,10.75,.52,22,RGBColor(210,218,214),False)
    text(s,"Chunyu Qu",.88,4.42,4.30,.42,18,GREEN,True); text(s,"45",12.45,7.05,.40,.20,8,RGBColor(190,200,195),False,PP_ALIGN.RIGHT)


def build():
    make_assets(); shutil.copy2(SOURCE,OUTPUT); prs=Presentation(OUTPUT)
    while len(prs.slides)>25: delete_slide(prs,25)
    add_part3(prs); add_part4(prs); add_tutorial(prs)
    prs.save(OUTPUT); print(OUTPUT)


if __name__=="__main__": build()
